"""
Report generation utilities for Excel and PowerPoint outputs.

Generates clean summary reports with tables, KPIs, and embedded plots.
"""

from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any

import numpy as np
import pandas as pd

from .config import AnalysisConfig, GroundStationConfig, ImagingModeConfig


def _convert_datetimes_to_string(df: pd.DataFrame) -> pd.DataFrame:
    """Convert all datetime columns in a DataFrame to strings for Excel compatibility."""
    df = df.copy()
    for col in df.columns:
        # Handle datetime64 columns
        if 'datetime' in str(df[col].dtype):
            try:
                # Remove timezone if present, then format
                if df[col].dt.tz is not None:
                    df[col] = df[col].dt.tz_localize(None)
                df[col] = df[col].dt.strftime('%Y-%m-%d %H:%M:%S')
            except Exception:
                df[col] = df[col].astype(str)
        # Handle object columns that might contain Timestamps
        elif df[col].dtype == 'object' and len(df) > 0:
            try:
                first_valid = df[col].dropna().iloc[0] if not df[col].dropna().empty else None
                if first_valid is not None and hasattr(first_valid, 'strftime'):
                    df[col] = df[col].apply(
                        lambda x: x.strftime('%Y-%m-%d %H:%M:%S') if pd.notna(x) and hasattr(x, 'strftime') else x
                    )
            except (IndexError, TypeError):
                pass
    return df


def generate_excel_report(
    config: AnalysisConfig,
    access_df: pd.DataFrame,
    mode_dfs: Dict[str, pd.DataFrame],
    ttnc_df: pd.DataFrame,
    backlog_df: pd.DataFrame,
    coverage_kpis: Dict[str, Any],
    contact_kpis: Dict[str, Any],
    downlink_kpis: Dict[str, Any],
    output_path: Optional[Path] = None,
    optimization_results: Optional[Dict[str, Any]] = None,
    propulsion_df: Optional[pd.DataFrame] = None,
) -> Path:
    """
    Generate Excel report with all analysis results.

    Parameters
    ----------
    config : AnalysisConfig
        Analysis configuration.
    access_df : pd.DataFrame
        Access windows DataFrame.
    mode_dfs : Dict[str, pd.DataFrame]
        Dictionary of mode DataFrames.
    ttnc_df : pd.DataFrame
        TTNC DataFrame.
    backlog_df : pd.DataFrame
        Backlog time series.
    coverage_kpis : Dict
        Coverage KPIs from validation.
    contact_kpis : Dict
        Contact KPIs from validation.
    downlink_kpis : Dict
        Downlink KPIs.
    output_path : Path, optional
        Output path for the Excel file.

    Returns
    -------
    Path
        Path to the generated Excel file.
    """
    if output_path is None:
        output_path = config.output_path / config.excel_filename

    output_path.parent.mkdir(parents=True, exist_ok=True)

    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        # Sheet 1: Configuration
        config_data = {
            'Parameter': [
                'Start Date',
                'Duration (days)',
                'Time Step (s)',
                'Number of Satellites',
                'Number of Ground Stations',
                'Number of Imaging Modes',
                'Min Access Duration (s)',
                'Slew Rate (deg/s)',
                'Mode Switch Time (s)',
                'Contact Buffer (s)',
                'TTNC Max Search (hours)',
            ],
            'Value': [
                config.start_date,
                config.duration_days,
                config.time_step_s,
                len(config.satellites),
                len(config.ground_stations),
                len(config.imaging_modes),
                config.min_access_duration_s,
                config.slew_rate_deg_per_s,
                config.mode_switch_time_s,
                config.contact_buffer_s,
                config.ttnc_max_search_hours,
            ]
        }
        pd.DataFrame(config_data).to_excel(writer, sheet_name='Config', index=False)

        # Satellite details
        sat_data = []
        for sat in config.satellites:
            sat_data.append({
                'Satellite ID': sat.sat_id,
                'Inclination (deg)': sat.inclination_deg,
                'Altitude (km)': sat.altitude_km,
                'RAAN (deg)': sat.raan_deg,
            })
        if sat_data:
            pd.DataFrame(sat_data).to_excel(writer, sheet_name='Satellites', index=False)

        # Ground station details
        gs_data = []
        for gs in config.ground_stations:
            gs_data.append({
                'Name': gs.name,
                'Latitude': gs.lat,
                'Longitude': gs.lon,
                'Elevation (m)': gs.elevation_m,
                'Min Elevation (deg)': gs.min_elevation_deg,
                'Ka Capable': gs.ka_capable,
                'Downlink Rate (Mbps)': gs.downlink_rate_mbps,
            })
        if gs_data:
            pd.DataFrame(gs_data).to_excel(writer, sheet_name='Ground_Stations', index=False)

        # Imaging mode details
        mode_data = []
        for mode in config.imaging_modes:
            mode_data.append({
                'Name': mode.name,
                'FOV Half Angle (deg)': mode.fov_half_angle_deg,
                'Dwell Time (s)': mode.collect_dwell_time_s,
                'Collection Rate (Gbps)': mode.collection_rate_gbps,
                'Processing Time (min)': mode.processing_time_min,
                'Off-Nadir Min (deg)': mode.off_nadir_min_deg,
                'Off-Nadir Max (deg)': mode.off_nadir_max_deg,
            })
        if mode_data:
            pd.DataFrame(mode_data).to_excel(writer, sheet_name='Imaging_Modes', index=False)

        # Sheet 2: Coverage KPIs
        kpi_data = []
        if coverage_kpis and 'stats' in coverage_kpis:
            stats = coverage_kpis['stats']
            kpi_data.extend([
                {'KPI': 'Total Accesses', 'Value': stats.get('total_accesses', 0), 'Unit': 'count'},
                {'KPI': 'Accesses per Day', 'Value': round(stats.get('accesses_per_day', 0), 2), 'Unit': 'count/day'},
                {'KPI': 'Mean Access Duration', 'Value': round(stats.get('mean_duration_s', 0), 1), 'Unit': 'seconds'},
                {'KPI': 'Max Access Duration', 'Value': round(stats.get('max_duration_s', 0), 1), 'Unit': 'seconds'},
                {'KPI': 'Mean Revisits per Target', 'Value': round(stats.get('mean_revisits', 0), 1), 'Unit': 'count'},
            ])

        if contact_kpis and 'stats' in contact_kpis:
            stats = contact_kpis['stats']
            kpi_data.extend([
                {'KPI': 'Total Valid Contacts', 'Value': stats.get('total_valid_contacts', 0), 'Unit': 'count'},
                {'KPI': 'Valid Contact Percentage', 'Value': round(stats.get('valid_percentage', 0), 1), 'Unit': '%'},
            ])

        if downlink_kpis:
            kpi_data.extend([
                {'KPI': 'Total Data Collected', 'Value': round(downlink_kpis.get('total_data_collected_gb', 0), 2), 'Unit': 'GB'},
                {'KPI': 'Total Data Downlinked', 'Value': round(downlink_kpis.get('total_data_downlinked_gb', 0), 2), 'Unit': 'GB'},
                {'KPI': 'Peak Backlog', 'Value': round(downlink_kpis.get('peak_backlog_gb', 0), 2), 'Unit': 'GB'},
                {'KPI': 'Downlink Efficiency', 'Value': round(downlink_kpis.get('downlink_efficiency_pct', 0), 1), 'Unit': '%'},
            ])

        if kpi_data:
            pd.DataFrame(kpi_data).to_excel(writer, sheet_name='Coverage_KPIs', index=False)

        # Sheet 3: Access Windows
        if not access_df.empty:
            access_export = access_df.copy()
            # Convert datetime columns to string for Excel
            access_export = _convert_datetimes_to_string(access_export)
            access_export.to_excel(writer, sheet_name='Access_Windows', index=False)

        # Sheet 4: Contacts (consolidated from all modes)
        contacts_list = []
        for mode_name, df in mode_dfs.items():
            mode_contacts = df.copy()
            mode_contacts['Imaging_Mode'] = mode_name
            # Select key columns
            key_cols = ['Imaging_Mode', 'Satellite', 'Start_Time', 'End_Time', 'AOI_Lat', 'AOI_Lon',
                        'Access_Duration', 'Total_Valid_Contact', 'Collect_to_Delivery_Timeline']
            # Add ground station columns
            for gs in config.ground_stations:
                for col_suffix in ['Contact_Start', 'Contact_End', 'Contact_Duration']:
                    col = f'Next_{gs.name}_{col_suffix}'
                    if col in mode_contacts.columns:
                        key_cols.append(col)
                valid_col = f'Valid_Contact_{gs.name}'
                if valid_col in mode_contacts.columns:
                    key_cols.append(valid_col)

            available_cols = [c for c in key_cols if c in mode_contacts.columns]
            contacts_list.append(mode_contacts[available_cols])

        if contacts_list:
            contacts_df = pd.concat(contacts_list, ignore_index=True)
            contacts_df = _convert_datetimes_to_string(contacts_df)
            contacts_df.to_excel(writer, sheet_name='Contacts', index=False)

        # Sheet 5: Downlink KPIs by Ground Station
        if downlink_kpis and 'ground_station_stats' in downlink_kpis:
            gs_stats_list = []
            for gs_name, stats in downlink_kpis['ground_station_stats'].items():
                gs_stats_list.append({
                    'Ground Station': gs_name,
                    'Total Contacts': stats.get('total_contacts', 0),
                    'Valid Contacts': stats.get('valid_contacts', 0),
                    'Contact Duration (hours)': round(stats.get('total_contact_duration_hours', 0), 2),
                    'Capacity (GB)': round(stats.get('total_capacity_gb', 0), 2),
                })
            if gs_stats_list:
                pd.DataFrame(gs_stats_list).to_excel(writer, sheet_name='Downlink_KPIs', index=False)

        # Sheet 6: Backlog Time Series
        if not backlog_df.empty:
            backlog_export = _convert_datetimes_to_string(backlog_df)
            backlog_export.to_excel(writer, sheet_name='Backlog_TimeSeries', index=False)

        # Sheet 7: TTNC Ka
        if not ttnc_df.empty:
            ttnc_export = _convert_datetimes_to_string(ttnc_df)
            ttnc_export.to_excel(writer, sheet_name='TTNC_Ka', index=False)

            # TTNC Summary
            ttnc_summary = []
            for mode in ttnc_df['Imaging_Mode'].unique():
                mode_data = ttnc_df[ttnc_df['Imaging_Mode'] == mode]['TTNC_Ka_minutes'].dropna()
                if not mode_data.empty:
                    ttnc_summary.append({
                        'Imaging Mode': mode,
                        'Count': len(mode_data),
                        'Median (min)': round(mode_data.median(), 1),
                        'P90 (min)': round(mode_data.quantile(0.90), 1),
                        'P95 (min)': round(mode_data.quantile(0.95), 1),
                        'Max (min)': round(mode_data.max(), 1),
                    })
            if ttnc_summary:
                pd.DataFrame(ttnc_summary).to_excel(writer, sheet_name='TTNC_Summary', index=False)

        # Sheet 8: Optimization Results (if available)
        if optimization_results:
            from .optimization import optimization_results_to_dataframe
            opt_summary_df = optimization_results_to_dataframe(optimization_results)
            opt_summary_df.to_excel(writer, sheet_name='Optimization_Summary', index=False)

            # Detailed station lists for each provider
            opt_details = []
            for provider, result in optimization_results.items():
                for station in result.ttc_stations:
                    opt_details.append({
                        'Provider': provider,
                        'Type': 'TT&C',
                        'Station Name': station.name,
                        'Latitude': station.lat,
                        'Longitude': station.lon,
                        'S-Band': station.s_band,
                        'X-Band': station.x_band,
                        'Ka-Band': station.ka_band,
                    })
                for station in result.ka_stations:
                    opt_details.append({
                        'Provider': provider,
                        'Type': 'Ka Payload',
                        'Station Name': station.name,
                        'Latitude': station.lat,
                        'Longitude': station.lon,
                        'S-Band': station.s_band,
                        'X-Band': station.x_band,
                        'Ka-Band': station.ka_band,
                    })
            if opt_details:
                pd.DataFrame(opt_details).to_excel(writer, sheet_name='Optimization_Stations', index=False)

        # Sheet 9: Propulsion/Station-Keeping Analysis (if available)
        if propulsion_df is not None and not propulsion_df.empty:
            propulsion_df.to_excel(writer, sheet_name='Propulsion_Analysis', index=False)

    print(f"Excel report saved to: {output_path}")
    return output_path


def generate_ppt_report(
    config: AnalysisConfig,
    coverage_kpis: Dict[str, Any],
    contact_kpis: Dict[str, Any],
    downlink_kpis: Dict[str, Any],
    ttnc_df: pd.DataFrame,
    plot_paths: Dict[str, Path],
    output_path: Optional[Path] = None,
    optimization_results: Optional[Dict[str, Any]] = None,
) -> Path:
    """
    Generate PowerPoint report with plots and KPIs.

    Parameters
    ----------
    config : AnalysisConfig
        Analysis configuration.
    coverage_kpis : Dict
        Coverage KPIs.
    contact_kpis : Dict
        Contact KPIs.
    downlink_kpis : Dict
        Downlink KPIs.
    ttnc_df : pd.DataFrame
        TTNC DataFrame.
    plot_paths : Dict[str, Path]
        Dictionary mapping plot names to file paths.
    output_path : Path, optional
        Output path for the PowerPoint file.

    Returns
    -------
    Path
        Path to the generated PowerPoint file.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("Warning: python-pptx not installed. Skipping PowerPoint generation.")
        print("Install with: pip install python-pptx")
        return None

    if output_path is None:
        output_path = config.output_path / config.ppt_filename

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "VLEO EO Coverage Analysis"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"Analysis Period: {config.start_date} to {config.end_datetime.strftime('%Y-%m-%d')}"
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    config_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    config_frame = config_box.text_frame
    config_para = config_frame.paragraphs[0]
    config_para.text = (
        f"Satellites: {len(config.satellites)} | "
        f"Ground Stations: {len(config.ground_stations)} | "
        f"Imaging Modes: {len(config.imaging_modes)}"
    )
    config_para.font.size = Pt(18)
    config_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Coverage KPIs
    slide = prs.slides.add_slide(slide_layout)
    _add_slide_title(slide, "Coverage Summary")

    kpi_text = []
    if coverage_kpis and 'stats' in coverage_kpis:
        stats = coverage_kpis['stats']
        kpi_text.append(f"Total Access Windows: {stats.get('total_accesses', 0)}")
        kpi_text.append(f"Accesses per Day: {stats.get('accesses_per_day', 0):.1f}")
        kpi_text.append(f"Mean Duration: {stats.get('mean_duration_s', 0):.0f} seconds")
        kpi_text.append(f"Mean Revisits per Target: {stats.get('mean_revisits', 0):.1f}")

    _add_kpi_box(slide, kpi_text, Inches(0.5), Inches(1.5))

    # Add coverage plot if available
    if 'coverage_map' in plot_paths and plot_paths['coverage_map'].exists():
        slide.shapes.add_picture(str(plot_paths['coverage_map']),
                                 Inches(6), Inches(1.5), width=Inches(6.5))

    # Slide 3: Contact KPIs
    slide = prs.slides.add_slide(slide_layout)
    _add_slide_title(slide, "Downlink Contact Summary")

    kpi_text = []
    if contact_kpis and 'stats' in contact_kpis:
        stats = contact_kpis['stats']
        kpi_text.append(f"Valid Contacts: {stats.get('total_valid_contacts', 0)}")
        kpi_text.append(f"Success Rate: {stats.get('valid_percentage', 0):.1f}%")

    if downlink_kpis:
        kpi_text.append(f"Total Data Collected: {downlink_kpis.get('total_data_collected_gb', 0):.1f} GB")
        kpi_text.append(f"Total Data Downlinked: {downlink_kpis.get('total_data_downlinked_gb', 0):.1f} GB")
        kpi_text.append(f"Downlink Efficiency: {downlink_kpis.get('downlink_efficiency_pct', 0):.1f}%")

    _add_kpi_box(slide, kpi_text, Inches(0.5), Inches(1.5))

    if 'contact_validity' in plot_paths and plot_paths['contact_validity'].exists():
        slide.shapes.add_picture(str(plot_paths['contact_validity']),
                                 Inches(5), Inches(1.5), width=Inches(7.5))

    # Slide 4: TTNC
    slide = prs.slides.add_slide(slide_layout)
    _add_slide_title(slide, "Time to Next Ka Contact (TTNC)")

    kpi_text = []
    if not ttnc_df.empty:
        valid_ttnc = ttnc_df['TTNC_Ka_minutes'].dropna()
        if not valid_ttnc.empty:
            kpi_text.append(f"Median TTNC: {valid_ttnc.median():.1f} minutes")
            kpi_text.append(f"P90 TTNC: {valid_ttnc.quantile(0.90):.1f} minutes")
            kpi_text.append(f"P95 TTNC: {valid_ttnc.quantile(0.95):.1f} minutes")
            kpi_text.append(f"Missing Data: {ttnc_df['TTNC_Ka_minutes'].isna().sum()}/{len(ttnc_df)}")

    _add_kpi_box(slide, kpi_text, Inches(0.5), Inches(1.5))

    if 'ttnc_distribution' in plot_paths and plot_paths['ttnc_distribution'].exists():
        slide.shapes.add_picture(str(plot_paths['ttnc_distribution']),
                                 Inches(5), Inches(1.5), width=Inches(7.5))

    # Slide 5: Backlog
    slide = prs.slides.add_slide(slide_layout)
    _add_slide_title(slide, "Data Backlog Analysis")

    kpi_text = []
    if downlink_kpis:
        kpi_text.append(f"Peak Backlog: {downlink_kpis.get('peak_backlog_gb', 0):.2f} GB")
        kpi_text.append(f"Mean Backlog: {downlink_kpis.get('mean_backlog_gb', 0):.2f} GB")
        kpi_text.append(f"Final Backlog: {downlink_kpis.get('final_backlog_gb', 0):.2f} GB")

    _add_kpi_box(slide, kpi_text, Inches(0.5), Inches(1.5))

    if 'backlog_timeseries' in plot_paths and plot_paths['backlog_timeseries'].exists():
        slide.shapes.add_picture(str(plot_paths['backlog_timeseries']),
                                 Inches(4.5), Inches(1.5), width=Inches(8))

    # Slide 6: Ground Tracks (if available)
    if 'ground_tracks' in plot_paths and plot_paths['ground_tracks'].exists():
        slide = prs.slides.add_slide(slide_layout)
        _add_slide_title(slide, "Satellite Ground Tracks")
        slide.shapes.add_picture(str(plot_paths['ground_tracks']),
                                 Inches(1), Inches(1.5), width=Inches(11))

    # Slide 7: Comm Cones (if available)
    if 'comm_cones' in plot_paths and plot_paths['comm_cones'].exists():
        slide = prs.slides.add_slide(slide_layout)
        _add_slide_title(slide, "Ground Station Communication Cones")
        slide.shapes.add_picture(str(plot_paths['comm_cones']),
                                 Inches(1), Inches(1.5), width=Inches(11))

    # Slide 8: Optimization Results (if available)
    if optimization_results:
        slide = prs.slides.add_slide(slide_layout)
        _add_slide_title(slide, "RGT-Min Optimization Results")

        opt_text = []
        for provider, result in optimization_results.items():
            opt_text.append(f"{provider}:")
            opt_text.append(f"  TT&C: {len(result.ttc_stations)} stations, {result.ttc_coverage_percent:.1f}% coverage")
            if result.ka_stations:
                opt_text.append(f"  Ka: {len(result.ka_stations)} stations, {result.ka_delivery_percent:.1f}% delivery")
            status = "SATISFIED" if (result.ttc_satisfied and result.ka_satisfied) else "NOT MET"
            opt_text.append(f"  SLA Status: {status}")
            opt_text.append("")

        _add_kpi_box(slide, opt_text[:12], Inches(0.5), Inches(1.5))  # Limit to fit

    prs.save(output_path)
    print(f"PowerPoint report saved to: {output_path}")
    return output_path


def _add_slide_title(slide, title_text: str) -> None:
    """Add title to a slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True


def _add_kpi_box(slide, kpi_text: List[str], left: float, top: float) -> None:
    """Add KPI text box to a slide."""
    from pptx.util import Inches, Pt

    kpi_box = slide.shapes.add_textbox(left, top, Inches(4.5), Inches(4))
    kpi_frame = kpi_box.text_frame
    kpi_frame.word_wrap = True

    for i, text in enumerate(kpi_text):
        if i == 0:
            para = kpi_frame.paragraphs[0]
        else:
            para = kpi_frame.add_paragraph()
        para.text = f"• {text}"
        para.font.size = Pt(18)
        para.space_after = Pt(12)
