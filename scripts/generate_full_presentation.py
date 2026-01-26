#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation for satellite coverage analysis.

Supports parallel plot generation for improved performance.
"""

import sys
import random
import time
from pathlib import Path
from datetime import timedelta
from typing import Dict, List, Tuple, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
from pyproj import Geod

import numpy as np
import pandas as pd
import geopandas as gpd
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend for thread safety
import matplotlib.pyplot as plt
import cartopy.crs as ccrs
import cartopy.feature as cfeature
from matplotlib.patches import Patch
from matplotlib.lines import Line2D
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.parent))

from src.vleo_eo.orbits import create_tle_data, calculate_ground_track, ltdn_to_raan, calculate_sso_inclination
from src.vleo_eo.contacts import calculate_comm_range_km, calculate_ground_distance_km
from src.vleo_eo.config import load_config
from src.vleo_eo.constants import EARTH_RADIUS_KM, EARTH_MU_KM3_S2


def create_comm_cone_geodesic(lon: float, lat: float, min_elevation_deg: float, sat_altitude_km: float):
    """Create comm cone using geodesic geometry.

    Handles antimeridian crossing for stations near ±180° longitude by
    normalizing longitudes to center around the station.
    """
    from shapely.geometry import Polygon

    earth_radius_km = 6378.137
    elevation_rad = np.radians(min_elevation_deg)
    central_angle = np.arccos(
        earth_radius_km * np.cos(elevation_rad) / (earth_radius_km + sat_altitude_km)
    ) - elevation_rad
    max_ground_range_m = earth_radius_km * central_angle * 1000

    # Use geodesic forward calculation to create circle points
    geod = Geod(ellps='WGS84')
    circle_lons = []
    circle_lats = []
    for azimuth in np.linspace(0, 360, 73):  # 5-degree increments
        dest_lon, dest_lat, _ = geod.fwd(lon, lat, azimuth, max_ground_range_m)
        circle_lons.append(dest_lon)
        circle_lats.append(dest_lat)

    # Normalize longitudes to avoid antimeridian discontinuity
    # Shift all longitudes to be within 180° of the center longitude
    normalized_lons = []
    for clon in circle_lons:
        # Calculate offset from center
        offset = clon - lon
        # Normalize to [-180, 180] range relative to center
        while offset > 180:
            offset -= 360
        while offset < -180:
            offset += 360
        normalized_lons.append(lon + offset)

    # Create polygon from normalized points
    coords = list(zip(normalized_lons, circle_lats))
    poly = Polygon(coords)

    if not poly.is_valid:
        poly = poly.buffer(0)

    return poly


def calculate_coverage_radius(sat_altitude_km: float, min_elevation_deg: float) -> float:
    """Calculate ground coverage radius in km for given satellite altitude and minimum elevation."""
    earth_radius_km = 6378.137
    elevation_rad = np.radians(min_elevation_deg)
    central_angle = np.arccos(
        earth_radius_km * np.cos(elevation_rad) / (earth_radius_km + sat_altitude_km)
    ) - elevation_rad
    return earth_radius_km * central_angle


def get_elevation_for_band(config, is_ka_band: bool) -> float:
    """Get the appropriate elevation mask from config based on band.

    Parameters
    ----------
    config : AnalysisConfig
        Configuration object with raw_config dict.
    is_ka_band : bool
        True for Ka-band (payload downlink), False for S/X-band (TT&C).

    Returns
    -------
    float
        Elevation mask in degrees.
    """
    if is_ka_band:
        return config.raw_config.get('min_elevation_ka_deg', 10.0)
    else:
        return config.raw_config.get('min_elevation_sx_deg', 5.0)


def generate_coverage_circle(lat: float, lon: float, radius_km: float, num_points: int = 72) -> Tuple[List[float], List[float]]:
    """Generate circle coordinates around a point with given radius."""
    geod = Geod(ellps='WGS84')
    lons = []
    lats = []
    for azimuth in np.linspace(0, 360, num_points):
        dest_lon, dest_lat, _ = geod.fwd(lon, lat, azimuth, radius_km * 1000)
        lons.append(dest_lon)
        lats.append(dest_lat)
    lons.append(lons[0])  # Close the circle
    lats.append(lats[0])
    return lons, lats


def setup_map_axes(ax, land_color='#fafaf9', ocean_color='#d4dbdc', border_color='#ebd6d9'):
    """Setup map with standard styling."""
    land = cfeature.NaturalEarthFeature('physical', 'land', '50m', facecolor=land_color)
    ocean = cfeature.NaturalEarthFeature('physical', 'ocean', '50m', facecolor=ocean_color)
    ax.add_feature(ocean, zorder=0)
    ax.add_feature(land, zorder=1)
    ax.add_feature(cfeature.BORDERS, linewidth=0.5, edgecolor=border_color, zorder=2)
    ax.add_feature(cfeature.COASTLINE, linewidth=0.5, edgecolor='gray', zorder=2)
    ax.gridlines(draw_labels=True, alpha=0.3)


def plot_ground_stations_and_cones(ax, ground_stations: List, sat_alt_km: float):
    """Plot ground stations with comm cones."""
    for gs in ground_stations:
        is_ka = gs.get('ka_band', False)
        marker_color = 'green' if is_ka else 'black'
        cone_color = 'green' if is_ka else 'black'
        cone_elev = 10.0 if is_ka else 5.0

        ax.plot(gs['lon'], gs['lat'], '^', color=marker_color, markersize=6,
                transform=ccrs.PlateCarree(), zorder=8)

        try:
            comm_cone = create_comm_cone_geodesic(gs['lon'], gs['lat'], cone_elev, sat_alt_km)
            if comm_cone is not None and comm_cone.is_valid:
                ax.add_geometries([comm_cone], crs=ccrs.PlateCarree(),
                                facecolor='none', edgecolor=cone_color, alpha=0.5,
                                linewidth=0.8, zorder=2)
        except Exception:
            pass


def add_slide1_text(prs, config, tle_data: Dict, targets_gdf: gpd.GeoDataFrame):
    """Add mission parameters as text slide directly to PowerPoint."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Mission Configuration"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Get config data
    orbit_config = config.raw_config.get('orbit', {})
    ltdn = orbit_config.get('ltdn_hours')
    ltdn_str = f'{int(ltdn)}:{int((ltdn%1)*60):02d}' if ltdn else 'N/A'
    ka_count = sum(1 for gs in config.ground_stations if gs.ka_capable)
    im = config.imaging_modes[0] if config.imaging_modes else None

    # Get unique providers
    providers = set()
    for gs in config.ground_stations:
        if hasattr(gs, 'provider') and gs.provider:
            providers.add(gs.provider)
    providers_str = ', '.join(sorted(providers)) if providers else 'N/A'

    # Get target region from geojson name or infer from config
    num_targets = len(targets_gdf) if targets_gdf is not None else 0
    targets_config = config.raw_config.get('targets', {})
    geojson_path = targets_config.get('geojson_path', '')
    # Extract region from filename like "target_aoi_deck_apac.geojson"
    if 'apac' in geojson_path.lower():
        target_region = 'APAC'
    elif 'europe' in geojson_path.lower():
        target_region = 'Europe'
    elif 'middleeast' in geojson_path.lower() or 'middle_east' in geojson_path.lower():
        target_region = 'Middle East'
    else:
        target_region = 'Global'

    # Column 1: Analysis Period + Constellation + Ground Network
    left_text = []
    left_text.append("ANALYSIS PERIOD")
    left_text.append(f"  Start Date: {config.start_date}")
    left_text.append(f"  Duration: {config.duration_days} days")
    left_text.append(f"  Time Step: {config.time_step_s} s")
    left_text.append(f"  End Date: {config.end_datetime.strftime('%Y-%m-%d')}")
    left_text.append("")
    left_text.append("CONSTELLATION")
    left_text.append(f"  Satellites: {len(config.satellites)}")
    left_text.append(f"  Altitude: {config.satellites[0].altitude_km} km")
    left_text.append(f"  Inclination: {config.satellites[0].inclination_deg}")
    left_text.append(f"  Orbit Type: {'SSO' if orbit_config.get('auto_sso') else 'Custom'}")
    left_text.append(f"  LTDN: {ltdn_str}")
    left_text.append("")
    left_text.append("GROUND NETWORK")
    left_text.append(f"  Provider(s): {providers_str}")
    left_text.append(f"  Total Stations: {len(config.ground_stations)}")
    left_text.append(f"  Ka-band: {ka_count}")
    left_text.append(f"  TT&C Only: {len(config.ground_stations) - ka_count}")
    left_text.append("")
    left_text.append("TARGET DECK")
    left_text.append(f"  {num_targets} targets in {target_region}")

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(4), Inches(6))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        if line and not line.startswith(' '):
            p.font.bold = True

    # Column 2: TLE Data + Imaging
    # Show up to 3 random satellites
    sat_ids = list(tle_data.keys())
    if len(sat_ids) > 3:
        import random
        sat_ids_to_show = random.sample(sat_ids, 3)
    else:
        sat_ids_to_show = sat_ids

    mid_text = []
    mid_text.append("TLE DATA")
    for sat_id in sat_ids_to_show:
        tle = tle_data[sat_id]
        mid_text.append(f"  Satellite {sat_id}")
        mid_text.append(f"  {tle.line1}")
        mid_text.append(f"  {tle.line2}")
        mid_text.append("")
    if len(sat_ids) > 3:
        mid_text.append(f"  ... and {len(sat_ids) - 3} more satellites")
        mid_text.append("")
    mid_text.append("IMAGING")
    mid_text.append(f"  Mode: Frame - Single Target")
    max_ona = config.imaging_modes[0].off_nadir_max_deg if config.imaging_modes else 60.0
    mid_text.append(f"  Max ONA: {max_ona:.0f} deg")

    mid_box = slide.shapes.add_textbox(Inches(5.0), Inches(1.0), Inches(7.5), Inches(6))
    tf = mid_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(mid_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        if line and not line.startswith(' '):
            p.font.bold = True

    return slide


def generate_slide2_gs_map(config, ground_stations: List, sat_alt_km: float,
                          targets_gdf: gpd.GeoDataFrame, output_path: Path,
                          excel_path: Path = None):
    """Generate ground stations, comm cones, and targets map."""
    fig = plt.figure(figsize=(16, 9))
    ax = plt.axes(projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_global()

    plot_ground_stations_and_cones(ax, ground_stations, sat_alt_km)

    # Add station labels
    for gs in ground_stations:
        is_ka = gs.get('ka_band', False)
        label = gs['name'].replace('KSAT-', '').replace('-', '\n')
        ax.text(gs['lon'], gs['lat'] + 3, label, transform=ccrs.PlateCarree(),
                fontsize=6, ha='center', va='bottom',
                color='green' if is_ka else 'black')

    # Plot targets as blue X markers
    if targets_gdf is not None and len(targets_gdf) > 0:
        for idx, row in targets_gdf.iterrows():
            lat = row.geometry.y
            lon = row.geometry.x
            ax.plot(lon, lat, 'x', color='blue', markersize=8, markeredgewidth=2,
                    transform=ccrs.PlateCarree(), zorder=10)

    # Legend
    legend_elements = [
        Line2D([0], [0], marker='^', color='green', markersize=8, linestyle='None', label='Ka Station'),
        Patch(facecolor='none', edgecolor='green', label='Ka Comm Cone (10°)'),
        Line2D([0], [0], marker='^', color='black', markersize=8, linestyle='None', label='TT&C Station'),
        Patch(facecolor='none', edgecolor='black', label='TT&C Comm Cone (5°)'),
        Line2D([0], [0], marker='x', color='blue', markersize=8, linestyle='None',
               markeredgewidth=2, label=f'Targets ({len(targets_gdf)})'),
    ]
    ax.legend(handles=legend_elements, loc='lower left', fontsize=9)
    ax.set_title('Ground Station Network, Communication Cones, and Target Locations', fontsize=14, fontweight='bold')

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


def generate_slide3_single_orbit(tle_data: Dict, config, output_path: Path,
                                  ground_stations: List[Dict] = None, sat_alt_km: float = None):
    """Generate single orbit ground track with ground stations and contacts."""
    fig = plt.figure(figsize=(16, 10))
    ax = plt.axes([0.05, 0.15, 0.90, 0.80], projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_global()

    track_color = '#1f77b4'

    # Calculate orbital period from altitude
    if sat_alt_km is None:
        sat_alt_km = config.satellites[0].altitude_km
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + sat_alt_km
    orbital_period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60

    # Generate ground track for first satellite
    sat_id = list(tle_data.keys())[0]
    ground_track = calculate_ground_track(
        tle_data, sat_id, config.start_datetime,
        duration=timedelta(minutes=orbital_period_min), step_s=10.0
    )

    if not ground_track:
        plt.close()
        return

    lons, lats = zip(*ground_track)

    # Plot ground track
    ax.plot(lons, lats, color=track_color, linewidth=1.5, transform=ccrs.PlateCarree(),
            label='Ground Track')

    # Start marker (circle)
    ax.plot(lons[0], lats[0], 'o', color=track_color, markersize=10,
            markeredgecolor='black', markeredgewidth=1,
            transform=ccrs.PlateCarree(), zorder=10)

    # End marker (square)
    ax.plot(lons[-1], lats[-1], 's', color=track_color, markersize=8,
            markeredgecolor='black', markeredgewidth=1,
            transform=ccrs.PlateCarree(), zorder=10)

    # Add direction arrows along track
    arrow_indices = [len(ground_track) // 4, len(ground_track) // 2, 3 * len(ground_track) // 4]
    for idx in arrow_indices:
        if idx + 1 < len(ground_track):
            dx = lons[idx + 1] - lons[idx]
            dy = lats[idx + 1] - lats[idx]
            ax.annotate('', xy=(lons[idx] + dx*0.5, lats[idx] + dy*0.5),
                       xytext=(lons[idx], lats[idx]),
                       arrowprops=dict(arrowstyle='->', color=track_color, lw=1.5),
                       transform=ccrs.PlateCarree())

    # Track contacts for summary
    ka_contacts = []
    ttc_contacts = []

    # Plot ground stations and calculate contacts
    if ground_stations:
        for gs in ground_stations:
            gs_lat = gs['lat']
            gs_lon = gs['lon']
            gs_name = gs['name']
            is_ka = gs.get('ka_band', False)
            # Use per-band elevation mask from config
            min_elev = get_elevation_for_band(config, is_ka)

            # Colors: green for Ka, black for TT&C
            gs_color = 'green' if is_ka else 'black'

            # Plot ground station marker
            ax.plot(gs_lon, gs_lat, '^', color=gs_color, markersize=7,
                    markeredgecolor='black', markeredgewidth=0.5,
                    transform=ccrs.PlateCarree(), zorder=15)

            # Calculate and plot coverage circle (solid line)
            coverage_radius_km = calculate_coverage_radius(sat_alt_km, min_elev)
            circle_lons, circle_lats = generate_coverage_circle(gs_lat, gs_lon, coverage_radius_km)
            ax.plot(circle_lons, circle_lats, color=gs_color, linewidth=1.0, alpha=0.6,
                    linestyle='-', transform=ccrs.PlateCarree())

            # Find contacts during this orbit
            coverage_radius_deg = coverage_radius_km / 111.0
            in_contact = False
            contact_start_idx = None

            for i, (lon, lat) in enumerate(ground_track):
                dist = np.sqrt((lat - gs_lat)**2 + ((lon - gs_lon) * np.cos(np.radians(gs_lat)))**2)

                if dist < coverage_radius_deg:
                    if not in_contact:
                        in_contact = True
                        contact_start_idx = i
                    contact_end_idx = i
                else:
                    if in_contact:
                        # Contact ended - record and plot
                        duration_sec = (contact_end_idx - contact_start_idx) * 10
                        if duration_sec >= 30:
                            start_lon, start_lat = ground_track[contact_start_idx]
                            end_lon, end_lat = ground_track[contact_end_idx]

                            # Plot contact segment
                            contact_lons = [ground_track[j][0] for j in range(contact_start_idx, contact_end_idx + 1)]
                            contact_lats = [ground_track[j][1] for j in range(contact_start_idx, contact_end_idx + 1)]
                            ax.plot(contact_lons, contact_lats, color=gs_color, linewidth=4, alpha=0.8,
                                    transform=ccrs.PlateCarree(), zorder=8)

                            # Add symbol at start of contact
                            ax.plot(start_lon, start_lat, 'D', color=gs_color, markersize=6,
                                    markeredgecolor='white', markeredgewidth=1,
                                    transform=ccrs.PlateCarree(), zorder=12)

                            # Add label - black text, aligned left at start
                            short_name = gs_name.replace('Viasat-', '').replace('KSAT-', '')
                            label_text = f"{short_name} ({duration_sec//60}m{duration_sec%60}s)"
                            ax.annotate(label_text,
                                       xy=(start_lon, start_lat),
                                       xytext=(8, 3), textcoords='offset points',
                                       fontsize=7, ha='left', va='bottom',
                                       color='black', fontweight='bold',
                                       transform=ccrs.PlateCarree(), zorder=20)

                            # Record for summary
                            contact_info = {'station': short_name, 'duration_sec': duration_sec}
                            if is_ka:
                                ka_contacts.append(contact_info)
                            else:
                                ttc_contacts.append(contact_info)

                        in_contact = False

            # Handle contact at end of orbit
            if in_contact and contact_start_idx is not None:
                duration_sec = (contact_end_idx - contact_start_idx) * 10
                if duration_sec >= 30:
                    start_lon, start_lat = ground_track[contact_start_idx]
                    end_lon, end_lat = ground_track[contact_end_idx]

                    contact_lons = [ground_track[j][0] for j in range(contact_start_idx, contact_end_idx + 1)]
                    contact_lats = [ground_track[j][1] for j in range(contact_start_idx, contact_end_idx + 1)]
                    ax.plot(contact_lons, contact_lats, color=gs_color, linewidth=4, alpha=0.8,
                            transform=ccrs.PlateCarree(), zorder=8)

                    ax.plot(start_lon, start_lat, 'D', color=gs_color, markersize=6,
                            markeredgecolor='white', markeredgewidth=1,
                            transform=ccrs.PlateCarree(), zorder=12)

                    short_name = gs_name.replace('Viasat-', '').replace('KSAT-', '')
                    label_text = f"{short_name} ({duration_sec//60}m{duration_sec%60}s)"
                    ax.annotate(label_text,
                               xy=(start_lon, start_lat),
                               xytext=(8, 3), textcoords='offset points',
                               fontsize=7, ha='left', va='bottom',
                               color='black', fontweight='bold',
                               transform=ccrs.PlateCarree(), zorder=20)

                    contact_info = {'station': short_name, 'duration_sec': duration_sec}
                    if is_ka:
                        ka_contacts.append(contact_info)
                    else:
                        ttc_contacts.append(contact_info)

    # Legend with start/end markers and contact types
    legend_elements = [
        Line2D([0], [0], color=track_color, linewidth=1.5, label='Ground Track'),
        Line2D([0], [0], marker='o', color=track_color, markersize=8, linestyle='None',
               markeredgecolor='black', label='Orbit Start'),
        Line2D([0], [0], marker='s', color=track_color, markersize=7, linestyle='None',
               markeredgecolor='black', label='Orbit End'),
        Line2D([0], [0], marker='^', color='green', markersize=8, linestyle='None',
               markeredgecolor='black', label='Ka Station'),
        Line2D([0], [0], color='green', linewidth=1.0, label='Ka Comm Cone'),
        Line2D([0], [0], color='green', linewidth=4, alpha=0.8, label='Ka Contact'),
        Line2D([0], [0], marker='^', color='black', markersize=8, linestyle='None',
               markeredgecolor='black', label='TT&C Station'),
        Line2D([0], [0], color='black', linewidth=1.0, label='TT&C Comm Cone'),
        Line2D([0], [0], color='black', linewidth=4, alpha=0.8, label='TT&C Contact'),
        Line2D([0], [0], marker='D', color='gray', markersize=6, linestyle='None',
               markeredgecolor='white', label='Contact Start'),
    ]
    ax.legend(handles=legend_elements, loc='upper right', fontsize=8)
    ax.set_title(f'Single Orbit Ground Track (~{orbital_period_min:.1f} minutes)', fontsize=14, fontweight='bold')

    # Add summary text below the plot
    total_ka = len(ka_contacts)
    total_ttc = len(ttc_contacts)
    avg_ka_dur = np.mean([c['duration_sec'] for c in ka_contacts]) if ka_contacts else 0
    avg_ttc_dur = np.mean([c['duration_sec'] for c in ttc_contacts]) if ttc_contacts else 0

    summary_text = f"Contact Summary:  "
    summary_text += f"Ka-band: {total_ka} contacts"
    if total_ka > 0:
        summary_text += f" (avg {avg_ka_dur/60:.1f} min)"
    summary_text += f"  |  TT&C: {total_ttc} contacts"
    if total_ttc > 0:
        summary_text += f" (avg {avg_ttc_dur/60:.1f} min)"

    fig.text(0.5, 0.05, summary_text, ha='center', va='center', fontsize=11, fontweight='bold')

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


def generate_slide_1day_tracks(tle_data: Dict, config, ground_stations: List[Dict], sat_alt_km: float, output_path: Path):
    """Generate all ground tracks for 1 day (24 hours) with ground stations and contacts."""
    fig = plt.figure(figsize=(16, 10))
    ax = plt.axes([0.05, 0.15, 0.90, 0.80], projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_global()

    track_color = '#1f77b4'

    # Calculate orbital period from altitude
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + sat_alt_km
    orbital_period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60
    orbits_per_day = 24 * 60 / orbital_period_min

    # Generate 24-hour ground track with finer time step for contact detection
    all_track_points = []  # (time, lon, lat) tuples
    for i, (sat_id, tle) in enumerate(tle_data.items()):
        # Generate 24-hour ground track with 10-second steps for accurate contact detection
        ground_track = calculate_ground_track(
            tle_data, sat_id, config.start_datetime,
            duration=timedelta(hours=24), step_s=10.0
        )
        if not ground_track:
            continue

        lons, lats = zip(*ground_track)

        # Store track points with timestamps
        for j, (lon, lat) in enumerate(ground_track):
            t = config.start_datetime + timedelta(seconds=j * 10)
            all_track_points.append((t, lon, lat, sat_id))

        # Plot the full day track
        ax.plot(lons, lats, color=track_color, linewidth=0.8, alpha=0.6,
                transform=ccrs.PlateCarree(), label=f'Satellite {sat_id}')

        # Mark the start
        ax.plot(lons[0], lats[0], 'o', color=track_color, markersize=8,
                markeredgecolor='black', markeredgewidth=1,
                transform=ccrs.PlateCarree(), zorder=10)

    # Track contacts for summary
    ka_contacts = []
    ttc_contacts = []

    # Plot ground stations with coverage circles (Ka=green, TT&C=black)
    for gs in ground_stations:
        gs_lat = gs['lat']
        gs_lon = gs['lon']
        is_ka = gs.get('ka_band', False)
        # Use per-band elevation mask from config
        min_elev = get_elevation_for_band(config, is_ka)

        # Colors: green for Ka, black for TT&C
        gs_color = 'green' if is_ka else 'black'

        # Plot ground station marker
        ax.plot(gs_lon, gs_lat, '^', color=gs_color, markersize=7,
                markeredgecolor='black', markeredgewidth=0.5,
                transform=ccrs.PlateCarree(), zorder=15)

        # Calculate and plot coverage circle (solid line)
        coverage_radius_km = calculate_coverage_radius(sat_alt_km, min_elev)
        circle_lons, circle_lats = generate_coverage_circle(gs_lat, gs_lon, coverage_radius_km)
        ax.plot(circle_lons, circle_lats, color=gs_color, linewidth=1.0, alpha=0.6,
                linestyle='-', transform=ccrs.PlateCarree())

    # Calculate contacts for each ground station
    contact_annotations = []
    for gs in ground_stations:
        gs_lat = gs['lat']
        gs_lon = gs['lon']
        gs_name = gs['name']
        is_ka = gs.get('ka_band', False)
        # Use per-band elevation mask from config
        min_elev = get_elevation_for_band(config, is_ka)

        # Colors: green for Ka, black for TT&C
        gs_color = 'green' if is_ka else 'black'

        coverage_radius_km = calculate_coverage_radius(sat_alt_km, min_elev)
        coverage_radius_deg = coverage_radius_km / 111.0  # Approximate conversion

        # Find contacts by checking if satellite is within coverage
        in_contact = False
        contact_start = None
        contact_start_pos = None

        for t, lon, lat, sat_id in all_track_points:
            # Calculate distance to ground station
            dist = np.sqrt((lat - gs_lat)**2 + ((lon - gs_lon) * np.cos(np.radians(gs_lat)))**2)

            if dist < coverage_radius_deg:
                if not in_contact:
                    # Contact starts
                    in_contact = True
                    contact_start = t
                    contact_start_pos = (lon, lat)
                contact_end = t
                contact_end_pos = (lon, lat)
            else:
                if in_contact:
                    # Contact ends - record it
                    duration = (contact_end - contact_start).total_seconds()
                    if duration >= 30:  # Only show contacts >= 30 seconds
                        contact_annotations.append({
                            'gs_name': gs_name.replace('Viasat-', '').replace('KSAT-', ''),
                            'aos': contact_start,
                            'los': contact_end,
                            'duration': duration,
                            'color': gs_color,
                            'is_ka': is_ka,
                            'start_pos': contact_start_pos,
                            'end_pos': contact_end_pos,
                        })

                        # Highlight contact segment on track
                        ax.plot([contact_start_pos[0], contact_end_pos[0]],
                               [contact_start_pos[1], contact_end_pos[1]],
                               color=gs_color, linewidth=4, alpha=0.8,
                               transform=ccrs.PlateCarree(), zorder=8)

                        # Add diamond marker at start of contact
                        ax.plot(contact_start_pos[0], contact_start_pos[1], 'D', color=gs_color,
                                markersize=5, markeredgecolor='white', markeredgewidth=0.5,
                                transform=ccrs.PlateCarree(), zorder=12)

                        # Record for summary
                        if is_ka:
                            ka_contacts.append({'duration': duration})
                        else:
                            ttc_contacts.append({'duration': duration})

                    in_contact = False
                    contact_start = None

        # Handle contact that extends to end of period
        if in_contact and contact_start is not None:
            duration = (contact_end - contact_start).total_seconds()
            if duration >= 30:
                contact_annotations.append({
                    'gs_name': gs_name.replace('Viasat-', '').replace('KSAT-', ''),
                    'aos': contact_start,
                    'los': contact_end,
                    'duration': duration,
                    'color': gs_color,
                    'is_ka': is_ka,
                    'start_pos': contact_start_pos,
                    'end_pos': contact_end_pos,
                })
                ax.plot([contact_start_pos[0], contact_end_pos[0]],
                       [contact_start_pos[1], contact_end_pos[1]],
                       color=gs_color, linewidth=4, alpha=0.8,
                       transform=ccrs.PlateCarree(), zorder=8)

                ax.plot(contact_start_pos[0], contact_start_pos[1], 'D', color=gs_color,
                        markersize=5, markeredgecolor='white', markeredgewidth=0.5,
                        transform=ccrs.PlateCarree(), zorder=12)

                if is_ka:
                    ka_contacts.append({'duration': duration})
                else:
                    ttc_contacts.append({'duration': duration})

    # Calculate revolution number for each contact based on orbital period
    total_revs = int(np.ceil(orbits_per_day))

    # Add contact annotations with rev number and duration only
    for i, contact in enumerate(contact_annotations):
        dur_min = int(contact['duration'] // 60)
        dur_sec = int(contact['duration'] % 60)

        # Calculate which rev this contact is in (1-indexed)
        time_since_start = (contact['aos'] - config.start_datetime).total_seconds()
        rev_num = int(time_since_start / (orbital_period_min * 60)) + 1

        # Position annotation at start of contact, offset to the side
        start_lon, start_lat = contact['start_pos']
        offset_x = 15 if (i % 2 == 0) else -15
        offset_y = 5 if (i % 3 == 0) else -5

        annotation_text = f"Rev {rev_num}/{total_revs} • {dur_min}m{dur_sec}s"

        ax.annotate(annotation_text,
                   xy=(start_lon, start_lat),
                   xytext=(offset_x, offset_y), textcoords='offset points',
                   fontsize=7, ha='left' if offset_x > 0 else 'right', va='center',
                   color='black', fontweight='bold',
                   transform=ccrs.PlateCarree(), zorder=20)

    # Legend with Ka/TT&C distinction
    legend_elements = [
        Line2D([0], [0], color=track_color, linewidth=1.5, label='Ground Track'),
        Line2D([0], [0], marker='o', color=track_color, markersize=8, linestyle='None',
               markeredgecolor='black', label='Day Start'),
        Line2D([0], [0], marker='^', color='green', markersize=8, linestyle='None',
               markeredgecolor='black', label='Ka Station'),
        Line2D([0], [0], color='green', linewidth=1.0, label='Ka Comm Cone'),
        Line2D([0], [0], color='green', linewidth=4, alpha=0.8, label='Ka Contact'),
        Line2D([0], [0], marker='^', color='black', markersize=8, linestyle='None',
               markeredgecolor='black', label='TT&C Station'),
        Line2D([0], [0], color='black', linewidth=1.0, label='TT&C Comm Cone'),
        Line2D([0], [0], color='black', linewidth=4, alpha=0.8, label='TT&C Contact'),
        Line2D([0], [0], marker='D', color='gray', markersize=6, linestyle='None',
               markeredgecolor='white', label='Contact Start'),
    ]
    ax.legend(handles=legend_elements, loc='upper right', fontsize=8)
    ax.set_title(f'24-Hour Ground Track with Contacts (~{orbits_per_day:.1f} orbits/day)',
                 fontsize=14, fontweight='bold')

    # Add summary text below the plot
    total_ka = len(ka_contacts)
    total_ttc = len(ttc_contacts)
    avg_ka_dur = np.mean([c['duration'] for c in ka_contacts]) if ka_contacts else 0
    avg_ttc_dur = np.mean([c['duration'] for c in ttc_contacts]) if ttc_contacts else 0

    summary_text = f"Contact Summary:  "
    summary_text += f"Ka-band: {total_ka} contacts"
    if total_ka > 0:
        summary_text += f" (avg {avg_ka_dur/60:.1f} min)"
    summary_text += f"  |  TT&C: {total_ttc} contacts"
    if total_ttc > 0:
        summary_text += f" (avg {avg_ttc_dur/60:.1f} min)"

    fig.text(0.5, 0.05, summary_text, ha='center', va='center', fontsize=11, fontweight='bold')

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


def generate_slide_gt_walking(tle_data: Dict, config, output_path: Path):
    """Generate ground track walking visualization - first orbit of each day over mission duration."""
    fig = plt.figure(figsize=(16, 9))
    ax = plt.axes(projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_global()

    # Create colormap for day progression
    cmap = plt.cm.viridis
    num_days = min(config.duration_days, 30)  # Cap at 30 days for clarity

    # Calculate orbital period
    sat_alt_km = config.satellites[0].altitude_km
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + sat_alt_km
    orbital_period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60

    # Use first satellite
    sat_id = list(tle_data.keys())[0]

    for day in range(num_days):
        # Calculate start time for this day's first orbit
        day_start = config.start_datetime + timedelta(days=day)

        # Generate single orbit ground track
        ground_track = calculate_ground_track(
            tle_data, sat_id, day_start,
            duration=timedelta(minutes=orbital_period_min), step_s=30.0
        )
        if not ground_track:
            continue

        lons, lats = zip(*ground_track)
        color = cmap(day / num_days)

        # Plot the orbit
        ax.plot(lons, lats, color=color, linewidth=1.2, alpha=0.8,
                transform=ccrs.PlateCarree())

        # Mark the start of each orbit
        ax.plot(lons[0], lats[0], 'o', color=color, markersize=5,
                markeredgecolor='black', markeredgewidth=0.5,
                transform=ccrs.PlateCarree(), zorder=10)

    # Add colorbar to show day progression
    sm = plt.cm.ScalarMappable(cmap=cmap, norm=plt.Normalize(vmin=1, vmax=num_days))
    sm.set_array([])
    cbar = plt.colorbar(sm, ax=ax, orientation='vertical', shrink=0.6, pad=0.02)
    cbar.set_label('Day of Mission', fontsize=10)

    ax.set_title(f'Ground Track Walking Over {num_days} Days\n(First orbit of each day showing longitude drift)',
                 fontsize=14, fontweight='bold')

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


def generate_slide_ttc_coverage_analysis(tle_data: Dict, config, ground_stations: List[Dict],
                                          sat_alt_km: float, output_path: Path):
    """Generate TT&C coverage analysis slide.

    Analyzes TT&C pass availability over the mission duration.
    Note: Ka-band stations are tri-band, so all stations can support TT&C.
    Goal: At least 1 TT&C pass per revolution.
    """
    fig = plt.figure(figsize=(16, 10))

    # Calculate orbital parameters
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + sat_alt_km
    orbital_period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60
    total_revs = int(config.duration_days * 24 * 60 / orbital_period_min)

    # Use first satellite
    sat_id = list(tle_data.keys())[0]

    # Generate full mission ground track with coarser step for efficiency
    print("    Calculating TT&C contacts over full mission...")
    ground_track = calculate_ground_track(
        tle_data, sat_id, config.start_datetime,
        duration=timedelta(days=config.duration_days), step_s=30.0
    )

    if not ground_track:
        plt.close()
        return None

    # Build time series of track points
    track_times = [config.start_datetime + timedelta(seconds=i * 30) for i in range(len(ground_track))]

    # Calculate contacts for ALL stations (Ka stations are tri-band, so all do TT&C)
    all_contacts = []  # List of (start_time, end_time, station_name, is_ka)
    station_contact_counts = {}

    # For TT&C analysis, use S/X-band elevation mask
    min_elev_ttc = get_elevation_for_band(config, is_ka_band=False)

    for gs in ground_stations:
        gs_lat = gs['lat']
        gs_lon = gs['lon']
        gs_name = gs['name'].replace('Viasat-', '').replace('KSAT-', '')
        is_ka = gs.get('ka_band', False)

        coverage_radius_km = calculate_coverage_radius(sat_alt_km, min_elev_ttc)
        coverage_radius_deg = coverage_radius_km / 111.0

        station_contacts = 0
        in_contact = False
        contact_start_time = None

        for i, (lon, lat) in enumerate(ground_track):
            dist = np.sqrt((lat - gs_lat)**2 + ((lon - gs_lon) * np.cos(np.radians(gs_lat)))**2)

            if dist < coverage_radius_deg:
                if not in_contact:
                    in_contact = True
                    contact_start_time = track_times[i]
                contact_end_time = track_times[i]
            else:
                if in_contact:
                    duration = (contact_end_time - contact_start_time).total_seconds()
                    if duration >= 30:
                        all_contacts.append({
                            'start': contact_start_time,
                            'end': contact_end_time,
                            'station': gs_name,
                            'is_ka': is_ka,
                            'duration': duration
                        })
                        station_contacts += 1
                    in_contact = False

        # Handle contact at end
        if in_contact and contact_start_time:
            duration = (contact_end_time - contact_start_time).total_seconds()
            if duration >= 30:
                all_contacts.append({
                    'start': contact_start_time,
                    'end': contact_end_time,
                    'station': gs_name,
                    'is_ka': is_ka,
                    'duration': duration
                })
                station_contacts += 1

        station_contact_counts[gs_name] = {'count': station_contacts, 'is_ka': is_ka}

    # Sort contacts by start time
    all_contacts.sort(key=lambda x: x['start'])

    # Analyze gaps between consecutive TT&C contacts
    gaps = []
    for i in range(1, len(all_contacts)):
        gap = (all_contacts[i]['start'] - all_contacts[i-1]['end']).total_seconds() / 60  # in minutes
        gaps.append({
            'gap_min': gap,
            'after_station': all_contacts[i-1]['station'],
            'before_station': all_contacts[i]['station'],
            'time': all_contacts[i-1]['end']
        })

    # Analyze per-revolution coverage
    revs_without_contact = []
    for rev in range(total_revs):
        rev_start = config.start_datetime + timedelta(minutes=rev * orbital_period_min)
        rev_end = rev_start + timedelta(minutes=orbital_period_min)

        # Check if any contact overlaps this revolution
        has_contact = False
        for contact in all_contacts:
            if contact['start'] < rev_end and contact['end'] > rev_start:
                has_contact = True
                break

        if not has_contact:
            revs_without_contact.append(rev + 1)  # 1-indexed

    # Create subplots
    # Top left: Gap histogram
    ax1 = fig.add_axes([0.06, 0.55, 0.28, 0.38])
    if gaps:
        gap_values = [g['gap_min'] for g in gaps]
        ax1.hist(gap_values, bins=30, color='#3498db', edgecolor='black', alpha=0.7)
        ax1.axvline(orbital_period_min, color='red', linestyle='--', linewidth=2,
                    label=f'Orbital Period ({orbital_period_min:.1f} min)')
        ax1.axvline(np.max(gap_values), color='orange', linestyle='--', linewidth=2,
                    label=f'Max Gap ({np.max(gap_values):.1f} min)')
        ax1.set_xlabel('Gap Between TT&C Contacts (minutes)', fontsize=9)
        ax1.set_ylabel('Frequency', fontsize=9)
        ax1.set_title('TT&C Contact Gap Distribution', fontsize=11, fontweight='bold')
        ax1.legend(fontsize=8)
        ax1.grid(True, alpha=0.3)

    # Top right: Station contact counts (bar chart)
    ax2 = fig.add_axes([0.40, 0.55, 0.55, 0.38])
    stations = list(station_contact_counts.keys())
    counts = [station_contact_counts[s]['count'] for s in stations]
    colors = ['green' if station_contact_counts[s]['is_ka'] else 'black' for s in stations]

    bars = ax2.barh(range(len(stations)), counts, color=colors, alpha=0.7, edgecolor='black')
    ax2.set_yticks(range(len(stations)))
    ax2.set_yticklabels(stations, fontsize=8)
    ax2.set_xlabel('Number of TT&C Contacts', fontsize=9)
    ax2.set_title(f'TT&C Contacts per Station ({config.duration_days} days)', fontsize=11, fontweight='bold')
    ax2.grid(True, alpha=0.3, axis='x')

    # Add count labels on bars
    for i, (bar, count) in enumerate(zip(bars, counts)):
        ax2.text(count + 0.5, i, str(count), va='center', fontsize=8)

    # Add legend for Ka vs TT&C-only
    from matplotlib.patches import Patch
    legend_elements = [
        Patch(facecolor='green', alpha=0.7, edgecolor='black', label='Ka-band (tri-band)'),
        Patch(facecolor='black', alpha=0.7, edgecolor='black', label='TT&C only'),
    ]
    ax2.legend(handles=legend_elements, loc='lower right', fontsize=8)

    # Bottom: Summary statistics text
    ax3 = fig.add_axes([0.06, 0.08, 0.88, 0.40])
    ax3.axis('off')

    # Calculate statistics
    total_contacts = len(all_contacts)
    avg_contacts_per_day = total_contacts / config.duration_days
    avg_contacts_per_rev = total_contacts / total_revs
    max_gap = max(gap_values) if gaps else 0
    avg_gap = np.mean(gap_values) if gaps else 0
    median_gap = np.median(gap_values) if gaps else 0
    gaps_exceeding_rev = len([g for g in gaps if g['gap_min'] > orbital_period_min])
    pct_revs_without = 100 * len(revs_without_contact) / total_revs

    # Find worst gaps
    worst_gaps = sorted(gaps, key=lambda x: -x['gap_min'])[:5] if gaps else []

    summary_text = f"""TT&C COVERAGE ANALYSIS SUMMARY
{'='*60}
Note: All Ka-band stations are tri-band and support TT&C operations.

MISSION PARAMETERS
  Duration: {config.duration_days} days
  Total Revolutions: {total_revs}
  Orbital Period: {orbital_period_min:.1f} minutes
  Ground Stations: {len(ground_stations)} (Ka: {sum(1 for g in ground_stations if g.get('ka_band', False))}, TT&C-only: {sum(1 for g in ground_stations if not g.get('ka_band', False))})

CONTACT STATISTICS
  Total TT&C Contacts: {total_contacts}
  Average Contacts/Day: {avg_contacts_per_day:.1f}
  Average Contacts/Rev: {avg_contacts_per_rev:.2f}

GAP ANALYSIS
  Maximum Gap: {max_gap:.1f} minutes ({max_gap/orbital_period_min:.2f} revs)
  Average Gap: {avg_gap:.1f} minutes
  Median Gap: {median_gap:.1f} minutes
  Gaps Exceeding 1 Rev: {gaps_exceeding_rev} ({100*gaps_exceeding_rev/len(gaps):.1f}% of gaps)

REQUIREMENT: 1 TT&C PASS PER REVOLUTION
  Revolutions Without Contact: {len(revs_without_contact)} of {total_revs} ({pct_revs_without:.1f}%)
  Requirement Met: {'YES' if len(revs_without_contact) == 0 else 'NO - ' + str(len(revs_without_contact)) + ' violations'}
"""

    if worst_gaps:
        summary_text += f"\nWORST GAPS (Top 5):\n"
        for i, gap in enumerate(worst_gaps):
            gap_time = gap['time'].strftime('%Y-%m-%d %H:%M')
            summary_text += f"  {i+1}. {gap['gap_min']:.1f} min after {gap['after_station']} ({gap_time} UTC)\n"

    ax3.text(0.0, 1.0, summary_text, transform=ax3.transAxes, fontsize=9,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='white', alpha=0.9))

    fig.suptitle('TT&C Coverage Analysis - All Stations (Ka tri-band + TT&C)',
                 fontsize=14, fontweight='bold', y=0.98)

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {
        'total_contacts': total_contacts,
        'total_revs': total_revs,
        'contacts_per_day': avg_contacts_per_day,
        'contacts_per_rev': avg_contacts_per_rev,
        'revs_without_contact': len(revs_without_contact),
        'coverage_pct': 100 * (total_revs - len(revs_without_contact)) / total_revs if total_revs > 0 else 0,
        'requirement_met': len(revs_without_contact) == 0,
        'max_gap_min': max_gap,
        'avg_gap_min': avg_gap,
        'median_gap_min': median_gap,
        'p95_gap_min': np.percentile(gap_values, 95) if gaps else 0,
        'station_contacts': station_contact_counts,
        'gaps': gaps,  # Raw gap data with timestamps for periodicity analysis
        'orbital_period_min': orbital_period_min,
    }


def generate_slide_ka_coverage_analysis(tle_data: Dict, config, ground_stations: List[Dict],
                                         sat_alt_km: float, output_path: Path, excel_path: Path = None):
    """Generate Ka-band post-collect downlink analysis slide.

    Analyzes VIABLE Ka-band downlink opportunities after each collect:
    - Includes 5-minute on-board processing buffer before downlink can start
    - Validates contact duration is sufficient to downlink image file (size from config)
    - Shows statistics on viable vs non-viable first Ka contacts
    """
    fig = plt.figure(figsize=(16, 10))

    # Configuration parameters
    PROCESSING_BUFFER_MIN = 5.0  # Minutes of on-board processing before downlink
    IMAGE_SIZE_GB = config.raw_config.get('image_size_gb', 3.0)
    downlink_rate_mbps = config.raw_config.get('downlink_rate_mbps', 1000)

    # Calculate minimum contact duration needed to downlink image
    # 6 GB = 48 Gbit, at 800 Mbps = 60 seconds, add 20% margin for overhead
    min_contact_duration_min = (IMAGE_SIZE_GB * 8 * 1000 / downlink_rate_mbps / 60) * 1.2

    # Calculate orbital parameters
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + sat_alt_km
    orbital_period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60

    # Load access data from Excel to get collect end times
    access_df = None
    if excel_path and excel_path.exists():
        try:
            access_df = pd.read_excel(excel_path, sheet_name='Access_Windows')
        except Exception:
            pass

    if access_df is None or len(access_df) == 0:
        plt.close()
        return None

    # Build list of Ka contacts for each station over mission duration
    print("    Building Ka contact windows for full mission...")
    sat_id = list(tle_data.keys())[0]

    # Generate full mission ground track
    ground_track = calculate_ground_track(
        tle_data, sat_id, config.start_datetime,
        duration=timedelta(days=config.duration_days), step_s=30.0
    )

    if not ground_track:
        plt.close()
        return None

    track_times = [config.start_datetime + timedelta(seconds=i * 30) for i in range(len(ground_track))]

    # Filter for Ka-capable stations only
    ka_stations = [gs for gs in ground_stations if gs.get('ka_band', False)]

    if not ka_stations:
        plt.close()
        return None

    # Build contact windows for each Ka station
    ka_contact_windows = []  # List of (start_time, end_time, station_name, duration_min)

    # For Ka-band analysis, use Ka-band elevation mask
    min_elev_ka = get_elevation_for_band(config, is_ka_band=True)

    for gs in ka_stations:
        gs_lat = gs['lat']
        gs_lon = gs['lon']
        gs_name = gs['name']

        coverage_radius_km = calculate_coverage_radius(sat_alt_km, min_elev_ka)
        coverage_radius_deg = coverage_radius_km / 111.0

        in_contact = False
        contact_start_time = None
        contact_end_time = None

        for i, (lon, lat) in enumerate(ground_track):
            dist = np.sqrt((lat - gs_lat)**2 + ((lon - gs_lon) * np.cos(np.radians(gs_lat)))**2)

            if dist < coverage_radius_deg:
                if not in_contact:
                    in_contact = True
                    contact_start_time = track_times[i]
                contact_end_time = track_times[i]
            else:
                if in_contact and contact_start_time:
                    duration_min = (contact_end_time - contact_start_time).total_seconds() / 60
                    if duration_min >= 0.5:  # At least 30 seconds
                        ka_contact_windows.append({
                            'start': contact_start_time,
                            'end': contact_end_time,
                            'station': gs_name,
                            'station_clean': gs_name.replace('Viasat-', '').replace('KSAT-', ''),
                            'duration_min': duration_min
                        })
                    in_contact = False

        # Handle contact at end
        if in_contact and contact_start_time:
            duration_min = (contact_end_time - contact_start_time).total_seconds() / 60
            if duration_min >= 0.5:
                ka_contact_windows.append({
                    'start': contact_start_time,
                    'end': contact_end_time,
                    'station': gs_name,
                    'station_clean': gs_name.replace('Viasat-', '').replace('KSAT-', ''),
                    'duration_min': duration_min
                })

    # Sort contacts by start time
    ka_contact_windows.sort(key=lambda x: x['start'])

    # Analyze each collect to find viable first Ka downlink
    print("    Analyzing viable Ka downlinks for each collect...")
    collect_analysis = []

    for _, access_row in access_df.iterrows():
        # Get collect end time (use Imaging_Stop if available, otherwise estimate)
        if 'Imaging_Stop' in access_row and pd.notna(access_row['Imaging_Stop']):
            collect_end = pd.to_datetime(access_row['Imaging_Stop'])
        else:
            collect_end = pd.to_datetime(access_row['Start_Time']) + timedelta(seconds=access_row.get('Access_Duration', 30))

        if collect_end.tzinfo is None:
            collect_end = collect_end.tz_localize('UTC')

        # Processing complete time (collect end + 5 min buffer)
        processing_complete = collect_end + timedelta(minutes=PROCESSING_BUFFER_MIN)

        # Find first Ka contact that:
        # 1. Starts after processing is complete, OR
        # 2. Is still ongoing when processing completes (with enough remaining time)
        first_viable_ka = None
        first_any_ka = None  # Track first Ka even if not viable (for comparison)

        for contact in ka_contact_windows:
            contact_start = contact['start']
            contact_end = contact['end']

            # Make timezone-aware if needed
            if contact_start.tzinfo is None:
                contact_start = contact_start.replace(tzinfo=timezone.utc)
            if contact_end.tzinfo is None:
                contact_end = contact_end.replace(tzinfo=timezone.utc)

            # Skip contacts that end before collect even starts
            if contact_end < collect_end:
                continue

            # Track first Ka contact after collect (regardless of viability)
            if first_any_ka is None and contact_start > collect_end:
                time_to_ka = (contact_start - collect_end).total_seconds() / 60
                first_any_ka = {
                    'station': contact['station_clean'],
                    'time_to_ka_min': time_to_ka,
                    'contact_duration_min': contact['duration_min'],
                    'viable': False,
                    'reason': 'checking'
                }

            # Check if contact is viable for downlink
            if contact_start >= processing_complete:
                # Contact starts after processing complete - full duration available
                available_duration = contact['duration_min']
                time_to_viable = (contact_start - collect_end).total_seconds() / 60

                if available_duration >= min_contact_duration_min:
                    first_viable_ka = {
                        'station': contact['station_clean'],
                        'time_to_viable_min': time_to_viable,
                        'available_duration_min': available_duration,
                        'can_downlink_gb': (available_duration * 60) * (downlink_rate_mbps / 8) / 1000
                    }
                    break

            elif contact_end > processing_complete:
                # Contact is ongoing when processing completes - partial duration available
                available_duration = (contact_end - processing_complete).total_seconds() / 60
                time_to_viable = PROCESSING_BUFFER_MIN  # Downlink starts right when processing completes

                if available_duration >= min_contact_duration_min:
                    first_viable_ka = {
                        'station': contact['station_clean'],
                        'time_to_viable_min': time_to_viable,
                        'available_duration_min': available_duration,
                        'can_downlink_gb': (available_duration * 60) * (downlink_rate_mbps / 8) / 1000
                    }
                    break

        # Record analysis result
        collect_analysis.append({
            'access_id': access_row.get('Access_ID', 0),
            'collect_end': collect_end,
            'has_viable_ka': first_viable_ka is not None,
            'viable_ka': first_viable_ka,
            'first_any_ka': first_any_ka
        })

    # Convert to DataFrame for analysis
    analysis_df = pd.DataFrame(collect_analysis)

    # Calculate statistics
    total_collects = len(analysis_df)
    viable_collects = analysis_df['has_viable_ka'].sum()
    non_viable_collects = total_collects - viable_collects

    # Extract viable contact data
    viable_data = [c['viable_ka'] for c in collect_analysis if c['viable_ka'] is not None]
    viable_df = pd.DataFrame(viable_data) if viable_data else pd.DataFrame()

    # =========================================================================
    # Create plots
    # =========================================================================

    # Top left: Time to Viable Ka Distribution
    ax1 = fig.add_axes([0.06, 0.55, 0.28, 0.38])
    if len(viable_df) > 0:
        time_to_viable = viable_df['time_to_viable_min']
        ax1.hist(time_to_viable, bins=15, color='green', edgecolor='black', alpha=0.7)
        ax1.axvline(time_to_viable.median(), color='red', linestyle='--', linewidth=2,
                    label=f'Median: {time_to_viable.median():.1f} min')
        ax1.axvline(time_to_viable.quantile(0.95), color='orange', linestyle='--', linewidth=2,
                    label=f'P95: {time_to_viable.quantile(0.95):.1f} min')
        ax1.axvline(PROCESSING_BUFFER_MIN, color='blue', linestyle=':', linewidth=2,
                    label=f'Min (processing): {PROCESSING_BUFFER_MIN} min')
    ax1.set_xlabel('Time to Viable Ka Downlink (minutes)', fontsize=9)
    ax1.set_ylabel('Number of Collects', fontsize=9)
    ax1.set_title('Post-Collect Downlink Latency\n(includes 5-min processing buffer)', fontsize=10, fontweight='bold')
    ax1.legend(fontsize=7)
    ax1.grid(True, alpha=0.3)

    # Top middle: First Viable Ka Station Distribution (pie chart)
    ax2 = fig.add_axes([0.38, 0.55, 0.28, 0.38])
    if len(viable_df) > 0:
        station_counts = viable_df['station'].value_counts()
        top_stations = station_counts.head(8)
        colors = plt.cm.Greens(np.linspace(0.3, 0.9, len(top_stations)))

        wedges, texts, autotexts = ax2.pie(
            top_stations.values,
            labels=top_stations.index,
            autopct='%1.0f%%',
            colors=colors,
            pctdistance=0.75,
            labeldistance=1.1
        )
        for text in texts:
            text.set_fontsize(8)
        for autotext in autotexts:
            autotext.set_fontsize(8)
            autotext.set_color('white')
            autotext.set_fontweight('bold')
    ax2.set_title('First Viable Downlink Station', fontsize=10, fontweight='bold')

    # Top right: Available Duration by Station (box plot)
    ax3 = fig.add_axes([0.72, 0.55, 0.26, 0.38])
    if len(viable_df) > 0:
        station_data = []
        station_labels = []
        for station in viable_df['station'].unique():
            data = viable_df[viable_df['station'] == station]['available_duration_min'].values
            if len(data) > 0:
                station_data.append(data)
                station_labels.append(station)

        if station_data:
            bp = ax3.boxplot(station_data, vert=False, patch_artist=True)
            for patch in bp['boxes']:
                patch.set_facecolor('green')
                patch.set_alpha(0.7)
            ax3.axvline(min_contact_duration_min, color='red', linestyle='--', linewidth=2,
                        label=f'Min for 6GB: {min_contact_duration_min:.1f} min')
            ax3.set_yticklabels(station_labels, fontsize=8)
            ax3.set_xlabel('Available Contact Duration (minutes)', fontsize=9)
            ax3.set_title('Contact Duration by Station', fontsize=10, fontweight='bold')
            ax3.legend(fontsize=7, loc='lower right')
            ax3.grid(True, alpha=0.3, axis='x')

    # Bottom: Summary statistics
    ax4 = fig.add_axes([0.06, 0.05, 0.88, 0.44])
    ax4.axis('off')

    # Calculate detailed statistics
    if len(viable_df) > 0:
        time_to_viable = viable_df['time_to_viable_min']
        available_dur = viable_df['available_duration_min']
        can_downlink = viable_df['can_downlink_gb']

        # Station breakdown
        station_stats = viable_df.groupby('station').agg({
            'time_to_viable_min': ['count', 'mean', 'median'],
            'available_duration_min': 'mean',
            'can_downlink_gb': 'mean'
        }).round(1)

        station_table = "VIABLE DOWNLINK STATION BREAKDOWN\n"
        station_table += f"{'Station':<18} {'Count':>5} {'%':>5} {'Avg Wait':>9} {'Avg Dur':>8} {'Avg Data':>9}\n"
        station_table += "-" * 60 + "\n"

        for station in station_stats.index:
            count = int(station_stats.loc[station, ('time_to_viable_min', 'count')])
            pct = 100 * count / viable_collects if viable_collects > 0 else 0
            avg_wait = station_stats.loc[station, ('time_to_viable_min', 'mean')]
            avg_dur = station_stats.loc[station, ('available_duration_min', 'mean')]
            avg_data = station_stats.loc[station, ('can_downlink_gb', 'mean')]
            station_table += f"{station:<18} {count:>5} {pct:>4.0f}% {avg_wait:>8.1f}m {avg_dur:>7.1f}m {avg_data:>8.1f}GB\n"

        summary_text = f"""VIABLE KA-BAND POST-COLLECT DOWNLINK ANALYSIS
{'='*72}
Viable = Ka contact with sufficient duration to downlink {IMAGE_SIZE_GB}GB image after {PROCESSING_BUFFER_MIN}-min processing

VIABILITY SUMMARY                            TIMING STATISTICS
  Total Collects: {total_collects:<26} Time to Viable Ka Downlink:
  Viable Ka Downlinks: {viable_collects} ({100*viable_collects/total_collects:.0f}%)                  Minimum: {time_to_viable.min():.1f} min
  Non-Viable: {non_viable_collects} ({100*non_viable_collects/total_collects:.0f}%)                          Median: {time_to_viable.median():.1f} min
                                             P95: {time_to_viable.quantile(0.95):.1f} min
                                             Maximum: {time_to_viable.max():.1f} min

DOWNLINK PARAMETERS                          CONTACT DURATION (available for downlink)
  Processing Buffer: {PROCESSING_BUFFER_MIN} min                      Minimum: {available_dur.min():.1f} min
  Downlink Rate: {downlink_rate_mbps} Mbps                       Median: {available_dur.median():.1f} min
  Image Size: {IMAGE_SIZE_GB} GB                             Average: {available_dur.mean():.1f} min
  Min Contact for {IMAGE_SIZE_GB}GB: {min_contact_duration_min:.1f} min              Maximum: {available_dur.max():.1f} min

DATA DELIVERY CAPACITY                       DELIVERY SLA (time to viable downlink)
  Avg Data per Contact: {can_downlink.mean():.1f} GB              ≤ 10 min: {len(time_to_viable[time_to_viable <= 10])}/{viable_collects} ({100*len(time_to_viable[time_to_viable <= 10])/viable_collects if viable_collects > 0 else 0:.0f}%)
  Min Data per Contact: {can_downlink.min():.1f} GB               ≤ 30 min: {len(time_to_viable[time_to_viable <= 30])}/{viable_collects} ({100*len(time_to_viable[time_to_viable <= 30])/viable_collects if viable_collects > 0 else 0:.0f}%)
  Max Data per Contact: {can_downlink.max():.1f} GB               ≤ 1 hr: {len(time_to_viable[time_to_viable <= 60])}/{viable_collects} ({100*len(time_to_viable[time_to_viable <= 60])/viable_collects if viable_collects > 0 else 0:.0f}%)

{station_table}"""
    else:
        summary_text = f"""VIABLE KA-BAND POST-COLLECT DOWNLINK ANALYSIS
{'='*72}
Viable = Ka contact with sufficient duration to downlink {IMAGE_SIZE_GB}GB image after {PROCESSING_BUFFER_MIN}-min processing

WARNING: No viable Ka downlinks found for any collect!

  Total Collects: {total_collects}
  Viable Ka Downlinks: 0 (0%)
  Non-Viable: {total_collects} (100%)

  Processing Buffer: {PROCESSING_BUFFER_MIN} min
  Downlink Rate: {downlink_rate_mbps} Mbps
  Image Size: {IMAGE_SIZE_GB} GB
  Min Contact Duration Needed: {min_contact_duration_min:.1f} min
"""

    ax4.text(0.0, 1.0, summary_text, transform=ax4.transAxes, fontsize=9,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#e8f5e9', alpha=0.9))

    fig.suptitle('Viable Ka-Band Post-Collect Downlink Analysis',
                 fontsize=14, fontweight='bold', y=0.98)

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    # Build station contact counts for Excel
    station_contact_dict = station_counts.to_dict() if len(viable_df) > 0 else {}

    return {
        'total_collects': total_collects,
        'viable_collects': viable_collects,
        'viable_pct': 100 * viable_collects / total_collects if total_collects > 0 else 0,
        'median_time_to_viable': viable_df['time_to_viable_min'].median() if len(viable_df) > 0 else None,
        'p95_time_to_viable': viable_df['time_to_viable_min'].quantile(0.95) if len(viable_df) > 0 else None,
        'min_contact_duration_needed': min_contact_duration_min,
        'image_size_gb': IMAGE_SIZE_GB,
        'downlink_rate_mbps': downlink_rate_mbps,
        'processing_buffer_min': PROCESSING_BUFFER_MIN,
        'station_contacts': station_contact_dict,
    }


def generate_slide_regional_capacity(tle_data: Dict, config, targets_gdf, access_df: pd.DataFrame,
                                      output_path: Path):
    """Generate regional capacity analysis slide.

    Explains why access opportunities may be limited for a regional target deck:
    - Analyzes orbit geometry vs target distribution
    - Shows monthly access patterns over 1 year
    - Provides orbit design recommendations
    """
    fig = plt.figure(figsize=(16, 10))

    # Get orbital parameters
    sat_alt_km = config.satellites[0].altitude_km if config.satellites else 250
    inclination = config.satellites[0].inclination_deg if config.satellites else 96.5
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + sat_alt_km
    period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60

    # Calculate swath width from imaging mode
    fov_half_angle = 2.5  # degrees, default panchromatic
    if hasattr(config, 'imaging_modes') and config.imaging_modes:
        fov_half_angle = config.imaging_modes[0].fov_half_angle_deg

    # Swath width calculation
    swath_half_km = sat_alt_km * np.tan(np.radians(fov_half_angle))
    swath_width_km = 2 * swath_half_km

    # Get off-nadir capability
    max_ona = 30.0  # default
    if hasattr(config, 'imaging_modes') and config.imaging_modes:
        max_ona = config.imaging_modes[0].off_nadir_max_deg

    # Access swath (with pointing)
    access_half_km = sat_alt_km * np.tan(np.radians(max_ona))
    access_swath_km = 2 * access_half_km

    # Get target statistics
    target_lats = []
    target_lons = []
    for idx, row in targets_gdf.iterrows():
        geom = row.geometry
        if geom.geom_type == 'Point':
            target_lats.append(geom.y)
            target_lons.append(geom.x)
        else:
            centroid = geom.centroid
            target_lats.append(centroid.y)
            target_lons.append(centroid.x)

    lat_min, lat_max = min(target_lats), max(target_lats)
    lon_min, lon_max = min(target_lons), max(target_lons)
    lat_span = lat_max - lat_min
    lon_span = lon_max - lon_min

    # Calculate ground track repeat cycle
    # Key question: how many days until ground track returns to same longitude at descending node?
    orbits_per_day = 24 * 60 / period_min  # e.g., 16.107 for 89.4 min period
    nodal_shift_per_orbit = 360.0 / orbits_per_day  # degrees longitude shift per orbit

    # Daily ground track drift at descending node:
    # Each day the satellite completes a fractional number of orbits
    # The fractional part determines how much the track shifts day-to-day
    minutes_per_day = 24 * 60
    remainder_minutes = minutes_per_day % period_min  # extra minutes beyond complete orbits
    daily_drift_deg = (remainder_minutes / period_min) * nodal_shift_per_orbit  # westward drift per day

    # Repeat cycle: days until ground track returns to same longitude
    if daily_drift_deg > 0.01:
        repeat_cycle_days = 360.0 / daily_drift_deg
    else:
        repeat_cycle_days = 1  # Nearly exact repeat each day

    # Simulate monthly access over 1 year
    print("    Simulating 1-year access patterns...")
    monthly_accesses = []
    sat_id = list(tle_data.keys())[0]

    # Use 30-day periods starting from different months
    base_date = config.start_datetime

    for month_offset in range(12):
        # Simulate each month by shifting the start date
        month_start = base_date + timedelta(days=30 * month_offset)

        # Generate ground track for this month (use coarse step for speed)
        track = calculate_ground_track(
            tle_data, sat_id, month_start,
            duration=timedelta(days=30), step_s=60.0
        )

        if not track:
            monthly_accesses.append(0)
            continue

        # Count unique target accesses this month
        month_access_count = 0
        for t_lat, t_lon in zip(target_lats, target_lons):
            # Check if any track point comes within access swath of target
            for lon, lat in track[::10]:  # Sample every 10th point for speed
                dist_km = np.sqrt(
                    ((lat - t_lat) * 111)**2 +
                    ((lon - t_lon) * 111 * np.cos(np.radians(t_lat)))**2
                )
                if dist_km <= access_swath_km / 2:
                    month_access_count += 1
                    break

        monthly_accesses.append(month_access_count)

    # Current month stats
    current_accesses = len(access_df) if access_df is not None else 0
    current_targets_accessed = access_df['AOI_Name'].nunique() if access_df is not None and 'AOI_Name' in access_df.columns else current_accesses

    # =========================================================================
    # Create plots
    # =========================================================================

    # Generate all descending passes over mission duration
    print("    Identifying descending passes over mission...")
    full_track = calculate_ground_track(
        tle_data, sat_id, config.start_datetime,
        duration=timedelta(days=config.duration_days), step_s=30.0
    )

    # Identify descending pass segments (latitude decreasing)
    descending_passes = []
    current_pass = []
    prev_lat = None

    if full_track:
        for i, (lon, lat) in enumerate(full_track):
            if prev_lat is not None:
                if lat < prev_lat:  # Descending
                    current_pass.append((lon, lat))
                else:  # Ascending - end current pass
                    if len(current_pass) > 10:  # Minimum pass length
                        descending_passes.append(current_pass)
                    current_pass = []
            prev_lat = lat
        # Don't forget last pass
        if len(current_pass) > 10:
            descending_passes.append(current_pass)

    # For each pass, check if it has access to any target
    passes_with_access = []
    passes_without_access = []

    for pass_track in descending_passes:
        has_access = False
        for t_lat, t_lon in zip(target_lats, target_lons):
            for lon, lat in pass_track[::3]:  # Sample for speed
                dist_km = np.sqrt(
                    ((lat - t_lat) * 111)**2 +
                    ((lon - t_lon) * 111 * np.cos(np.radians(t_lat)))**2
                )
                if dist_km <= access_swath_km / 2:
                    has_access = True
                    break
            if has_access:
                break

        if has_access:
            passes_with_access.append(pass_track)
        else:
            passes_without_access.append(pass_track)

    total_passes = len(descending_passes)
    access_passes = len(passes_with_access)

    # Top left: Target distribution map with all descending passes
    ax1 = fig.add_axes([0.03, 0.52, 0.46, 0.43], projection=ccrs.PlateCarree())
    ax1.set_extent([lon_min - 15, lon_max + 15, lat_min - 10, lat_max + 10], crs=ccrs.PlateCarree())
    ax1.add_feature(cfeature.LAND, facecolor='#f5f5f5')
    ax1.add_feature(cfeature.OCEAN, facecolor='#e0e8ef')
    ax1.add_feature(cfeature.COASTLINE, linewidth=0.5)
    ax1.add_feature(cfeature.BORDERS, linewidth=0.3, linestyle=':')
    ax1.gridlines(draw_labels=True, linewidth=0.3, alpha=0.5)

    # Plot passes without access (gray, behind)
    for pass_track in passes_without_access:
        pass_lons, pass_lats = zip(*pass_track)
        ax1.plot(pass_lons, pass_lats, '-', color='gray', linewidth=0.3, alpha=0.4,
                 transform=ccrs.PlateCarree(), zorder=3)

    # Plot passes with access (green, on top)
    for pass_track in passes_with_access:
        pass_lons, pass_lats = zip(*pass_track)
        ax1.plot(pass_lons, pass_lats, '-', color='green', linewidth=0.8, alpha=0.7,
                 transform=ccrs.PlateCarree(), zorder=4)

    # Plot targets
    ax1.scatter(target_lons, target_lats, c='red', s=60, marker='x', linewidths=2,
                transform=ccrs.PlateCarree(), zorder=5)

    # Custom legend
    from matplotlib.lines import Line2D
    legend_elements = [
        Line2D([0], [0], color='green', linewidth=2, label=f'Pass w/ Access ({access_passes})'),
        Line2D([0], [0], color='gray', linewidth=1, alpha=0.5, label=f'Pass w/o Access ({total_passes - access_passes})'),
        Line2D([0], [0], marker='x', color='red', linestyle='None', markersize=8, label=f'Targets ({len(target_lats)})'),
    ]
    ax1.legend(handles=legend_elements, loc='lower left', fontsize=8)
    ax1.set_title(f'Descending Passes ({config.duration_days} days) - 1 Collect Max per Pass', fontsize=10, fontweight='bold')

    # Top right: Monthly access pattern (bar chart)
    ax2 = fig.add_axes([0.54, 0.52, 0.43, 0.43])
    months = ['M1', 'M2', 'M3', 'M4', 'M5', 'M6', 'M7', 'M8', 'M9', 'M10', 'M11', 'M12']
    colors = ['green' if m == monthly_accesses[0] else 'steelblue' for m in monthly_accesses]
    colors[0] = 'orange'  # Highlight current month

    bars = ax2.bar(months, monthly_accesses, color=colors, edgecolor='black', alpha=0.7)
    ax2.axhline(np.mean(monthly_accesses), color='red', linestyle='--',
                label=f'Avg: {np.mean(monthly_accesses):.1f} targets/month')
    ax2.axhline(len(target_lats), color='gray', linestyle=':',
                label=f'Total Targets: {len(target_lats)}')

    ax2.set_xlabel('Month (from mission start)', fontsize=9)
    ax2.set_ylabel('Targets Accessible', fontsize=9)
    ax2.set_title('Projected Monthly Target Access (1 Year)', fontsize=10, fontweight='bold')
    ax2.legend(fontsize=8)
    ax2.set_ylim(0, len(target_lats) + 2)

    # Add value labels on bars
    for bar, val in zip(bars, monthly_accesses):
        ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.5,
                 str(val), ha='center', va='bottom', fontsize=8)

    # Bottom left: Analysis text
    ax3 = fig.add_axes([0.03, 0.05, 0.46, 0.42])
    ax3.axis('off')

    # Calculate statistics
    avg_monthly = np.mean(monthly_accesses)
    max_monthly = max(monthly_accesses)
    min_monthly = min(monthly_accesses)
    variation_pct = (max_monthly - min_monthly) / avg_monthly * 100 if avg_monthly > 0 else 0

    analysis_text = f"""REGIONAL CAPACITY ANALYSIS
{'='*50}

TARGET DISTRIBUTION
  Total Targets: {len(target_lats)}
  Latitude Range: {lat_min:.1f}° to {lat_max:.1f}° ({lat_span:.1f}° span)
  Longitude Range: {lon_min:.1f}° to {lon_max:.1f}° ({lon_span:.1f}° span)
  Region Width: ~{lon_span * 111 * np.cos(np.radians((lat_min+lat_max)/2)):.0f} km

DESCENDING PASSES ({config.duration_days} days)
  Total Descending Passes: {total_passes}
  Passes with Target Access: {access_passes} ({100*access_passes/total_passes:.1f}%)
  Passes without Access: {total_passes - access_passes}
  Note: Max 1 collect per pass (single satellite)

IMAGING CAPABILITY
  Max Off-Nadir: {max_ona}°
  Access Swath: {access_swath_km:.1f} km
  Orbital Period: {period_min:.1f} min

CURRENT PERIOD RESULTS
  Collect Opportunities: {current_accesses}
  Unique Targets Accessed: {current_targets_accessed}/{len(target_lats)}
  Target Coverage: {100*current_targets_accessed/len(target_lats):.1f}%
"""

    ax3.text(0, 0.98, analysis_text, transform=ax3.transAxes, fontsize=9,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#e3f2fd', alpha=0.8))

    # Bottom right: Recommendations
    ax4 = fig.add_axes([0.54, 0.05, 0.43, 0.42])
    ax4.axis('off')

    rec_text = f"""GROUND TRACK REPEAT CYCLE
{'='*50}

DAILY DRIFT AT DESCENDING NODE
  Orbits per Day: {orbits_per_day:.3f}
  Nodal Shift per Orbit: {nodal_shift_per_orbit:.2f}° (westward)
  Daily Drift: {daily_drift_deg:.2f}° (westward)
  Repeat Cycle: ~{repeat_cycle_days:.0f} days

WHAT THIS MEANS
  • Ground track shifts {daily_drift_deg:.1f}° west each day
  • After {repeat_cycle_days:.0f} days, track returns to same longitude
  • Target region spans {lon_span:.0f}° longitude
  • Access swath covers {access_swath_km/111:.1f}° per pass
  • Region needs ~{int(lon_span / (access_swath_km/111))+1} passes to fully cover

1-YEAR ACCESS PROJECTION
  Average: {avg_monthly:.1f} targets/month
  Min: {min_monthly} (Month {monthly_accesses.index(min_monthly)+1})
  Max: {max_monthly} (Month {monthly_accesses.index(max_monthly)+1})

ORBIT DESIGN OPTIONS:
  1. Different LTDN for different phasing
  2. Higher altitude for wider swath
  3. Add 2nd satellite (180° RAAN offset)
  4. Increase off-nadir capability
"""

    ax4.text(0, 0.98, rec_text, transform=ax4.transAxes, fontsize=9,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#fff3e0', alpha=0.8))

    fig.suptitle('Regional Capacity Analysis - Why Limited Access?',
                 fontsize=14, fontweight='bold', y=0.98)

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {
        'current_accesses': current_accesses,
        'avg_monthly': avg_monthly,
        'variation_pct': variation_pct
    }


def generate_slide4_access_results(access_df: pd.DataFrame, tle_data: Dict, output_path: Path):
    """Generate access results overview plot."""
    fig = plt.figure(figsize=(16, 9))
    ax = plt.axes(projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_global()

    colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd']

    for _, access in access_df.iterrows():
        sat_id = access['Satellite']
        start_time = pd.to_datetime(access['Start_Time'])
        if start_time.tzinfo is None:
            start_time = start_time.tz_localize('UTC')

        duration_s = access.get('Access_Duration', 60)
        ground_track = calculate_ground_track(
            tle_data, sat_id, start_time.to_pydatetime(),
            duration=timedelta(seconds=duration_s), step_s=10.0
        )
        if not ground_track:
            continue

        lons, lats = zip(*ground_track)
        color = colors[sat_id % len(colors)]

        ax.plot(lons, lats, color=color, linewidth=1, alpha=0.6, transform=ccrs.PlateCarree())
        ax.plot(access['AOI_Lon'], access['AOI_Lat'], 'x', color='blue', markersize=6,
                markeredgewidth=1.5, transform=ccrs.PlateCarree())

    ax.set_title(f'Access Windows Overview ({len(access_df)} accesses)', fontsize=14, fontweight='bold')

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


def generate_slide5_access_with_cones(access_df: pd.DataFrame, tle_data: Dict,
                                      ground_stations: List, sat_alt_km: float, output_path: Path):
    """Generate access plot with 60-min post-collect tracks and comm cones."""
    fig = plt.figure(figsize=(16, 9))
    ax = plt.axes(projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_global()

    # Plot comm cones
    plot_ground_stations_and_cones(ax, ground_stations, sat_alt_km)

    colors = ['#1f77b4', '#ff7f0e', '#2ca02c', '#d62728', '#9467bd']

    for _, access in access_df.iterrows():
        sat_id = access['Satellite']
        start_time = pd.to_datetime(access['Start_Time'])
        if start_time.tzinfo is None:
            start_time = start_time.tz_localize('UTC')

        # 60 min ground track after collect
        ground_track = calculate_ground_track(
            tle_data, sat_id, start_time.to_pydatetime(),
            duration=timedelta(minutes=60), step_s=30.0
        )
        if not ground_track:
            continue

        lons, lats = zip(*ground_track)
        color = colors[sat_id % len(colors)]

        ax.plot(lons, lats, color=color, linewidth=1, alpha=0.7, transform=ccrs.PlateCarree())
        ax.plot(access['AOI_Lon'], access['AOI_Lat'], 'x', color='blue', markersize=6,
                markeredgewidth=1.5, transform=ccrs.PlateCarree())

    ax.set_title('Access Windows with 60-min Post-Collect Tracks', fontsize=14, fontweight='bold')

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()


def calculate_swath_polygon(ground_track: List[Tuple], sat_alt_km: float,
                            off_nadir_min_deg: float, off_nadir_max_deg: float):
    """Calculate imaging swath polygon based on off-nadir angle limits."""
    from pyproj import Geod
    from shapely.geometry import Polygon

    if len(ground_track) < 2:
        return None

    geod = Geod(ellps='WGS84')
    earth_radius_km = 6378.137

    # Calculate ground range for min and max off-nadir angles
    def ona_to_ground_range(ona_deg):
        ona_rad = np.radians(ona_deg)
        # Simplified geometry: ground_range ≈ altitude * tan(ona)
        return sat_alt_km * np.tan(ona_rad) * 1000  # meters

    range_min = ona_to_ground_range(off_nadir_min_deg)
    range_max = ona_to_ground_range(off_nadir_max_deg)

    # Build swath polygon by projecting perpendicular to track
    left_points = []
    right_points = []

    for i, (lon, lat) in enumerate(ground_track):
        # Calculate track heading
        if i < len(ground_track) - 1:
            next_lon, next_lat = ground_track[i + 1]
            fwd_az, _, _ = geod.inv(lon, lat, next_lon, next_lat)
        else:
            prev_lon, prev_lat = ground_track[i - 1]
            fwd_az, _, _ = geod.inv(prev_lon, prev_lat, lon, lat)

        # Perpendicular azimuths (left and right of track)
        left_az = (fwd_az - 90) % 360
        right_az = (fwd_az + 90) % 360

        # Project to max off-nadir range on both sides
        left_lon, left_lat, _ = geod.fwd(lon, lat, left_az, range_max)
        right_lon, right_lat, _ = geod.fwd(lon, lat, right_az, range_max)

        left_points.append((left_lon, left_lat))
        right_points.append((right_lon, right_lat))

    # Create polygon: left side forward, right side backward
    polygon_coords = left_points + right_points[::-1] + [left_points[0]]

    try:
        poly = Polygon(polygon_coords)
        if not poly.is_valid:
            poly = poly.buffer(0)
        return poly
    except:
        return None


def generate_broadside_collect_plot(access_row: pd.Series, tle_data: Dict, ground_stations: List,
                                    sat_alt_km: float, output_path: Path,
                                    off_nadir_min: float = 0.0, off_nadir_max: float = 30.0):
    """Generate broadside collect opportunity plot with swath visualization."""
    sat_id = access_row['Satellite']
    start_time = pd.to_datetime(access_row['Start_Time'])
    if start_time.tzinfo is None:
        start_time = start_time.tz_localize('UTC')

    aoi_lat = access_row['AOI_Lat']
    aoi_lon = access_row['AOI_Lon']

    fig = plt.figure(figsize=(13, 6))
    ax = fig.add_axes([0.05, 0.05, 0.90, 0.82], projection=ccrs.PlateCarree())
    setup_map_axes(ax)
    ax.set_aspect('auto')  # Don't preserve geographic aspect ratio

    # Use orange color for satellite track
    track_color = '#ff7f0e'  # Orange

    # Calculate extended track for context (before and after)
    track_before = calculate_ground_track(
        tle_data, sat_id,
        (start_time - timedelta(minutes=3)).to_pydatetime(),
        duration=timedelta(minutes=3), step_s=10.0
    )
    track_during = calculate_ground_track(
        tle_data, sat_id, start_time.to_pydatetime(),
        duration=timedelta(seconds=access_row.get('Access_Duration', 60)), step_s=5.0
    )
    track_after = calculate_ground_track(
        tle_data, sat_id,
        (start_time + timedelta(seconds=access_row.get('Access_Duration', 60))).to_pydatetime(),
        duration=timedelta(minutes=3), step_s=10.0
    )

    # Calculate and plot swath polygon
    if track_during:
        swath_poly = calculate_swath_polygon(track_during, sat_alt_km, off_nadir_min, off_nadir_max)
        if swath_poly and swath_poly.is_valid:
            ax.add_geometries([swath_poly], crs=ccrs.PlateCarree(),
                            facecolor=track_color, edgecolor='none',
                            alpha=0.15, zorder=3, label='Imaging Swath')

    # Plot tracks - gray for before/after, orange for access window
    if track_before:
        lons, lats = zip(*track_before)
        ax.plot(lons, lats, color='gray', linewidth=1.5, alpha=0.5, transform=ccrs.PlateCarree())

    # Calculate actual ONA to target at broadside
    target_ona_deg = None
    mid_lon, mid_lat = None, None

    if track_during:
        lons, lats = zip(*track_during)
        ax.plot(lons, lats, color=track_color, linewidth=2.5, transform=ccrs.PlateCarree(),
                label='Access Window')

        # Mark broadside (midpoint) with a diamond marker
        mid_idx = len(track_during) // 2
        mid_lon, mid_lat = track_during[mid_idx]
        ax.plot(mid_lon, mid_lat, 'D', color=track_color, markersize=10,
                markeredgecolor='black', markeredgewidth=1.5,
                transform=ccrs.PlateCarree(), label='Broadside', zorder=10)

        # Calculate ground distance from nadir to target at broadside
        ground_dist_km = np.sqrt(
            ((mid_lat - aoi_lat) * 111)**2 +
            ((mid_lon - aoi_lon) * 111 * np.cos(np.radians(aoi_lat)))**2
        )
        # ONA = arctan(ground_distance / altitude)
        target_ona_deg = np.degrees(np.arctan(ground_dist_km / sat_alt_km))

    if track_after:
        lons, lats = zip(*track_after)
        ax.plot(lons, lats, color='gray', linewidth=1.5, alpha=0.5, transform=ccrs.PlateCarree())

    # Plot target as blue X
    ax.plot(aoi_lon, aoi_lat, 'x', color='blue', markersize=12, markeredgewidth=3,
            transform=ccrs.PlateCarree(), label='Target', zorder=11)

    # Set extent around the access
    all_lons = [aoi_lon]
    all_lats = [aoi_lat]
    if track_during:
        lons, lats = zip(*track_during)
        all_lons.extend(lons)
        all_lats.extend(lats)

    margin = 5
    ax.set_extent([min(all_lons) - margin, max(all_lons) + margin,
                   min(all_lats) - margin, max(all_lats) + margin], crs=ccrs.PlateCarree())

    # Legend
    legend_elements = [
        Line2D([0], [0], color='gray', linewidth=1.5, alpha=0.5, label='Ground Track'),
        Line2D([0], [0], color=track_color, linewidth=2.5, label='Access Window'),
        Patch(facecolor=track_color, alpha=0.15, label=f'Imaging Swath (ONA {off_nadir_min}-{off_nadir_max}°)'),
        Line2D([0], [0], marker='D', color=track_color, markersize=8, linestyle='None',
               markeredgecolor='black', label='Broadside'),
        Line2D([0], [0], marker='x', color='blue', markersize=10, linestyle='None',
               markeredgewidth=3, label='Target'),
    ]
    ax.legend(handles=legend_elements, loc='upper right', fontsize=9)

    ona_str = f"ONA to Target: {target_ona_deg:.1f}°" if target_ona_deg is not None else ""
    ax.set_title(f'Broadside Collect Opportunity\n'
                 f'Target: ({aoi_lat:.1f}°, {aoi_lon:.1f}°) | Satellite {sat_id} | {ona_str}\n'
                 f'{start_time.strftime("%Y-%m-%d %H:%M:%S")} UTC',
                 fontsize=11, fontweight='bold')

    plt.savefig(output_path, dpi=150, facecolor='white')
    plt.close()


def generate_pre_collect_diagram(access_row: pd.Series, tle_data: Dict, ground_stations: List,
                                  sat_alt_km: float, output_path: Path):
    """Generate diagram showing ground track BEFORE collect, back to last TT&C contact."""
    sat_id = access_row['Satellite']
    start_time = pd.to_datetime(access_row['Start_Time'])
    if start_time.tzinfo is None:
        start_time = start_time.tz_localize('UTC')

    aoi_lat = access_row['AOI_Lat']
    aoi_lon = access_row['AOI_Lon']

    fig = plt.figure(figsize=(16, 10))
    ax = plt.axes([0.05, 0.05, 0.70, 0.90], projection=ccrs.PlateCarree())
    setup_map_axes(ax)

    # Calculate ground track going BACKWARDS from collect to find last TT&C
    # Search back 12 hours to ensure we find a contact (covers ~8 revolutions)
    SEARCH_HOURS = 12
    search_start_ts = start_time - timedelta(hours=SEARCH_HOURS)
    search_start = search_start_ts.to_pydatetime()
    if search_start.tzinfo is not None:
        search_start = search_start.replace(tzinfo=None)

    pre_track = calculate_ground_track(
        tle_data, sat_id, search_start,
        duration=timedelta(hours=SEARCH_HOURS), step_s=30.0
    )

    if not pre_track:
        plt.close()
        return None

    # Find last TT&C contact before collect
    # Note: Ka-band stations are tri-band, so ALL stations can support TT&C
    # Exclude Tokyo from uplink analysis (too close to APAC targets)
    # Require TT&C contact to end at least 2 minutes before collect
    MIN_UPLINK_BUFFER_MIN = 2.0
    min_buffer_idx = int(MIN_UPLINK_BUFFER_MIN * 60 / 30)  # Convert to track indices (30s step)
    max_valid_idx = len(pre_track) - min_buffer_idx  # TT&C must end before this index

    last_ttc_contact = None
    last_ttc_station = None
    last_ttc_end_idx = 0

    for gs in ground_stations:
        # Exclude Tokyo from TT&C uplink analysis
        if 'Tokyo' in gs['name']:
            continue

        contact_elev = 5.0
        max_range_km = calculate_comm_range_km(contact_elev, sat_alt_km)

        in_contact = False
        contact_end_idx = None

        for i, (lon, lat) in enumerate(pre_track):
            dist = calculate_ground_distance_km(lat, lon, gs['lat'], gs['lon'])
            currently_in = dist <= max_range_km

            if currently_in:
                in_contact = True
                contact_end_idx = i
            elif not currently_in and in_contact:
                # Contact ended - only count if it ends at least 2 min before collect
                if contact_end_idx is not None and contact_end_idx > last_ttc_end_idx and contact_end_idx <= max_valid_idx:
                    last_ttc_end_idx = contact_end_idx
                    last_ttc_station = gs['name']
                    last_ttc_contact = search_start + timedelta(seconds=contact_end_idx * 30)
                in_contact = False

        # Handle contact still ongoing at end of track - only if it ends before the buffer
        if in_contact and contact_end_idx is not None and contact_end_idx > last_ttc_end_idx and contact_end_idx <= max_valid_idx:
            last_ttc_end_idx = contact_end_idx
            last_ttc_station = gs['name']
            last_ttc_contact = search_start + timedelta(seconds=contact_end_idx * 30)

    # Plot the pre-collect track from last TT&C to collect
    if last_ttc_end_idx > 0:
        track_segment = pre_track[last_ttc_end_idx:]
    else:
        track_segment = pre_track

    lons, lats = zip(*track_segment)
    ax.plot(lons, lats, color='purple', linewidth=1.5, alpha=0.7,
            transform=ccrs.PlateCarree(), label='Track from Last TT&C')

    # Mark last TT&C contact point
    if last_ttc_end_idx > 0:
        ttc_lon, ttc_lat = pre_track[last_ttc_end_idx]
        ax.plot(ttc_lon, ttc_lat, 'o', color='purple', markersize=10,
                transform=ccrs.PlateCarree(), label=f'Last TT&C ({last_ttc_station})')

    # Mark collect point
    ax.plot(aoi_lon, aoi_lat, 'x', color='blue', markersize=10, markeredgewidth=2,
            transform=ccrs.PlateCarree(), label='Collect Target')

    # Plot ground stations and cones
    plot_ground_stations_and_cones(ax, ground_stations, sat_alt_km)

    # Calculate uplink-to-collect time
    if last_ttc_contact:
        collect_time = start_time.to_pydatetime()
        if collect_time.tzinfo is not None:
            collect_time = collect_time.replace(tzinfo=None)
        uplink_to_collect_min = (collect_time - last_ttc_contact).total_seconds() / 60
    else:
        uplink_to_collect_min = float('nan')

    # Set extent
    margin = 20
    ax.set_extent([min(lons) - margin, max(lons) + margin,
                   max(min(lats) - margin, -85), min(max(lats) + margin, 85)], crs=ccrs.PlateCarree())

    # Legend on the right
    legend_elements = [
        Line2D([0], [0], color='purple', linewidth=1.5, label='Track from Last TT&C'),
        Line2D([0], [0], marker='o', color='purple', markersize=8, linestyle='None', label='Last TT&C Contact'),
        Line2D([0], [0], marker='x', color='blue', markersize=8, linestyle='None', markeredgewidth=2, label='Collect Target'),
        Line2D([0], [0], marker='^', color='green', markersize=8, linestyle='None', label='Ka Station'),
        Line2D([0], [0], marker='^', color='black', markersize=8, linestyle='None', label='TT&C Station'),
    ]
    ax.legend(handles=legend_elements, loc='upper left', bbox_to_anchor=(1.02, 1.0), fontsize=9)

    title = f'Pre-Collect Track - Target ({aoi_lat:.1f}°, {aoi_lon:.1f}°)\n'
    title += f'Last TT&C: {last_ttc_station or "None"}'
    if not np.isnan(uplink_to_collect_min):
        title += f' | Uplink-to-Collect: {uplink_to_collect_min:.1f} min'
    ax.set_title(title, fontsize=11, fontweight='bold')

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {
        'last_ttc_station': last_ttc_station,
        'uplink_to_collect_min': uplink_to_collect_min,
    }


def generate_collect_diagram_3h(access_row: pd.Series, tle_data: Dict, ground_stations: List,
                                sat_alt_km: float, output_path: Path,
                                onboard_processing_time_min: float = 5.0):
    """Generate 3-hour post-collect diagram with on-board processing buffer.

    Args:
        onboard_processing_time_min: Minimum time after collect before Ka downlink can begin (default 5 min)
    """
    sat_id = access_row['Satellite']
    start_time = pd.to_datetime(access_row['Start_Time'])
    if start_time.tzinfo is None:
        start_time = start_time.tz_localize('UTC')

    aoi_lat = access_row['AOI_Lat']
    aoi_lon = access_row['AOI_Lon']

    fig = plt.figure(figsize=(16, 10))
    ax = plt.axes([0.05, 0.05, 0.70, 0.90], projection=ccrs.PlateCarree())
    setup_map_axes(ax)

    # 3h ground track
    extended_track = calculate_ground_track(
        tle_data, sat_id, start_time.to_pydatetime(),
        duration=timedelta(hours=3), step_s=30.0
    )

    if not extended_track:
        plt.close()
        return None

    ext_lons, ext_lats = zip(*extended_track)
    ax.plot(ext_lons, ext_lats, color='gray', linewidth=1.5, alpha=0.6,
            transform=ccrs.PlateCarree(), label='Ground Track (3h)')

    # Target
    ax.plot(aoi_lon, aoi_lat, 'x', color='blue', markersize=8, markeredgewidth=2,
            transform=ccrs.PlateCarree(), label='Target')

    # Color map for contact times
    colors = ['#00ff00', '#2ecc71', '#f1c40f', '#e67e22', '#e74c3c']
    custom_cmap = LinearSegmentedColormap.from_list("contact_time", colors, N=256)

    # Track contacts and plot
    contact_stats = []
    first_ka_contact = None

    for gs in ground_stations:
        is_ka = gs.get('ka_band', False)
        marker_color = 'green' if is_ka else 'black'
        cone_color = 'green' if is_ka else 'black'
        cone_elev = 10.0 if is_ka else 5.0

        ax.plot(gs['lon'], gs['lat'], '^', color=marker_color, markersize=6,
                transform=ccrs.PlateCarree(), zorder=8)

        try:
            comm_cone = create_comm_cone_geodesic(gs['lon'], gs['lat'], cone_elev, sat_alt_km)
            if comm_cone and comm_cone.is_valid:
                ax.add_geometries([comm_cone], crs=ccrs.PlateCarree(),
                                facecolor='none', edgecolor=cone_color, alpha=0.5,
                                linewidth=0.8, zorder=2)
        except:
            pass

        # Find contacts
        contact_elev = 10.0 if is_ka else 5.0
        max_range_km = calculate_comm_range_km(contact_elev, sat_alt_km)

        in_contact = False
        contact_start_idx = None
        contacts = []

        for i, (lon, lat) in enumerate(extended_track):
            dist = calculate_ground_distance_km(lat, lon, gs['lat'], gs['lon'])
            currently_in = dist <= max_range_km

            if currently_in and not in_contact:
                contact_start_idx = i
                in_contact = True
            elif not currently_in and in_contact:
                contacts.append((contact_start_idx, i - 1))
                in_contact = False

        if in_contact:
            contacts.append((contact_start_idx, len(extended_track) - 1))

        for start_idx, end_idx in contacts:
            if end_idx <= start_idx:
                continue

            contact_lons = [extended_track[j][0] for j in range(start_idx, end_idx + 1)]
            contact_lats = [extended_track[j][1] for j in range(start_idx, end_idx + 1)]

            time_to_contact_min = start_idx * 30 / 60
            contact_duration_min = (end_idx - start_idx) * 30 / 60

            norm_time = min(time_to_contact_min / 180.0, 1.0)

            # Only color-highlight Ka-band contacts (used for payload downlink)
            # Non-Ka contacts shown in muted gray since they can't be used for payload downlink
            ka_valid = True
            if not is_ka:
                # Non-Ka station - show in muted gray (TT&C only, not for payload downlink)
                contact_color = '#aaaaaa'
                line_style = '-'
            elif time_to_contact_min < onboard_processing_time_min:
                # Ka contact too soon after collect - data still processing
                ka_valid = False
                contact_color = '#cccccc'
                line_style = '--'
            else:
                # Valid Ka contact - color by time
                contact_color = custom_cmap(norm_time)
                line_style = '-'

            # Plot the contact segment
            line_width = 3 if is_ka else 2  # Thinner lines for non-Ka contacts
            ax.plot(contact_lons, contact_lats, color=contact_color, linewidth=line_width,
                    linestyle=line_style, transform=ccrs.PlateCarree(), zorder=5 if is_ka else 4)

            contact_stats.append({
                'station': gs['name'],
                'is_ka': is_ka,
                'ka_valid': ka_valid,
                'time_to_contact_min': time_to_contact_min,
                'duration_min': contact_duration_min,
                'color': contact_color,
                'start_idx': start_idx,
                'end_idx': end_idx,
            })

            # Only consider valid Ka contacts for first_ka_contact
            if is_ka and ka_valid and (first_ka_contact is None or time_to_contact_min < first_ka_contact['time_to_contact_min']):
                first_ka_contact = contact_stats[-1]

    # Sort contacts by time
    contact_stats.sort(key=lambda x: x['time_to_contact_min'])

    # Set extent
    margin = 20
    ax.set_extent([min(ext_lons) - margin, max(ext_lons) + margin,
                   max(min(ext_lats) - margin, -85), min(max(ext_lats) + margin, 85)], crs=ccrs.PlateCarree())

    # Legend
    legend_elements = [
        Line2D([0], [0], color='gray', linewidth=1.5, alpha=0.6, label='Ground Track (3h)'),
        Line2D([0], [0], marker='x', color='blue', markersize=8, linestyle='None', markeredgewidth=2, label='Target'),
        Line2D([0], [0], marker='^', color='green', markersize=8, linestyle='None', label='Ka Station'),
        Patch(facecolor='none', edgecolor='green', label='Ka Comm Cone'),
        Line2D([0], [0], marker='^', color='black', markersize=8, linestyle='None', label='TT&C Station'),
        Patch(facecolor='none', edgecolor='black', label='TT&C Comm Cone'),
        Line2D([0], [0], color='#cccccc', linewidth=3, linestyle='--', label=f'Invalid Ka (<{onboard_processing_time_min:.0f}m)'),
        Line2D([0], [0], color='#aaaaaa', linewidth=2, linestyle='-', label='TT&C Contact (no Ka)'),
    ]

    # Add contact entries - only show Ka contacts in detail
    ka_contacts = [c for c in contact_stats if c['is_ka']]
    for contact in ka_contacts[:8]:  # Limit to 8 Ka contacts
        invalid_marker = " [INVALID]" if not contact.get('ka_valid', True) else ""
        label = f"{contact['station']} (Ka){invalid_marker}: C+{contact['time_to_contact_min']:.1f}m"
        linestyle = '--' if not contact.get('ka_valid', True) else '-'
        legend_elements.append(Line2D([0], [0], color=contact['color'], linewidth=3, linestyle=linestyle, label=label))

    ax.legend(handles=legend_elements, loc='upper left', bbox_to_anchor=(1.02, 1.0), fontsize=8)

    title = f'Post-Collect Track (3h) - Target ({aoi_lat:.1f}°, {aoi_lon:.1f}°)\n'
    title += f'On-board processing buffer: {onboard_processing_time_min:.0f} min\n'
    if first_ka_contact:
        title += f'First Valid Ka: {first_ka_contact["station"]} at C+{first_ka_contact["time_to_contact_min"]:.1f}min'
    else:
        title += 'No valid Ka contact within 3h (all within processing buffer)'
    ax.set_title(title, fontsize=11, fontweight='bold')

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {'first_ka_contact': first_ka_contact}


def generate_downlink_delay_summary_plot(excel_path: Path, output_path: Path):
    """Generate payload downlink delay summary analysis plot with station breakdown."""
    try:
        delay_df = pd.read_excel(excel_path, sheet_name='Downlink_Delay')
    except Exception:
        return None

    if len(delay_df) == 0:
        return None

    # Get downlink delay values in minutes
    if 'Downlink_Delay_minutes' in delay_df.columns:
        delay_values = delay_df['Downlink_Delay_minutes'].dropna()
    elif 'Downlink_Delay_seconds' in delay_df.columns:
        delay_values = delay_df['Downlink_Delay_seconds'].dropna() / 60
    else:
        return None

    fig = plt.figure(figsize=(16, 7))

    # Histogram (left)
    ax1 = fig.add_axes([0.05, 0.12, 0.28, 0.78])
    ax1.hist(delay_values, bins=20, color='#2ecc71', edgecolor='black', alpha=0.7)
    ax1.axvline(delay_values.median(), color='red', linestyle='--', linewidth=2, label=f'Median: {delay_values.median():.1f} min')
    ax1.axvline(delay_values.quantile(0.95), color='orange', linestyle='--', linewidth=2, label=f'P95: {delay_values.quantile(0.95):.1f} min')
    ax1.set_xlabel('Payload Downlink Delay (minutes)', fontsize=10)
    ax1.set_ylabel('Frequency', fontsize=10)
    ax1.set_title('Downlink Delay Distribution', fontsize=11, fontweight='bold')
    ax1.legend(fontsize=9)
    ax1.grid(True, alpha=0.3)

    # Box plot (middle)
    ax2 = fig.add_axes([0.38, 0.12, 0.18, 0.78])
    bp = ax2.boxplot(delay_values, vert=True, patch_artist=True)
    bp['boxes'][0].set_facecolor('#2ecc71')
    bp['boxes'][0].set_alpha(0.7)
    ax2.set_ylabel('Payload Downlink Delay (minutes)', fontsize=10)
    ax2.set_title('Overall Stats', fontsize=11, fontweight='bold')
    ax2.set_xticklabels(['All'])

    # Add stats text below box plot
    stats_text = f"Count: {len(delay_values)}\n"
    stats_text += f"Min: {delay_values.min():.1f} min\n"
    stats_text += f"Median: {delay_values.median():.1f} min\n"
    stats_text += f"Mean: {delay_values.mean():.1f} min\n"
    stats_text += f"P95: {delay_values.quantile(0.95):.1f} min\n"
    stats_text += f"Max: {delay_values.max():.1f} min"
    ax2.text(0.5, -0.25, stats_text, transform=ax2.transAxes, fontsize=9,
             verticalalignment='top', horizontalalignment='center',
             fontfamily='monospace')

    # Ka station breakdown table (right)
    ax3 = fig.add_axes([0.62, 0.12, 0.36, 0.78])
    ax3.axis('off')

    if 'Next_Ka_Station' in delay_df.columns:
        # Get stats per station
        station_stats = delay_df.groupby('Next_Ka_Station')['Downlink_Delay_minutes'].agg(
            ['count', 'mean', 'median', 'min', 'max']
        ).round(1)
        station_stats = station_stats.sort_values('count', ascending=False)

        # Build table text
        table_text = "First Ka Contact by Station\n"
        table_text += "=" * 50 + "\n\n"
        table_text += f"{'Station':<22} {'#':>4} {'Med':>6} {'Avg':>6} {'Min':>6} {'Max':>6}\n"
        table_text += "-" * 52 + "\n"

        for station in station_stats.index:
            count = int(station_stats.loc[station, 'count'])
            median = station_stats.loc[station, 'median']
            mean = station_stats.loc[station, 'mean']
            min_val = station_stats.loc[station, 'min']
            max_val = station_stats.loc[station, 'max']
            # Shorten station name
            short_name = station.replace('KSAT-', '').replace('-', ' ')[:20]
            table_text += f"{short_name:<22} {count:>4} {median:>5.0f}m {mean:>5.0f}m {min_val:>5.0f}m {max_val:>5.0f}m\n"

        table_text += "-" * 52 + "\n"
        total = len(delay_df)
        overall_median = delay_df['Downlink_Delay_minutes'].median()
        overall_mean = delay_df['Downlink_Delay_minutes'].mean()
        overall_min = delay_df['Downlink_Delay_minutes'].min()
        overall_max = delay_df['Downlink_Delay_minutes'].max()
        table_text += f"{'TOTAL':<22} {total:>4} {overall_median:>5.0f}m {overall_mean:>5.0f}m {overall_min:>5.0f}m {overall_max:>5.0f}m\n"

        ax3.text(0, 0.95, table_text, transform=ax3.transAxes, fontsize=10,
                 verticalalignment='top', fontfamily='monospace',
                 bbox=dict(boxstyle='round', facecolor='#e8f5e9', alpha=0.8))

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {'median': delay_values.median(), 'p95': delay_values.quantile(0.95)}


# Backward compatibility alias
generate_ttnc_summary_plot = generate_downlink_delay_summary_plot


def generate_uplink_to_collect_summary(access_df: pd.DataFrame, tle_data: Dict,
                                        ground_stations: List, sat_alt_km: float,
                                        output_path: Path):
    """Generate summary analysis of uplink-to-collect times (pre-collect TT&C) with station breakdown."""
    from src.vleo_eo.contacts import calculate_comm_range_km, calculate_ground_distance_km

    # Collect data with station info
    uplink_data = []

    for _, access_row in access_df.iterrows():
        sat_id = access_row['Satellite']
        start_time = pd.to_datetime(access_row['Start_Time'])
        if start_time.tzinfo is None:
            start_time = start_time.tz_localize('UTC')

        # Calculate track going backwards 12 hours to ensure we find a TT&C contact
        SEARCH_HOURS = 12
        search_start = (start_time - timedelta(hours=SEARCH_HOURS)).to_pydatetime()
        if search_start.tzinfo is not None:
            search_start = search_start.replace(tzinfo=None)

        pre_track = calculate_ground_track(
            tle_data, sat_id, search_start,
            duration=timedelta(hours=SEARCH_HOURS), step_s=30.0
        )

        if not pre_track:
            continue

        # Find last TT&C contact before collect
        # Note: Ka-band stations are tri-band, so ALL stations can support TT&C
        # Exclude Tokyo from uplink analysis (too close to APAC targets)
        # Require TT&C contact to end at least 2 minutes before collect
        MIN_UPLINK_BUFFER_MIN = 2.0
        min_buffer_idx = int(MIN_UPLINK_BUFFER_MIN * 60 / 30)  # Convert to track indices (30s step)
        max_valid_idx = len(pre_track) - min_buffer_idx  # TT&C must end before this index

        last_ttc_idx = None
        last_ttc_station = None
        for gs in ground_stations:
            # Exclude Tokyo from TT&C uplink analysis
            if 'Tokyo' in gs['name']:
                continue

            max_range_km = calculate_comm_range_km(5.0, sat_alt_km)

            for i, (lon, lat) in enumerate(pre_track):
                # Only consider contacts that end at least 2 min before collect
                if i > max_valid_idx:
                    break

                dist = calculate_ground_distance_km(lat, lon, gs['lat'], gs['lon'])
                if dist <= max_range_km:
                    if last_ttc_idx is None or i > last_ttc_idx:
                        last_ttc_idx = i
                        last_ttc_station = gs['name']

        if last_ttc_idx is not None and last_ttc_station is not None:
            # Time from last TT&C to collect (must be at least 2 min)
            uplink_to_collect_min = (len(pre_track) - last_ttc_idx) * 30 / 60
            uplink_data.append({
                'station': last_ttc_station,
                'uplink_time_min': uplink_to_collect_min
            })

    if len(uplink_data) == 0:
        return None

    uplink_df = pd.DataFrame(uplink_data)
    uplink_times = uplink_df['uplink_time_min'].values

    fig = plt.figure(figsize=(16, 7))

    # Histogram (left)
    ax1 = fig.add_axes([0.05, 0.12, 0.28, 0.78])
    ax1.hist(uplink_times, bins=20, color='#9b59b6', edgecolor='black', alpha=0.7)
    ax1.axvline(np.median(uplink_times), color='red', linestyle='--', linewidth=2,
                label=f'Median: {np.median(uplink_times):.1f} min')
    ax1.axvline(np.percentile(uplink_times, 95), color='orange', linestyle='--', linewidth=2,
                label=f'P95: {np.percentile(uplink_times, 95):.1f} min')
    ax1.set_xlabel('Time Since Last TT&C Contact (minutes)', fontsize=10)
    ax1.set_ylabel('Frequency', fontsize=10)
    ax1.set_title('Uplink-to-Collect Distribution', fontsize=11, fontweight='bold')
    ax1.legend(fontsize=9)
    ax1.grid(True, alpha=0.3)

    # Box plot (middle)
    ax2 = fig.add_axes([0.38, 0.12, 0.18, 0.78])
    bp = ax2.boxplot(uplink_times, vert=True, patch_artist=True)
    bp['boxes'][0].set_facecolor('#9b59b6')
    bp['boxes'][0].set_alpha(0.7)
    ax2.set_ylabel('Time Since Last TT&C (minutes)', fontsize=10)
    ax2.set_title('Overall Stats', fontsize=11, fontweight='bold')
    ax2.set_xticklabels(['All'])

    # Add stats text below box plot
    stats_text = f"Count: {len(uplink_times)}\n"
    stats_text += f"Min: {np.min(uplink_times):.1f} min\n"
    stats_text += f"Median: {np.median(uplink_times):.1f} min\n"
    stats_text += f"Mean: {np.mean(uplink_times):.1f} min\n"
    stats_text += f"P95: {np.percentile(uplink_times, 95):.1f} min\n"
    stats_text += f"Max: {np.max(uplink_times):.1f} min"
    ax2.text(0.5, -0.25, stats_text, transform=ax2.transAxes, fontsize=9,
             verticalalignment='top', horizontalalignment='center',
             fontfamily='monospace')

    # TT&C station breakdown table (right)
    ax3 = fig.add_axes([0.62, 0.12, 0.36, 0.78])
    ax3.axis('off')

    # Get stats per station
    station_stats = uplink_df.groupby('station')['uplink_time_min'].agg(
        ['count', 'mean', 'median', 'min', 'max']
    ).round(1)
    station_stats = station_stats.sort_values('count', ascending=False)

    # Build table text
    table_text = "Last TT&C Contact by Station\n"
    table_text += "=" * 50 + "\n\n"
    table_text += f"{'Station':<22} {'#':>4} {'Med':>6} {'Avg':>6} {'Min':>6} {'Max':>6}\n"
    table_text += "-" * 52 + "\n"

    for station in station_stats.index:
        count = int(station_stats.loc[station, 'count'])
        median = station_stats.loc[station, 'median']
        mean = station_stats.loc[station, 'mean']
        min_val = station_stats.loc[station, 'min']
        max_val = station_stats.loc[station, 'max']
        # Shorten station name
        short_name = station.replace('KSAT-', '').replace('-', ' ')[:20]
        table_text += f"{short_name:<22} {count:>4} {median:>5.0f}m {mean:>5.0f}m {min_val:>5.0f}m {max_val:>5.0f}m\n"

    table_text += "-" * 52 + "\n"
    total = len(uplink_df)
    overall_median = np.median(uplink_times)
    overall_mean = np.mean(uplink_times)
    overall_min = np.min(uplink_times)
    overall_max = np.max(uplink_times)
    table_text += f"{'TOTAL':<22} {total:>4} {overall_median:>5.0f}m {overall_mean:>5.0f}m {overall_min:>5.0f}m {overall_max:>5.0f}m\n"

    ax3.text(0, 0.95, table_text, transform=ax3.transAxes, fontsize=10,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#f3e5f5', alpha=0.8))

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {'median': np.median(uplink_times), 'p95': np.percentile(uplink_times, 95)}


def generate_optimization_slide(excel_path: Path, output_path: Path, config,
                                 tle_data: Dict = None, ground_stations: List[Dict] = None):
    """Generate ground station optimization analysis slide.

    Analyzes minimum station set needed to meet TT&C requirement:
    - Requirement: 1 TT&C pass per orbital revolution
    - Ka-band stations are tri-band (can do TT&C)
    - Shows gap analysis for station sets that don't meet requirement
    """
    # Get orbital parameters
    altitude_km = config.satellites[0].altitude_km if config.satellites else 250
    earth_radius_km = 6378.137
    mu = 398600.4418  # km³/s²
    semi_major_axis = earth_radius_km + altitude_km
    period_min = 2 * np.pi * np.sqrt(semi_major_axis**3 / mu) / 60
    duration_days = config.duration_days
    total_revs = int(duration_days * 24 * 60 / period_min)

    # If we don't have TLE data or ground stations, create a simpler slide
    if tle_data is None or ground_stations is None:
        fig = plt.figure(figsize=(16, 10))
        ax = fig.add_axes([0.1, 0.1, 0.8, 0.8])
        ax.axis('off')
        ax.text(0.5, 0.5, 'Ground Station Optimization\n\nInsufficient data for analysis',
                ha='center', va='center', fontsize=14)
        plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
        plt.close()
        return None

    # Generate mission ground track (use coarser step for efficiency)
    print("    Calculating TT&C coverage for station optimization...")
    sat_id = list(tle_data.keys())[0]
    ground_track = calculate_ground_track(
        tle_data, sat_id, config.start_datetime,
        duration=timedelta(days=duration_days), step_s=60.0  # 1-min steps for efficiency
    )

    if not ground_track:
        return None

    track_times = [config.start_datetime + timedelta(seconds=i * 60) for i in range(len(ground_track))]

    # Calculate TT&C contacts for each station (all stations can do TT&C since Ka is tri-band)
    # Use S/X-band elevation mask for TT&C
    min_elev_ttc = get_elevation_for_band(config, is_ka_band=False)

    station_contacts = {}
    for gs in ground_stations:
        gs_name = gs['name'].replace('Viasat-', '').replace('KSAT-', '')
        gs_lat = gs['lat']
        gs_lon = gs['lon']

        coverage_radius_km = calculate_coverage_radius(altitude_km, min_elev_ttc)
        coverage_radius_deg = coverage_radius_km / 111.0

        contacts = []
        in_contact = False
        contact_start_time = None

        for i, (lon, lat) in enumerate(ground_track):
            dist = np.sqrt((lat - gs_lat)**2 + ((lon - gs_lon) * np.cos(np.radians(gs_lat)))**2)

            if dist < coverage_radius_deg:
                if not in_contact:
                    in_contact = True
                    contact_start_time = track_times[i]
            else:
                if in_contact and contact_start_time:
                    contacts.append({
                        'start': contact_start_time,
                        'end': track_times[i-1] if i > 0 else contact_start_time
                    })
                    in_contact = False

        # Handle contact at end
        if in_contact and contact_start_time:
            contacts.append({
                'start': contact_start_time,
                'end': track_times[-1]
            })

        station_contacts[gs_name] = contacts

    def analyze_station_set(station_names):
        """Analyze TT&C coverage for a set of stations."""
        # Combine all contacts from selected stations
        all_contacts = []
        for name in station_names:
            if name in station_contacts:
                all_contacts.extend(station_contacts[name])

        if not all_contacts:
            return {'revs_without': total_revs, 'pct_coverage': 0, 'max_gap': float('inf'), 'gaps': []}

        # Sort by start time
        all_contacts.sort(key=lambda x: x['start'])

        # Calculate gaps
        gaps = []
        for i in range(1, len(all_contacts)):
            gap_min = (all_contacts[i]['start'] - all_contacts[i-1]['end']).total_seconds() / 60
            if gap_min > 0:
                gaps.append(gap_min)

        # Count revolutions without contact
        revs_without = 0
        for rev in range(total_revs):
            rev_start = config.start_datetime + timedelta(minutes=rev * period_min)
            rev_end = rev_start + timedelta(minutes=period_min)

            has_contact = False
            for contact in all_contacts:
                if contact['start'] < rev_end and contact['end'] > rev_start:
                    has_contact = True
                    break

            if not has_contact:
                revs_without += 1

        return {
            'total_contacts': len(all_contacts),
            'revs_without': revs_without,
            'pct_coverage': 100 * (total_revs - revs_without) / total_revs,
            'max_gap': max(gaps) if gaps else 0,
            'avg_gap': np.mean(gaps) if gaps else 0,
            'gaps': gaps
        }

    # Rank stations by contact count
    station_ranking = [(name, len(contacts)) for name, contacts in station_contacts.items()]
    station_ranking.sort(key=lambda x: -x[1])

    # Analyze full network
    all_station_names = [name for name, _ in station_ranking]
    full_network = analyze_station_set(all_station_names)

    # Find minimum station set using greedy approach
    # Start with highest-contact station, add stations until requirement met
    min_stations = []
    current_coverage = {'revs_without': total_revs, 'pct_coverage': 0}

    for station_name, _ in station_ranking:
        min_stations.append(station_name)
        current_coverage = analyze_station_set(min_stations)
        if current_coverage['revs_without'] == 0:
            break

    # Also test top 1, 2, 3 stations individually
    test_configs = [
        ('Top 1 Station', [station_ranking[0][0]] if station_ranking else []),
        ('Top 2 Stations', [s[0] for s in station_ranking[:2]] if len(station_ranking) >= 2 else []),
        ('Top 3 Stations', [s[0] for s in station_ranking[:3]] if len(station_ranking) >= 3 else []),
        ('Minimum Set', min_stations),
        ('Full Network', all_station_names),
    ]

    config_results = []
    for name, stations in test_configs:
        if stations:
            result = analyze_station_set(stations)
            result['name'] = name
            result['stations'] = stations
            config_results.append(result)

    # Create figure
    fig = plt.figure(figsize=(16, 10))

    # Top left: Requirements & Mission Parameters
    ax1 = fig.add_axes([0.04, 0.55, 0.44, 0.40])
    ax1.axis('off')

    req_text = "TT&C COVERAGE REQUIREMENT\n"
    req_text += "=" * 44 + "\n\n"
    req_text += "Requirement: 1 TT&C pass per revolution\n"
    req_text += "Note: Ka-band stations are tri-band (TT&C capable)\n\n"
    req_text += f"Mission Parameters:\n"
    req_text += f"  Duration: {duration_days} days\n"
    req_text += f"  Altitude: {altitude_km} km\n"
    req_text += f"  Orbital Period: {period_min:.1f} min\n"
    req_text += f"  Total Revolutions: {total_revs}\n\n"
    req_text += f"Current Network:\n"
    req_text += f"  Total Stations: {len(ground_stations)}\n"
    ka_count = sum(1 for gs in ground_stations if gs.get('ka_band', False))
    req_text += f"  Ka-band (tri-band): {ka_count}\n"
    req_text += f"  TT&C-only: {len(ground_stations) - ka_count}\n\n"
    req_text += f"Full Network Results:\n"
    req_text += f"  Total Contacts: {full_network['total_contacts']}\n"
    req_text += f"  Revs w/o Contact: {full_network['revs_without']}/{total_revs}\n"
    req_text += f"  Coverage: {full_network['pct_coverage']:.1f}%\n"
    req_text += f"  Max Gap: {full_network['max_gap']:.1f} min ({full_network['max_gap']/period_min:.2f} revs)\n"

    ax1.text(0, 0.98, req_text, transform=ax1.transAxes, fontsize=10,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#e3f2fd', alpha=0.8))

    # Top right: Station ranking by contacts
    ax2 = fig.add_axes([0.52, 0.55, 0.44, 0.40])
    ax2.axis('off')

    rank_text = "STATION RANKING (by TT&C contacts)\n"
    rank_text += "=" * 44 + "\n\n"
    rank_text += f"{'#':<3} {'Station':<22} {'Contacts':>8} {'Per Rev':>8}\n"
    rank_text += "-" * 44 + "\n"

    for i, (name, count) in enumerate(station_ranking[:12], 1):
        per_rev = count / total_revs
        rank_text += f"{i:<3} {name:<22} {count:>8} {per_rev:>7.2f}\n"

    if len(station_ranking) > 12:
        rank_text += f"    ... and {len(station_ranking) - 12} more stations\n"

    ax2.text(0, 0.98, rank_text, transform=ax2.transAxes, fontsize=10,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#fff3e0', alpha=0.8))

    # Bottom left: Configuration comparison
    ax3 = fig.add_axes([0.04, 0.05, 0.44, 0.45])
    ax3.axis('off')

    comp_text = "STATION SET COMPARISON\n"
    comp_text += "=" * 52 + "\n\n"
    comp_text += f"{'Configuration':<16} {'Stns':>4} {'Coverage':>8} {'Max Gap':>10} {'Status':>10}\n"
    comp_text += "-" * 52 + "\n"

    for result in config_results:
        num_stations = len(result['stations'])
        coverage = result['pct_coverage']
        max_gap = result['max_gap']
        status = "MEETS REQ" if result['revs_without'] == 0 else f"{result['revs_without']} gaps"
        comp_text += f"{result['name']:<16} {num_stations:>4} {coverage:>7.1f}% {max_gap:>9.1f}m {status:>10}\n"

    comp_text += "\n"

    # Show details of minimum set
    min_result = next((r for r in config_results if r['name'] == 'Minimum Set'), None)
    if min_result:
        comp_text += f"MINIMUM SET DETAILS:\n"
        comp_text += f"  Stations: {', '.join(min_result['stations'])}\n"
        if min_result['revs_without'] == 0:
            comp_text += f"  Status: MEETS 1-pass-per-rev requirement\n"
        else:
            comp_text += f"  Status: {min_result['revs_without']} revs without contact\n"
            comp_text += f"  Gap Analysis:\n"
            comp_text += f"    Max Gap: {min_result['max_gap']:.1f} min ({min_result['max_gap']/period_min:.2f} revs)\n"
            comp_text += f"    Avg Gap: {min_result['avg_gap']:.1f} min\n"

    ax3.text(0, 0.98, comp_text, transform=ax3.transAxes, fontsize=9,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#e8f5e9', alpha=0.8))

    # Bottom right: Recommendations
    ax4 = fig.add_axes([0.52, 0.05, 0.44, 0.45])
    ax4.axis('off')

    rec_text = "RECOMMENDATIONS\n"
    rec_text += "=" * 48 + "\n\n"

    # Find best configuration that meets requirement
    meeting_req = [r for r in config_results if r['revs_without'] == 0]

    if meeting_req:
        best = min(meeting_req, key=lambda x: len(x['stations']))
        rec_text += f"Minimum to meet 1-pass-per-rev: {len(best['stations'])} stations\n"
        rec_text += f"  Stations: {', '.join(best['stations'])}\n\n"

        # Cost savings
        savings = len(ground_stations) - len(best['stations'])
        rec_text += f"Potential Reduction:\n"
        rec_text += f"  Current: {len(ground_stations)} stations\n"
        rec_text += f"  Minimum: {len(best['stations'])} stations\n"
        rec_text += f"  Reduction: {savings} stations ({100*savings/len(ground_stations):.0f}%)\n\n"
    else:
        rec_text += "WARNING: No configuration meets 1-pass-per-rev!\n\n"
        # Find best available
        best_available = min(config_results, key=lambda x: x['revs_without'])
        rec_text += f"Best Available: {best_available['name']}\n"
        rec_text += f"  Coverage: {best_available['pct_coverage']:.1f}%\n"
        rec_text += f"  Gaps: {best_available['revs_without']} revs without contact\n"
        rec_text += f"  Max Gap: {best_available['max_gap']:.1f} min\n\n"

    # Gap distribution for full network
    if full_network['gaps']:
        gaps_over_1rev = len([g for g in full_network['gaps'] if g > period_min])
        gaps_over_2rev = len([g for g in full_network['gaps'] if g > 2 * period_min])
        rec_text += f"Full Network Gap Analysis:\n"
        rec_text += f"  Total Gaps: {len(full_network['gaps'])}\n"
        rec_text += f"  Gaps > 1 rev ({period_min:.0f} min): {gaps_over_1rev}\n"
        rec_text += f"  Gaps > 2 rev ({2*period_min:.0f} min): {gaps_over_2rev}\n"

    ax4.text(0, 0.98, rec_text, transform=ax4.transAxes, fontsize=9,
             verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#fce4ec', alpha=0.8))

    fig.suptitle('Ground Station Optimization - TT&C Coverage Analysis',
                 fontsize=14, fontweight='bold', y=0.98)

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {
        'min_stations': len(min_stations),
        'full_network_coverage': full_network['pct_coverage'],
        'requirement_met': full_network['revs_without'] == 0
    }


def add_slide2_results(prs, config, runtime_stats: Dict):
    """Add simulation results as text slide directly to PowerPoint."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Simulation Results"
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Column 1: Coverage Results
    left_text = []
    left_text.append("COVERAGE RESULTS")
    left_text.append(f"  Total Accesses: {runtime_stats.get('total_accesses', 'N/A')}")
    left_text.append(f"  Accesses/Day: {runtime_stats.get('accesses_per_day', 'N/A')}")
    left_text.append(f"  Valid Contacts: {runtime_stats.get('valid_contacts', 'N/A')}")
    left_text.append("")
    left_text.append("TT&C PASS GAPS")
    left_text.append(f"  Total TT&C Contacts: {runtime_stats.get('ttc_total_contacts', 'N/A')}")
    left_text.append(f"  Revs Without Contact: {runtime_stats.get('ttc_revs_without', 'N/A')}")
    left_text.append(f"  Max Gap: {runtime_stats.get('ttc_max_gap', 'N/A')}")
    left_text.append(f"  Avg Gap: {runtime_stats.get('ttc_avg_gap', 'N/A')}")
    left_text.append("")
    left_text.append("KA-BAND DOWNLINK (POST-COLLECT)")
    left_text.append(f"  Viable Downlinks: {runtime_stats.get('ka_viable_count', 'N/A')} ({runtime_stats.get('ka_viable_pct', 'N/A')})")
    left_text.append(f"  Median Latency: {runtime_stats.get('ka_median_latency', 'N/A')}")

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(5), Inches(6))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(left_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        if line and not line.startswith(' '):
            p.font.bold = True

    # Column 2: Runtime info
    right_text = []
    right_text.append("RUNTIME")
    right_text.append(f"  Processing Time: {runtime_stats.get('processing_time', 'N/A')}")

    right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.0), Inches(6), Inches(6))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(right_text):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        if line and not line.startswith(' '):
            p.font.bold = True

    return slide


def build_presentation(plots_dir: Path, config, pptx_path: Path, tle_data: Dict = None, targets_gdf = None, runtime_stats: Dict = None):
    """Build the PowerPoint presentation from generated plots."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    def add_image_slide(image_path: Path, title: str = None):
        slide = prs.slides.add_slide(blank_layout)
        if image_path.exists():
            slide.shapes.add_picture(str(image_path), Inches(0.2), Inches(0.4), width=Inches(12.9))
        if title:
            title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.05), Inches(12.9), Inches(0.35))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(18)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        return slide

    # Slide 1: Mission Parameters (text-based)
    if tle_data is not None:
        add_slide1_text(prs, config, tle_data, targets_gdf)
    else:
        add_image_slide(plots_dir / 'slide1_mission_params.png', 'Mission Parameters')

    # Slide 2: Simulation Results (text-based)
    if runtime_stats is not None:
        add_slide2_results(prs, config, runtime_stats)

    # Slide 3: Ground Stations
    add_image_slide(plots_dir / 'slide2_ground_stations.png', 'Ground Station Network')

    # Slide 4: Single Orbit
    add_image_slide(plots_dir / 'slide3_single_orbit.png', 'Single Orbit Ground Track')

    # Slide 5: 24-Hour Ground Tracks
    if (plots_dir / 'slide_1day_tracks.png').exists():
        add_image_slide(plots_dir / 'slide_1day_tracks.png', '24-Hour Ground Track Coverage')

    # Slide 6: Ground Track Walking
    if (plots_dir / 'slide_gt_walking.png').exists():
        add_image_slide(plots_dir / 'slide_gt_walking.png', 'Ground Track Walking Analysis')

    # Slide 7: TT&C Coverage Analysis
    if (plots_dir / 'slide_ttc_coverage.png').exists():
        add_image_slide(plots_dir / 'slide_ttc_coverage.png', 'TT&C Coverage Analysis')

    # Slide 8: Access Results
    if (plots_dir / 'slide4_access_results.png').exists():
        add_image_slide(plots_dir / 'slide4_access_results.png', 'Access Windows Overview')

    # Slide 9: Regional Capacity Analysis
    if (plots_dir / 'slide_regional_capacity.png').exists():
        add_image_slide(plots_dir / 'slide_regional_capacity.png', 'Regional Capacity Analysis')

    # Slide 10: Access with Cones
    if (plots_dir / 'slide5_access_with_cones.png').exists():
        add_image_slide(plots_dir / 'slide5_access_with_cones.png', 'Access Windows with 60-min Post-Collect Tracks')

    # Slides: Broadside Collect
    for i in range(1, 4):
        if (plots_dir / f'slide_broadside_{i}.png').exists():
            add_image_slide(plots_dir / f'slide_broadside_{i}.png', f'Broadside Collect Opportunity {i}')

    # Slides: 3h Post-Collect
    for i in range(1, 4):
        if (plots_dir / f'slide_postcollect_{i}.png').exists():
            add_image_slide(plots_dir / f'slide_postcollect_{i}.png', f'Post-Collect Analysis {i} (3h)')

    # Slide: Ka-Band Coverage Analysis (payload downlink delay)
    if (plots_dir / 'slide_ka_coverage.png').exists():
        add_image_slide(plots_dir / 'slide_ka_coverage.png', 'Ka-Band Post-Collect Coverage Analysis')

    # Slides: Pre-Collect
    for i in range(1, 4):
        if (plots_dir / f'slide_precollect_{i}.png').exists():
            add_image_slide(plots_dir / f'slide_precollect_{i}.png', f'Pre-Collect Analysis {i} (Uplink-to-Collect)')

    # Slide: Uplink-to-Collect Summary
    if (plots_dir / 'slide_uplink_summary.png').exists():
        add_image_slide(plots_dir / 'slide_uplink_summary.png', 'Pre-Collect TT&C Analysis')


    # Slide 19: Ground Station Optimization
    if (plots_dir / 'slide_optimization.png').exists():
        add_image_slide(plots_dir / 'slide_optimization.png', 'Ground Station Optimization Analysis')

    prs.save(str(pptx_path))
    print(f"Presentation saved to {pptx_path}")


def main():
    import argparse

    parser = argparse.ArgumentParser(description='Generate full presentation')
    parser.add_argument('--config', required=True, help='Path to config YAML')
    parser.add_argument('--results-dir', required=True, help='Path to results directory')
    parser.add_argument('--seed', type=int, default=42, help='Random seed')
    parser.add_argument('--parallel', action='store_true', help='Generate plots in parallel')
    parser.add_argument('--workers', type=int, default=4, help='Number of parallel workers')
    args = parser.parse_args()

    start_time = time.time()

    config = load_config(args.config)
    results_dir = Path(args.results_dir)
    plots_dir = results_dir / config.plots_subdir
    plots_dir.mkdir(parents=True, exist_ok=True)

    # Load data
    excel_path = results_dir / config.excel_filename
    access_df = pd.read_excel(excel_path, sheet_name='Access_Windows')
    print(f"Loaded {len(access_df)} access windows")

    # Setup TLE data
    orbit_config = config.raw_config.get('orbit', {})
    ltdn_hours = orbit_config.get('ltdn_hours')
    auto_sso = orbit_config.get('auto_sso', False)

    tle_data = {}
    for sat in config.satellites:
        inc_deg = sat.inclination_deg
        raan_deg = sat.raan_deg
        if auto_sso and ltdn_hours:
            inc_deg = calculate_sso_inclination(sat.altitude_km)
            raan_deg = ltdn_to_raan(ltdn_hours, config.start_datetime)

        tle = create_tle_data(
            sat_id=sat.sat_id, inclination_deg=inc_deg,
            altitude_km=sat.altitude_km, raan_deg=raan_deg,
            epoch=config.start_datetime,
        )
        tle_data[sat.sat_id] = tle

    # Ground stations
    ground_stations = [{'name': gs.name, 'lat': gs.lat, 'lon': gs.lon,
                        'min_elevation_deg': gs.min_elevation_deg, 'ka_band': gs.ka_capable}
                       for gs in config.ground_stations]

    sat_alt_km = config.satellites[0].altitude_km

    # Load targets GeoDataFrame
    targets_config = config.raw_config.get('targets', {})
    geojson_path = targets_config.get('geojson_path')
    targets_gdf = None
    if geojson_path:
        full_geojson_path = Path(args.config).parent.parent / geojson_path
        if full_geojson_path.exists():
            targets_gdf = gpd.read_file(full_geojson_path)
            print(f"Loaded {len(targets_gdf)} targets from {geojson_path}")
        else:
            # Try relative to current directory
            if Path(geojson_path).exists():
                targets_gdf = gpd.read_file(geojson_path)
                print(f"Loaded {len(targets_gdf)} targets from {geojson_path}")

    # Load runtime stats from Excel
    runtime_stats = {}
    try:
        # Get KPIs - format is KPI, Value, Unit columns
        kpi_df = pd.read_excel(excel_path, sheet_name='Coverage_KPIs')
        kpi_dict = dict(zip(kpi_df['KPI'], kpi_df['Value']))

        total_accesses = int(kpi_dict.get('Total Accesses', len(access_df)))
        valid_contacts = int(kpi_dict.get('Total Valid Contacts', 0))
        runtime_stats['total_accesses'] = total_accesses
        runtime_stats['valid_contacts'] = valid_contacts
        runtime_stats['accesses_per_day'] = f"{kpi_dict.get('Accesses Per Day', total_accesses / config.duration_days):.1f}"
        runtime_stats['contact_success'] = f"{(valid_contacts / total_accesses * 100) if total_accesses > 0 else 0:.1f}%"
        runtime_stats['data_collected'] = f"{kpi_dict.get('Total Data Collected', 0):.1f} GB"
        runtime_stats['data_downlinked'] = f"{kpi_dict.get('Total Data Downlinked', 0):.1f} GB"

        # Get downlink delay stats
        delay_df = pd.read_excel(excel_path, sheet_name='Downlink_Delay_Summary')
        if len(delay_df) > 0:
            runtime_stats['downlink_delay_median'] = f"{delay_df['Median (min)'].iloc[0]:.1f} min"
            runtime_stats['downlink_delay_p95'] = f"{delay_df['P95 (min)'].iloc[0]:.1f} min"

        # Get station-keeping stats if available
        try:
            sk_df = pd.read_excel(excel_path, sheet_name='Station_Keeping')
            if len(sk_df) > 0:
                sk_dict = dict(zip(sk_df['Parameter'], sk_df['Value']))
                runtime_stats['thruster'] = sk_dict.get('Thruster', 'N/A')
                runtime_stats['delta_v'] = f"{sk_dict.get('Delta-V Required (m/s)', 0):.2f} m/s"
                runtime_stats['propellant'] = f"{sk_dict.get('Propellant Mass (kg)', 0):.2f} kg"
                runtime_stats['propellant_margin'] = f"{sk_dict.get('Propellant with Margin (kg)', 0):.2f} kg"
                firing_hrs = sk_dict.get('Total Firing Time (hours)', 0)
                duty = sk_dict.get('Duty Cycle (%)', 0)
                runtime_stats['firing_time'] = f"{firing_hrs:.1f} hrs ({duty:.1f}% duty)"
                runtime_stats['annual_propellant'] = f"{sk_dict.get('Annualized Propellant (kg/yr)', 0):.1f} kg/yr"
        except Exception:
            pass  # Station-keeping sheet may not exist

    except Exception as e:
        print(f"Warning: Could not load runtime stats: {e}")
        runtime_stats['total_accesses'] = len(access_df)

    # Select 3 random accesses for detailed slides
    random.seed(args.seed)
    num_accesses = min(3, len(access_df))
    random_indices = random.sample(range(len(access_df)), num_accesses)

    # Get off-nadir limits from config
    off_nadir_min = config.imaging_modes[0].off_nadir_min_deg if config.imaging_modes else 0.0
    off_nadir_max = config.imaging_modes[0].off_nadir_max_deg if config.imaging_modes else 30.0

    # Define all plot generation tasks
    plot_tasks = [
        ('Ground Stations', generate_slide2_gs_map,
         (config, ground_stations, sat_alt_km, targets_gdf, plots_dir / 'slide2_ground_stations.png'),
         {'excel_path': excel_path}),
        ('Single Orbit', generate_slide3_single_orbit,
         (tle_data, config, plots_dir / 'slide3_single_orbit.png'), {}),
        ('Access Results', generate_slide4_access_results,
         (access_df, tle_data, plots_dir / 'slide4_access_results.png'), {}),
        ('Access with Cones', generate_slide5_access_with_cones,
         (access_df, tle_data, ground_stations, sat_alt_km, plots_dir / 'slide5_access_with_cones.png'), {}),
        ('Downlink Delay Summary', generate_downlink_delay_summary_plot,
         (excel_path, plots_dir / 'slide_downlink_delay_summary.png'), {}),
        ('Uplink Summary', generate_uplink_to_collect_summary,
         (access_df, tle_data, ground_stations, sat_alt_km, plots_dir / 'slide_uplink_summary.png'), {}),
        ('Optimization', generate_optimization_slide,
         (excel_path, plots_dir / 'slide_optimization.png', config), {}),
    ]

    # Add broadside collect tasks
    for i, idx in enumerate(random_indices):
        plot_tasks.append((
            f'Broadside {i+1}', generate_broadside_collect_plot,
            (access_df.iloc[idx], tle_data, ground_stations, sat_alt_km, plots_dir / f'slide_broadside_{i+1}.png'),
            {'off_nadir_min': off_nadir_min, 'off_nadir_max': off_nadir_max}
        ))

    # Add post-collect tasks
    for i, idx in enumerate(random_indices):
        plot_tasks.append((
            f'Post-Collect {i+1}', generate_collect_diagram_3h,
            (access_df.iloc[idx], tle_data, ground_stations, sat_alt_km, plots_dir / f'slide_postcollect_{i+1}.png'),
            {}
        ))

    # Add pre-collect tasks
    for i, idx in enumerate(random_indices):
        plot_tasks.append((
            f'Pre-Collect {i+1}', generate_pre_collect_diagram,
            (access_df.iloc[idx], tle_data, ground_stations, sat_alt_km, plots_dir / f'slide_precollect_{i+1}.png'),
            {}
        ))

    # Execute plot generation
    print(f"\nGenerating {len(plot_tasks)} plots" + (" in parallel..." if args.parallel else "..."))

    if args.parallel:
        # Parallel execution
        def execute_task(task):
            name, func, task_args, kwargs = task
            try:
                func(*task_args, **kwargs)
                return (name, True, None)
            except Exception as e:
                return (name, False, str(e))

        with ThreadPoolExecutor(max_workers=args.workers) as executor:
            futures = {executor.submit(execute_task, task): task[0] for task in plot_tasks}
            completed = 0
            for future in as_completed(futures):
                name, success, error = future.result()
                completed += 1
                status = "OK" if success else f"FAILED: {error}"
                print(f"  [{completed}/{len(plot_tasks)}] {name}: {status}")
    else:
        # Sequential execution with progress
        for i, (name, func, task_args, kwargs) in enumerate(plot_tasks):
            print(f"  [{i+1}/{len(plot_tasks)}] Generating {name}...")
            try:
                func(*task_args, **kwargs)
            except Exception as e:
                print(f"    ERROR: {e}")

    # Record timing
    plot_time = time.time() - start_time
    runtime_stats['processing_time'] = f"{plot_time:.1f} s"

    # Build presentation
    print("\nBuilding PowerPoint presentation...")
    pptx_path = results_dir / config.ppt_filename
    build_presentation(plots_dir, config, pptx_path, tle_data=tle_data, targets_gdf=targets_gdf, runtime_stats=runtime_stats)

    total_time = time.time() - start_time
    print(f"\nDone! Total time: {total_time:.1f}s")


if __name__ == '__main__':
    main()
