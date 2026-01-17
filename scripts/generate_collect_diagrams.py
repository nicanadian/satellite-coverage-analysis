#!/usr/bin/env python3
"""
Generate ground track + comm cone visualization diagrams for individual collects.
Shows extended ground track, imaging segment, comm cones, and contact windows.
"""

import sys
import random
from pathlib import Path
from datetime import timedelta

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import cartopy.crs as ccrs
import cartopy.feature as cfeature
from matplotlib.colors import LinearSegmentedColormap
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from src.vleo_eo.orbits import create_tle_data, calculate_ground_track
from src.vleo_eo.contacts import calculate_comm_range_km, calculate_ground_distance_km
from src.vleo_eo.config import load_config
from shapely.geometry import Polygon, Point
import geopandas as gpd


def create_comm_cone_geodesic(lon: float, lat: float, min_elevation_deg: float, sat_altitude_km: float):
    """
    Create a communication cone polygon using proper geodesic geometry.

    Uses a local azimuthal equidistant projection centered on the ground station
    to create an accurate circle, then transforms back to WGS84 for display.
    This correctly handles high-latitude distortion on Mercator projections.

    Returns a GeoDataFrame with the comm cone geometry.
    """
    earth_radius_km = 6378.137

    # Calculate maximum ground range
    elevation_rad = np.radians(min_elevation_deg)
    central_angle = np.arccos(
        earth_radius_km * np.cos(elevation_rad) / (earth_radius_km + sat_altitude_km)
    ) - elevation_rad
    max_ground_range_m = earth_radius_km * central_angle * 1000  # Convert to meters

    # Create a point at the ground station location
    gs_point = Point(lon, lat)
    gs_gdf = gpd.GeoDataFrame({'geometry': [gs_point]}, crs='EPSG:4326')

    # Use Azimuthal Equidistant projection centered on the ground station
    # This projection preserves distances from the center point
    aeqd_crs = f"+proj=aeqd +lat_0={lat} +lon_0={lon} +x_0=0 +y_0=0 +datum=WGS84 +units=m +no_defs"

    # Transform to the local projection
    gs_projected = gs_gdf.to_crs(aeqd_crs)

    # Buffer by the comm range (creates a true circle in the local projection)
    buffered = gs_projected.buffer(max_ground_range_m)

    # Transform back to WGS84
    comm_cone_gdf = gpd.GeoDataFrame({'geometry': buffered}, crs=aeqd_crs).to_crs('EPSG:4326')

    return comm_cone_gdf.geometry.iloc[0]


def generate_collect_diagram(
    access_row: pd.Series,
    tle_data: dict,
    ground_stations: list,
    ka_stations: list,
    output_path: Path,
    config: dict,
) -> dict:
    """
    Generate a single collect diagram showing ground track, comm cones, and contacts.

    Returns dict with statistics for the plot.
    """
    sat_id = access_row['Satellite']
    start_time = pd.to_datetime(access_row['Start_Time'])
    if start_time.tzinfo is None:
        start_time = start_time.tz_localize('UTC')

    # Get target location
    aoi_lat = access_row['AOI_Lat']
    aoi_lon = access_row['AOI_Lon']

    # Create figure with extra space on right for legend
    fig = plt.figure(figsize=(18, 10), dpi=150)
    ax = plt.axes([0.05, 0.05, 0.70, 0.90], projection=ccrs.PlateCarree())

    # Add map features
    land_color = config.get('land_color', '#fafaf9')
    ocean_color = config.get('ocean_color', '#d4dbdc')
    border_color = config.get('border_color', '#ebd6d9')

    ax.add_feature(cfeature.OCEAN, facecolor=ocean_color, zorder=0)
    ax.add_feature(cfeature.LAND, facecolor=land_color, zorder=1)
    ax.add_feature(cfeature.BORDERS, linewidth=0.5, edgecolor=border_color, zorder=2)
    ax.add_feature(cfeature.COASTLINE, linewidth=0.5, edgecolor='gray', zorder=2)
    ax.gridlines(draw_labels=True, alpha=0.3)

    # Calculate extended ground track (3 hours after collect)
    extended_track = calculate_ground_track(
        tle_data, sat_id, start_time.to_pydatetime(),
        duration=timedelta(hours=3), step_s=30.0
    )

    if not extended_track:
        print(f"Warning: Could not calculate ground track for access")
        return None

    ext_lons, ext_lats = zip(*extended_track)

    # Plot extended ground track in gray
    ax.plot(ext_lons, ext_lats, color='gray', linewidth=1.5, alpha=0.6,
            transform=ccrs.PlateCarree(), label='Ground Track (3h)', zorder=3)

    # Plot target location as blue X
    ax.plot(aoi_lon, aoi_lat, 'x', color='blue', markersize=8,
            markeredgewidth=2, transform=ccrs.PlateCarree(),
            label='Target', zorder=7)

    # Create colormap for contact times (green -> yellow -> orange -> red)
    colors = ['#00ff00', '#2ecc71', '#f1c40f', '#e67e22', '#e74c3c']
    custom_cmap = LinearSegmentedColormap.from_list("contact_time", colors, N=256)

    # Get satellite altitude
    sat_alt_km = tle_data[sat_id].altitude_km

    # Track contact statistics
    contact_stats = []
    first_ka_contact = None

    # Process each ground station
    for gs in ground_stations:
        gs_lat = gs['lat']
        gs_lon = gs['lon']
        gs_name = gs['name']
        min_elev = gs.get('min_elevation_deg', 5.0)
        is_ka = gs.get('ka_band', False)

        # Plot ground station - green triangle for Ka, black triangle for TT&C
        marker_color = 'green' if is_ka else 'black'
        ax.plot(gs_lon, gs_lat, '^', color=marker_color, markersize=8,
                transform=ccrs.PlateCarree(), zorder=8)

        # Create and plot comm cone using proper geodesic geometry
        # This correctly handles high-latitude distortion on Mercator projections
        # Use different elevation masks: Ka = 10 deg, TT&C = 5 deg
        cone_color = 'green' if is_ka else 'black'
        cone_elev = 10.0 if is_ka else 5.0

        try:
            comm_cone = create_comm_cone_geodesic(gs_lon, gs_lat, cone_elev, sat_alt_km)

            if comm_cone is not None and comm_cone.is_valid:
                ax.add_geometries(
                    [comm_cone], crs=ccrs.PlateCarree(),
                    facecolor='none',
                    edgecolor=cone_color,
                    alpha=0.6,
                    linewidth=0.8, zorder=2
                )
        except Exception as e:
            # Skip if geometry fails (shouldn't happen with geodesic approach)
            print(f"Warning: Could not create comm cone for {gs_name}: {e}")

        # Check for contact along the extended track
        # Use different elevation masks: Ka = 10 deg, TT&C = 5 deg
        contact_elev = 10.0 if is_ka else 5.0
        max_range_km = calculate_comm_range_km(contact_elev, sat_alt_km)

        # Find when satellite enters/exits comm cone
        in_contact = False
        contact_start_idx = None
        contacts = []

        for i, (lon, lat) in enumerate(extended_track):
            dist = calculate_ground_distance_km(lat, lon, gs_lat, gs_lon)
            currently_in = dist <= max_range_km

            if currently_in and not in_contact:
                contact_start_idx = i
                in_contact = True
            elif not currently_in and in_contact:
                contacts.append((contact_start_idx, i - 1))
                in_contact = False

        if in_contact:
            contacts.append((contact_start_idx, len(extended_track) - 1))

        # Plot contact segments
        for start_idx, end_idx in contacts:
            if end_idx <= start_idx:
                continue

            contact_lons = [extended_track[j][0] for j in range(start_idx, end_idx + 1)]
            contact_lats = [extended_track[j][1] for j in range(start_idx, end_idx + 1)]

            # Calculate time to contact from collect start
            time_to_contact_min = start_idx * 30 / 60  # 30 second steps
            contact_duration_min = (end_idx - start_idx) * 30 / 60

            # Color based on time to contact (0-180 min range)
            norm_time = min(time_to_contact_min / 180.0, 1.0)
            contact_color = custom_cmap(norm_time)

            # Record contact stats (including plot info for legend)
            contact_stats.append({
                'station': gs_name,
                'is_ka': is_ka,
                'time_to_contact_min': time_to_contact_min,
                'duration_min': contact_duration_min,
                'color': contact_color,
                'lons': contact_lons,
                'lats': contact_lats,
            })

            # Track first Ka contact
            if is_ka and (first_ka_contact is None or time_to_contact_min < first_ka_contact['time_to_contact_min']):
                first_ka_contact = contact_stats[-1]

    # Sort contacts by time to contact and plot them
    contact_stats.sort(key=lambda x: x['time_to_contact_min'])

    for contact in contact_stats:
        ax.plot(contact['lons'], contact['lats'], color=contact['color'], linewidth=3,
                transform=ccrs.PlateCarree(), zorder=4)

    # Set map extent based on ground track
    all_lons = list(ext_lons)
    all_lats = list(ext_lats)

    lon_margin = 20
    lat_margin = 15
    ax.set_extent([
        min(all_lons) - lon_margin, max(all_lons) + lon_margin,
        max(min(all_lats) - lat_margin, -90), min(max(all_lats) + lat_margin, 90)
    ], crs=ccrs.PlateCarree())

    # Build title with statistics
    title_lines = [
        f"Collect: {access_row.get('Imaging_Mode', 'Nominal')} - Satellite {sat_id}",
        f"Target: ({aoi_lat:.1f}°, {aoi_lon:.1f}°) | Start: {start_time.strftime('%Y-%m-%d %H:%M:%S')} UTC",
    ]

    if first_ka_contact:
        title_lines.append(
            f"First Ka Contact: {first_ka_contact['station']} at C+{first_ka_contact['time_to_contact_min']:.1f}min "
            f"(duration: {first_ka_contact['duration_min']:.1f}min)"
        )
    else:
        title_lines.append("No Ka contact within 3 hours")

    # Add contact summary
    ka_contacts = [c for c in contact_stats if c['is_ka']]
    all_contacts = contact_stats
    title_lines.append(
        f"Total contacts: {len(all_contacts)} ({len(ka_contacts)} Ka-band)"
    )

    ax.set_title('\n'.join(title_lines), fontsize=11, fontweight='bold')

    # Add legend
    from matplotlib.patches import Patch
    from matplotlib.lines import Line2D

    legend_elements = [
        Line2D([0], [0], color='gray', linewidth=1.5, alpha=0.6, label='Ground Track (3h)'),
        Line2D([0], [0], marker='x', color='blue', markersize=8, linestyle='None',
               markeredgewidth=2, label='Target'),
        Line2D([0], [0], marker='^', color='green', markersize=8, linestyle='None', label='Ka Station'),
        Patch(facecolor='none', edgecolor='green', alpha=0.5, linewidth=1.5, label='Ka Comm Cone'),
        Line2D([0], [0], marker='^', color='black', markersize=8, linestyle='None', label='TT&C Station'),
        Patch(facecolor='none', edgecolor='black', alpha=0.5, linewidth=1.5, label='TT&C Comm Cone'),
    ]

    # Add contact entries to legend with time info
    for contact in contact_stats:
        ka_marker = " (Ka)" if contact['is_ka'] else ""
        label = f"{contact['station']}{ka_marker}: C+{contact['time_to_contact_min']:.1f}m, {contact['duration_min']:.1f}m"
        legend_elements.append(
            Line2D([0], [0], color=contact['color'], linewidth=3, label=label)
        )

    # Position legend to the right of the map
    ax.legend(handles=legend_elements, loc='upper left', bbox_to_anchor=(1.02, 1.0),
              fontsize=8, frameon=True, fancybox=True)

    plt.savefig(output_path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()

    return {
        'path': output_path,
        'satellite': sat_id,
        'target': (aoi_lat, aoi_lon),
        'start_time': start_time,
        'first_ka_station': first_ka_contact['station'] if first_ka_contact else None,
        'first_ka_time_min': first_ka_contact['time_to_contact_min'] if first_ka_contact else None,
        'total_contacts': len(all_contacts),
        'ka_contacts': len(ka_contacts),
    }


def add_diagrams_to_pptx(pptx_path: Path, diagram_paths: list, diagram_stats: list):
    """Add collect diagrams to the beginning of an existing PowerPoint."""
    from copy import deepcopy
    import shutil

    # Create backup
    backup_path = pptx_path.with_suffix('.pptx.bak')
    shutil.copy(pptx_path, backup_path)

    # Load existing presentation
    existing_prs = Presentation(str(pptx_path))

    # Create new presentation with same slide size
    new_prs = Presentation()
    new_prs.slide_width = existing_prs.slide_width
    new_prs.slide_height = existing_prs.slide_height

    # Get blank layout
    blank_layout = new_prs.slide_layouts[6]

    # 1. Add summary slide
    summary_slide = new_prs.slides.add_slide(blank_layout)

    # Add title
    title_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Sample Collect Analysis"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"Ground Track + Communication Cone Analysis for {len(diagram_stats)} Random Collects"
    subtitle_para.font.size = Pt(16)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Add stats summary
    stats_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    stats_frame = stats_box.text_frame

    for i, stat in enumerate(diagram_stats):
        if stat:
            p = stats_frame.add_paragraph()
            ka_info = f"{stat['first_ka_station']} at C+{stat['first_ka_time_min']:.1f}min" if stat['first_ka_time_min'] else "None within 3h"
            p.text = f"Collect {i+1}: Target ({stat['target'][0]:.1f}°, {stat['target'][1]:.1f}°) | First Ka: {ka_info}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)

    # 2. Add each diagram as a new slide
    for i, (path, stat) in enumerate(zip(diagram_paths, diagram_stats)):
        if not path.exists():
            continue

        slide = new_prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(path), Inches(0.15), Inches(0.15), width=Inches(9.7))

    # 3. Copy slides from existing presentation
    # We can't directly copy slides between presentations, so we save and load
    # Instead, let's append the new slides to the existing presentation

    # Actually, let's use a simpler approach: add to existing and reorder
    prs = Presentation(str(backup_path))
    blank_layout = prs.slide_layouts[6]

    # Count existing slides
    num_existing = len(prs.slides)

    # Add summary slide
    summary_slide = prs.slides.add_slide(blank_layout)

    title_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Sample Collect Analysis"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    subtitle_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = f"Ground Track + Communication Cone Analysis for {len(diagram_stats)} Random Collects"
    subtitle_para.font.size = Pt(16)
    subtitle_para.alignment = PP_ALIGN.CENTER

    stats_box = summary_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(3))
    stats_frame = stats_box.text_frame

    for i, stat in enumerate(diagram_stats):
        if stat:
            p = stats_frame.add_paragraph()
            ka_info = f"{stat['first_ka_station']} at C+{stat['first_ka_time_min']:.1f}min" if stat['first_ka_time_min'] else "None within 3h"
            p.text = f"Collect {i+1}: Target ({stat['target'][0]:.1f}°, {stat['target'][1]:.1f}°) | First Ka: {ka_info}"
            p.font.size = Pt(14)
            p.space_after = Pt(8)

    # Add diagram slides
    for path in diagram_paths:
        if not path.exists():
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(str(path), Inches(0.15), Inches(0.15), width=Inches(9.7))

    # Reorder: move new slides to after the first slide (title)
    # New slides are at positions num_existing to end
    # We want them at positions 1 to (1 + num_new_slides)
    num_new_slides = 1 + len(diagram_paths)  # summary + diagrams

    # Get the slide ID list
    sldIdLst = prs.slides._sldIdLst

    # Extract new slide IDs (they are at the end)
    new_slide_ids = list(sldIdLst)[-num_new_slides:]

    # Remove them from the end
    for _ in range(num_new_slides):
        del sldIdLst[-1]

    # Insert them after position 0 (after title slide)
    for i, slide_id in enumerate(new_slide_ids):
        sldIdLst.insert(1 + i, slide_id)

    prs.save(str(pptx_path))
    print(f"Added {len(diagram_paths)} collect diagrams to {pptx_path}")

    # Clean up backup
    backup_path.unlink()


def main():
    import argparse

    parser = argparse.ArgumentParser(description='Generate collect diagram visualizations')
    parser.add_argument('--config', required=True, help='Path to config YAML file')
    parser.add_argument('--results-dir', required=True, help='Path to results directory')
    parser.add_argument('--num-collects', type=int, default=3, help='Number of random collects to visualize')
    parser.add_argument('--seed', type=int, default=42, help='Random seed for reproducibility')
    args = parser.parse_args()

    # Load config
    config = load_config(args.config)
    results_dir = Path(args.results_dir)
    plots_dir = results_dir / config.plots_subdir

    # Load access windows from Excel
    excel_path = results_dir / config.excel_filename
    if not excel_path.exists():
        print(f"Error: Excel file not found at {excel_path}")
        sys.exit(1)

    access_df = pd.read_excel(excel_path, sheet_name='Access_Windows')
    print(f"Loaded {len(access_df)} access windows")

    if len(access_df) == 0:
        print("No access windows to visualize")
        sys.exit(0)

    # Select random accesses
    random.seed(args.seed)
    num_to_select = min(args.num_collects, len(access_df))
    random_indices = random.sample(range(len(access_df)), num_to_select)
    print(f"Selected {num_to_select} random collects: indices {random_indices}")

    # Create TLE data
    from src.vleo_eo.propulsion import ltdn_to_raan, calculate_sso_inclination

    tle_data = {}
    orbit_config = config.raw_config.get('orbit', {})
    ltdn_hours = orbit_config.get('ltdn_hours')
    auto_sso = orbit_config.get('auto_sso', False)

    for sat in config.satellites:
        inc_deg = sat.inclination_deg
        raan_deg = sat.raan_deg

        if auto_sso and ltdn_hours is not None:
            inc_deg = calculate_sso_inclination(sat.altitude_km)
            raan_deg = ltdn_to_raan(ltdn_hours, config.start_datetime)

        tle = create_tle_data(
            sat_id=sat.sat_id,
            inclination_deg=inc_deg,
            altitude_km=sat.altitude_km,
            raan_deg=raan_deg,
            epoch=config.start_datetime,
        )
        tle_data[sat.sat_id] = tle

    # Convert ground stations to list of dicts
    ground_stations = []
    ka_stations = []
    for gs in config.ground_stations:
        gs_dict = {
            'name': gs.name,
            'lat': gs.lat,
            'lon': gs.lon,
            'min_elevation_deg': gs.min_elevation_deg,
            'ka_band': gs.ka_capable,
        }
        ground_stations.append(gs_dict)
        if gs.ka_capable:
            ka_stations.append(gs_dict)

    # Generate diagrams
    diagram_paths = []
    diagram_stats = []

    raw_config = config.raw_config
    plot_config = {
        'land_color': raw_config.get('land_color', '#fafaf9'),
        'ocean_color': raw_config.get('ocean_color', '#d4dbdc'),
        'border_color': raw_config.get('border_color', '#ebd6d9'),
    }

    for i, idx in enumerate(random_indices):
        access_row = access_df.iloc[idx]
        output_path = plots_dir / f'collect_diagram_{i+1}.png'

        print(f"\nGenerating diagram {i+1}/{num_to_select}...")
        print(f"  Target: ({access_row['AOI_Lat']:.1f}°, {access_row['AOI_Lon']:.1f}°)")
        print(f"  Start: {access_row['Start_Time']}")

        stats = generate_collect_diagram(
            access_row, tle_data, ground_stations, ka_stations,
            output_path, plot_config
        )

        if stats:
            diagram_paths.append(output_path)
            diagram_stats.append(stats)
            print(f"  Saved to: {output_path}")
            if stats['first_ka_station']:
                print(f"  First Ka: {stats['first_ka_station']} at C+{stats['first_ka_time_min']:.1f}min")

    # Add to PowerPoint
    pptx_path = results_dir / config.ppt_filename
    if pptx_path.exists() and diagram_paths:
        add_diagrams_to_pptx(pptx_path, diagram_paths, diagram_stats)

    print("\nDone!")


if __name__ == '__main__':
    main()
