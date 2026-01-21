#!/usr/bin/env python3
"""
Comparison Report Generator for VLEO EO Coverage Analysis

Generates comparison visualizations and PowerPoint presentation for
comparing multiple ground station provider configurations.
"""

from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
import numpy as np
import pandas as pd

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from matplotlib.lines import Line2D

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Import comparison types
from src.vleo_eo.comparison import (
    ComparisonResult, get_provider_color, PROVIDER_COLORS,
)


def generate_coverage_circle(lat: float, lon: float, radius_km: float, num_points: int = 72) -> Tuple[List[float], List[float]]:
    """Generate circle coordinates around a point with given radius using geodesic calculation."""
    try:
        from pyproj import Geod
        geod = Geod(ellps='WGS84')
        lons = []
        lats = []
        for azimuth in np.linspace(0, 360, num_points):
            dest_lon, dest_lat, _ = geod.fwd(lon, lat, azimuth, radius_km * 1000)
            lons.append(dest_lon)
            lats.append(dest_lat)
        lons.append(lons[0])  # Close the circle
        lats.append(lats[0])
        return lons, lats
    except ImportError:
        # Fallback to simple approximation
        lons = []
        lats = []
        for i in range(num_points + 1):
            bearing = 2 * np.pi * i / num_points
            d_lat = (radius_km / 111.0) * np.cos(bearing)
            d_lon = (radius_km / (111.0 * np.cos(np.radians(lat)))) * np.sin(bearing)
            lats.append(lat + d_lat)
            lons.append(lon + d_lon)
        return lons, lats


def calculate_coverage_radius(altitude_km: float, min_elevation_deg: float) -> float:
    """Calculate coverage radius on Earth's surface for given altitude and minimum elevation."""
    earth_radius_km = 6378.137
    elev_rad = np.radians(min_elevation_deg)
    # Central angle from satellite to horizon at min elevation
    central_angle = np.arccos(earth_radius_km * np.cos(elev_rad) / (earth_radius_km + altitude_km)) - elev_rad
    coverage_radius_km = earth_radius_km * central_angle
    return coverage_radius_km


def load_targets_from_config(config_path: str) -> Optional[Any]:
    """Load targets GeoDataFrame from a config file's targets section."""
    try:
        import yaml
        import geopandas as gpd

        config_path = Path(config_path)
        if not config_path.exists():
            return None

        with open(config_path) as f:
            config = yaml.safe_load(f)

        targets_config = config.get('targets', {})
        geojson_path = targets_config.get('geojson_path')

        if not geojson_path:
            return None

        # Try relative to config directory
        full_path = config_path.parent / geojson_path
        if not full_path.exists():
            # Try relative to repo root
            full_path = config_path.parent.parent / geojson_path
        if not full_path.exists():
            # Try as absolute path
            full_path = Path(geojson_path)

        if full_path.exists():
            return gpd.read_file(full_path)
        return None
    except Exception as e:
        print(f"  Warning: Could not load targets from {config_path}: {e}")
        return None


def get_satellite_altitude_from_config(config_path: str) -> float:
    """Get satellite altitude from config file. Returns default of 350km if not found."""
    try:
        import yaml
        with open(config_path) as f:
            config = yaml.safe_load(f)
        satellites = config.get('satellites', [])
        if satellites and len(satellites) > 0:
            return satellites[0].get('altitude_km', 350.0)
        return 350.0
    except Exception:
        return 350.0

# Slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def load_ttc_data_from_excel(results_info: Dict[str, Dict]) -> Dict[str, Dict]:
    """Load TT&C coverage data from Excel files."""
    ttc_data = {}
    for label, info in results_info.items():
        results_dir = Path(info.get('results_dir', ''))
        excel_files = list(results_dir.glob('*.xlsx'))
        if not excel_files:
            continue
        try:
            df = pd.read_excel(excel_files[0], sheet_name='TTC_Coverage')
            # Convert to dict with Metric as key
            data = {}
            for _, row in df.iterrows():
                metric = row.get('Metric', '')
                value = row.get('Value', '')
                data[metric] = value
            ttc_data[label] = data
        except Exception as e:
            print(f"  Warning: Could not load TT&C data for {label}: {e}")
    return ttc_data


def load_downlink_delay_stats(results_info: Dict[str, Dict]) -> Dict[str, Dict]:
    """Load downlink delay statistics from Excel files."""
    delay_stats = {}
    for label, info in results_info.items():
        results_dir = Path(info.get('results_dir', ''))
        excel_files = list(results_dir.glob('*.xlsx'))
        if not excel_files:
            continue
        try:
            # Load raw delay data
            df = pd.read_excel(excel_files[0], sheet_name='Downlink_Delay')
            if 'Downlink_Delay_minutes' in df.columns:
                delays = df['Downlink_Delay_minutes'].dropna()
                total = len(delays)
                stats = {
                    'total': total,
                    'median': delays.median(),
                    'p90': delays.quantile(0.90),
                    'p95': delays.quantile(0.95),
                    'max': delays.max(),
                    'under_30min': (delays < 30).sum(),
                    'under_30min_pct': 100 * (delays < 30).sum() / total if total > 0 else 0,
                    'under_60min': (delays < 60).sum(),
                    'under_60min_pct': 100 * (delays < 60).sum() / total if total > 0 else 0,
                }
                delay_stats[label] = stats
        except Exception as e:
            print(f"  Warning: Could not load delay stats for {label}: {e}")
    return delay_stats


def load_optimization_data(results_info: Dict[str, Dict]) -> Dict[str, Dict]:
    """Load ground station optimization data by analyzing TT&C by station."""
    opt_data = {}
    for label, info in results_info.items():
        results_dir = Path(info.get('results_dir', ''))
        excel_files = list(results_dir.glob('*.xlsx'))
        if not excel_files:
            continue
        try:
            # Load station-by-station TT&C data
            df = pd.read_excel(excel_files[0], sheet_name='TTC_By_Station')
            # Sort by contacts descending
            df_sorted = df.sort_values('TT&C Contacts', ascending=False)

            # Get top stations
            stations = []
            for _, row in df_sorted.iterrows():
                stations.append({
                    'name': row.get('Station', ''),
                    'contacts': row.get('TT&C Contacts', 0),
                    'ka_capable': row.get('Ka-Band Capable', 'No') == 'Yes',
                })

            opt_data[label] = {
                'total_stations': len(df),
                'stations_ranked': stations,
            }
        except Exception as e:
            print(f"  Warning: Could not load optimization data for {label}: {e}")
    return opt_data


def load_gap_data_from_excel(results_info: Dict[str, Dict]) -> Dict[str, Dict]:
    """Load TT&C gap data from Excel files for periodicity analysis."""
    gap_data = {}
    for label, info in results_info.items():
        results_dir = Path(info.get('results_dir', ''))
        excel_files = list(results_dir.glob('*.xlsx'))
        if not excel_files:
            continue
        try:
            # Try to load raw gap data
            df = pd.read_excel(excel_files[0], sheet_name='TTC_Gaps')
            if len(df) > 0:
                gap_durations = df['Gap_Duration_min'].values
                day_numbers = df['Day_Number'].dropna().values if 'Day_Number' in df.columns else []

                # Count gaps by threshold
                total_gaps = len(gap_durations)
                gaps_over_60 = (gap_durations > 60).sum()
                gaps_over_90 = (gap_durations > 90).sum()
                gaps_over_120 = (gap_durations > 120).sum()

                # Analyze periodicity - count gaps per day
                if len(day_numbers) > 0:
                    max_day = int(np.ceil(day_numbers.max())) if len(day_numbers) > 0 else 30
                    gaps_per_day = []
                    significant_gaps_per_day = []  # gaps > 90 min
                    for day in range(max_day + 1):
                        day_mask = (day_numbers >= day) & (day_numbers < day + 1)
                        gaps_per_day.append(day_mask.sum())
                        # Count significant gaps for this day
                        day_gaps = gap_durations[(df['Day_Number'] >= day) & (df['Day_Number'] < day + 1)]
                        significant_gaps_per_day.append((day_gaps > 90).sum() if len(day_gaps) > 0 else 0)
                else:
                    gaps_per_day = []
                    significant_gaps_per_day = []

                gap_data[label] = {
                    'total_gaps': total_gaps,
                    'gaps_over_60': gaps_over_60,
                    'gaps_over_90': gaps_over_90,
                    'gaps_over_120': gaps_over_120,
                    'gaps_per_day': gaps_per_day,
                    'significant_gaps_per_day': significant_gaps_per_day,
                    'raw_gaps': gap_durations,
                    'raw_days': day_numbers,
                }
        except Exception as e:
            # TTC_Gaps sheet may not exist in older results
            print(f"  Note: No TTC_Gaps data for {label} (may need re-run): {e}")
    return gap_data


def generate_gap_analysis_comparison(
    gap_data: Dict[str, Dict],
    ttc_data: Dict[str, Dict],
    output_path: Path,
) -> None:
    """
    Generate gap analysis comparison visualization.

    Shows:
    - Gap count comparison by threshold
    - Significant gaps over time (periodicity view)

    Args:
        gap_data: Dictionary of gap analysis data per provider
        ttc_data: TT&C coverage data for additional context
        output_path: Output path for the plot
    """
    if not gap_data:
        print("  Warning: No gap data available")
        return

    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    labels = list(gap_data.keys())
    colors = [get_provider_color(label) for label in labels]
    x = np.arange(len(labels))
    width = 0.2

    # Panel 1: Gap counts by threshold
    ax1 = axes[0]

    gaps_60 = [gap_data[label].get('gaps_over_60', 0) for label in labels]
    gaps_90 = [gap_data[label].get('gaps_over_90', 0) for label in labels]
    gaps_120 = [gap_data[label].get('gaps_over_120', 0) for label in labels]

    bars1 = ax1.bar(x - width, gaps_60, width, label='> 60 min', color='#fdae61', edgecolor='black')
    bars2 = ax1.bar(x, gaps_90, width, label='> 90 min (1+ rev)', color='#f46d43', edgecolor='black')
    bars3 = ax1.bar(x + width, gaps_120, width, label='> 120 min (2+ rev)', color='#d73027', edgecolor='black')

    ax1.set_ylabel('Number of Gaps', fontsize=11)
    ax1.set_title('TT&C Pass Gaps by Duration Threshold\n(30-day analysis period)', fontsize=12, fontweight='bold')
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=9)
    ax1.legend(fontsize=9)
    ax1.grid(axis='y', alpha=0.3)

    # Add value labels
    for bars in [bars1, bars2, bars3]:
        for bar in bars:
            height = bar.get_height()
            if height > 0:
                ax1.annotate(f'{int(height)}',
                            xy=(bar.get_x() + bar.get_width() / 2, height),
                            xytext=(0, 3), textcoords="offset points",
                            ha='center', va='bottom', fontsize=8)

    # Panel 2: Significant gaps over time (periodicity)
    ax2 = axes[1]

    for label, color in zip(labels, colors):
        sig_gaps = gap_data[label].get('significant_gaps_per_day', [])
        if len(sig_gaps) > 0:
            days = np.arange(len(sig_gaps))
            ax2.plot(days, sig_gaps, 'o-', color=color, label=label, markersize=4, linewidth=1.5, alpha=0.8)

    ax2.set_xlabel('Day of Analysis', fontsize=11)
    ax2.set_ylabel('Significant Gaps (> 90 min)', fontsize=11)
    ax2.set_title('Pass Gap Periodicity\n(significant gaps per day)', fontsize=12, fontweight='bold')
    ax2.legend(fontsize=9)
    ax2.grid(alpha=0.3)
    ax2.set_xlim(-0.5, 30.5)

    # Add weekly markers
    for week in [7, 14, 21, 28]:
        ax2.axvline(week, color='gray', linestyle=':', alpha=0.5)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, facecolor='white', bbox_inches='tight')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_network_comparison_map(
    results_info: Dict[str, Dict],
    output_path: Path,
) -> None:
    """
    Generate a single world map showing all ground station networks overlaid with comm cones.

    Args:
        results_info: Dictionary with config info including results_dir and config_path
        output_path: Output path for the plot
    """
    try:
        import cartopy.crs as ccrs
        import cartopy.feature as cfeature
    except ImportError:
        print("  Warning: Cartopy not available, skipping network map")
        return

    from src.vleo_eo.comparison import load_ground_stations_data

    # Load ground station data for all providers
    results_dirs = {label: info['results_dir'] for label, info in results_info.items()}
    gs_data = load_ground_stations_data(results_dirs)

    if not gs_data:
        print("  Warning: No ground station data available")
        return

    # Load targets from the first config (they're all the same APAC targets)
    targets_gdf = None
    sat_altitude_km = 350.0  # Default
    for label, info in results_info.items():
        config_path = info.get('config_path', '')
        if config_path:
            targets_gdf = load_targets_from_config(config_path)
            sat_altitude_km = get_satellite_altitude_from_config(config_path)
            if targets_gdf is not None:
                break

    # Create single map with PlateCarree projection (like ksat_baseline style)
    fig = plt.figure(figsize=(16, 9))
    ax = plt.axes(projection=ccrs.PlateCarree())
    ax.set_global()

    # Map styling (matching ksat_baseline style)
    land = cfeature.NaturalEarthFeature('physical', 'land', '50m', facecolor='#fafaf9')
    ocean = cfeature.NaturalEarthFeature('physical', 'ocean', '50m', facecolor='#d4dbdc')
    ax.add_feature(ocean, zorder=0)
    ax.add_feature(land, zorder=1)
    ax.add_feature(cfeature.COASTLINE, linewidth=0.5, edgecolor='gray', zorder=2)
    ax.add_feature(cfeature.BORDERS, linewidth=0.5, edgecolor='#ebd6d9', zorder=2)
    ax.gridlines(draw_labels=True, alpha=0.3)

    labels = list(gs_data.keys())
    legend_elements = []

    # Plot each provider's stations and comm cones
    for label in labels:
        df = gs_data[label].copy()
        provider_color = get_provider_color(label)

        # Normalize column names (handle both 'lat'/'lon' and 'Latitude'/'Longitude')
        col_mapping = {}
        for col in df.columns:
            if col.lower() == 'latitude':
                col_mapping[col] = 'lat'
            elif col.lower() == 'longitude':
                col_mapping[col] = 'lon'
            elif col.lower() in ['ka capable', 'ka_capable', 'ka_band']:
                col_mapping[col] = 'ka_band'
            elif col.lower() == 'name':
                col_mapping[col] = 'name'
        df = df.rename(columns=col_mapping)

        if 'lat' not in df.columns or 'lon' not in df.columns:
            print(f"    Warning: No lat/lon columns found for {label}, columns: {list(df.columns)}")
            continue

        # Separate Ka-band and TT&C-only stations
        if 'ka_band' in df.columns:
            ka_df = df[df['ka_band'] == True]
            non_ka_df = df[df['ka_band'] == False]
        else:
            ka_df = pd.DataFrame()
            non_ka_df = df

        # Plot comm cones for Ka stations (10° elevation, thicker lines)
        for _, row in ka_df.iterrows():
            try:
                coverage_radius = calculate_coverage_radius(sat_altitude_km, 10.0)
                circle_lons, circle_lats = generate_coverage_circle(row['lat'], row['lon'], coverage_radius)
                ax.plot(circle_lons, circle_lats, color=provider_color, linewidth=1.2, alpha=0.6,
                        linestyle='-', transform=ccrs.PlateCarree(), zorder=3)
            except Exception:
                pass

        # Plot comm cones for TT&C-only stations (5° elevation, thinner lines)
        for _, row in non_ka_df.iterrows():
            try:
                coverage_radius = calculate_coverage_radius(sat_altitude_km, 5.0)
                circle_lons, circle_lats = generate_coverage_circle(row['lat'], row['lon'], coverage_radius)
                ax.plot(circle_lons, circle_lats, color=provider_color, linewidth=0.6, alpha=0.4,
                        linestyle='-', transform=ccrs.PlateCarree(), zorder=3)
            except Exception:
                pass

        # Plot Ka-band station markers (larger, with labels)
        for _, row in ka_df.iterrows():
            ax.plot(row['lon'], row['lat'], '^', color=provider_color, markersize=10,
                    markeredgecolor='black', markeredgewidth=0.8,
                    transform=ccrs.PlateCarree(), zorder=8)
            # Add station name label
            station_name = row.get('name', row.get('station_name', ''))
            if station_name:
                ax.text(row['lon'], row['lat'] + 2, station_name, fontsize=6,
                        ha='center', va='bottom', color=provider_color,
                        transform=ccrs.PlateCarree(), zorder=9)

        # Plot TT&C-only station markers (smaller)
        for _, row in non_ka_df.iterrows():
            ax.plot(row['lon'], row['lat'], '^', color=provider_color, markersize=7,
                    markeredgecolor='black', markeredgewidth=0.5, alpha=0.8,
                    transform=ccrs.PlateCarree(), zorder=7)

        # Add legend entry for this provider
        ka_count = len(ka_df)
        ttc_count = len(non_ka_df)
        total = len(df)
        legend_elements.append(
            Line2D([0], [0], marker='^', color='w', markersize=10, linestyle='None',
                   markerfacecolor=provider_color, markeredgecolor='black',
                   label=f'{label} ({total} stations: {ka_count} Ka, {ttc_count} TT&C)')
        )
        legend_elements.append(
            Line2D([0], [0], color=provider_color, linewidth=1.2,
                   label=f'{label} Comm Cone')
        )

    # Plot targets as blue X markers
    if targets_gdf is not None and len(targets_gdf) > 0:
        target_lons = []
        target_lats = []
        for _, row in targets_gdf.iterrows():
            try:
                target_lats.append(row.geometry.y)
                target_lons.append(row.geometry.x)
            except Exception:
                pass
        if target_lons:
            ax.scatter(target_lons, target_lats, c='blue', s=60, marker='x', linewidths=2,
                       transform=ccrs.PlateCarree(), zorder=5)
            legend_elements.append(
                Line2D([0], [0], marker='x', color='blue', markersize=8, linestyle='None',
                       markeredgewidth=2, label=f'Targets ({len(targets_gdf)})')
            )

    ax.set_title('Ground Station Network Comparison', fontsize=14, fontweight='bold')
    ax.legend(handles=legend_elements, loc='lower left', fontsize=8, framealpha=0.9)

    plt.savefig(output_path, dpi=150, facecolor='white')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_metrics_comparison_table(
    results: List[ComparisonResult],
    ttc_data: Dict[str, Dict],
    delay_stats: Dict[str, Dict],
    output_path: Path,
    gap_data: Optional[Dict[str, Dict]] = None,
) -> None:
    """
    Generate a comparison table image showing TT&C and Downlink metrics.

    Args:
        results: List of ComparisonResult objects
        ttc_data: TT&C coverage data from Excel
        delay_stats: Downlink delay statistics
        output_path: Output path for the plot
        gap_data: Optional gap analysis data for periodicity metrics
    """
    fig, ax = plt.subplots(figsize=(14, 9))
    ax.axis('off')

    # Prepare data for table
    columns = [r.label for r in results]

    # Build metrics using TT&C and delay data
    metrics = [
        # Network info
        ('', ['NETWORK INFO']),  # Section header
        ('Ground Stations', [r.num_stations for r in results]),
        ('Ka-Band Stations', [r.num_ka_stations for r in results]),
        # TT&C metrics
        ('', ['TT&C COVERAGE']),  # Section header
        ('Contacts/Revolution', [f"{ttc_data.get(r.label, {}).get('Contacts per Revolution', 'N/A')}" for r in results]),
        ('Max TT&C Gap', [f"{ttc_data.get(r.label, {}).get('Maximum Gap', 'N/A')} min" for r in results]),
        ('P95 TT&C Gap', [f"{ttc_data.get(r.label, {}).get('P95 Gap', 'N/A')} min" for r in results]),
        ('1-Pass/Rev Requirement', [ttc_data.get(r.label, {}).get('1 Pass/Rev Requirement Met', 'N/A') for r in results]),
    ]

    # Add gap count metrics if available
    if gap_data:
        metrics.extend([
            ('', ['PASS GAP ANALYSIS']),  # Section header
            ('Total Gaps (30 days)', [f"{gap_data.get(r.label, {}).get('total_gaps', 'N/A')}" for r in results]),
            ('Gaps > 90 min', [f"{gap_data.get(r.label, {}).get('gaps_over_90', 'N/A')}" for r in results]),
            ('Gaps > 120 min', [f"{gap_data.get(r.label, {}).get('gaps_over_120', 'N/A')}" for r in results]),
        ])

    # Downlink metrics
    metrics.extend([
        ('', ['PAYLOAD DOWNLINK']),  # Section header
        ('Delay Median', [f"{delay_stats.get(r.label, {}).get('median', 0):.1f} min" for r in results]),
        ('Delay P90', [f"{delay_stats.get(r.label, {}).get('p90', 0):.1f} min" for r in results]),
        ('Delay P95', [f"{delay_stats.get(r.label, {}).get('p95', 0):.1f} min" for r in results]),
        ('Downlinks < 30 min', [f"{delay_stats.get(r.label, {}).get('under_30min_pct', 0):.0f}%" for r in results]),
        ('Downlinks < 60 min', [f"{delay_stats.get(r.label, {}).get('under_60min_pct', 0):.0f}%" for r in results]),
    ])

    # Filter out section headers for the actual table
    row_labels = []
    cell_data = []
    section_rows = []
    for i, (label, values) in enumerate(metrics):
        if label == '':
            section_rows.append(len(row_labels))
            row_labels.append(values[0])
            cell_data.append([''] * len(columns))
        else:
            row_labels.append(label)
            cell_data.append(values)

    # Create table
    table = ax.table(
        cellText=cell_data,
        rowLabels=row_labels,
        colLabels=columns,
        cellLoc='center',
        loc='center',
    )

    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1.2, 1.6)

    # Style header row with provider colors
    for j in range(len(columns)):
        table[(0, j)].set_facecolor(get_provider_color(results[j].label))
        table[(0, j)].set_text_props(color='white', fontweight='bold')

    # Style section header rows
    for row_idx in section_rows:
        table[(row_idx + 1, -1)].set_facecolor('#e0e0e0')
        table[(row_idx + 1, -1)].set_text_props(fontweight='bold')
        for j in range(len(columns)):
            table[(row_idx + 1, j)].set_facecolor('#e0e0e0')

    ax.set_title('Ground Station Provider Comparison - Key Metrics', fontsize=14, fontweight='bold', pad=20)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, facecolor='white', bbox_inches='tight')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_contact_comparison_chart(
    results: List[ComparisonResult],
    output_path: Path,
) -> None:
    """
    Generate grouped bar chart comparing contact metrics.

    Args:
        results: List of ComparisonResult objects
        output_path: Output path for the plot
    """
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    labels = [r.label for r in results]
    colors = [get_provider_color(label) for label in labels]
    x = np.arange(len(labels))
    width = 0.6

    # Panel 1: Total vs Valid Contacts
    ax1 = axes[0]
    total_contacts = [r.total_contacts for r in results]
    valid_contacts = [r.total_valid_contacts for r in results]

    bars1 = ax1.bar(x - width/4, total_contacts, width/2, label='Total Contacts', color='lightgray', edgecolor='black')
    bars2 = ax1.bar(x + width/4, valid_contacts, width/2, label='Valid Contacts', color=colors, edgecolor='black')

    ax1.set_ylabel('Number of Contacts', fontsize=11)
    ax1.set_title('Contact Volume Comparison', fontsize=12, fontweight='bold')
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=10)
    ax1.legend(fontsize=9)
    ax1.grid(axis='y', alpha=0.3)

    # Add value labels
    for bar in bars1:
        height = bar.get_height()
        ax1.annotate(f'{int(height)}',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', va='bottom', fontsize=9)

    for bar in bars2:
        height = bar.get_height()
        ax1.annotate(f'{int(height)}',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', va='bottom', fontsize=9)

    # Panel 2: Valid Contact Percentage
    ax2 = axes[1]
    valid_pct = [r.valid_contact_pct for r in results]

    bars = ax2.bar(x, valid_pct, width, color=colors, edgecolor='black')

    ax2.set_ylabel('Valid Contact %', fontsize=11)
    ax2.set_title('Contact Success Rate', fontsize=12, fontweight='bold')
    ax2.set_xticks(x)
    ax2.set_xticklabels(labels, fontsize=10)
    ax2.set_ylim(0, 100)
    ax2.grid(axis='y', alpha=0.3)

    # Add value labels
    for bar in bars:
        height = bar.get_height()
        ax2.annotate(f'{height:.1f}%',
                    xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points",
                    ha='center', va='bottom', fontsize=10, fontweight='bold')

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, facecolor='white', bbox_inches='tight')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_delay_distribution_comparison(
    delay_data: Dict[str, pd.DataFrame],
    output_path: Path,
) -> None:
    """
    Generate overlaid box plots or violin plots for delay comparison.

    Args:
        delay_data: Dictionary mapping labels to delay DataFrames
        output_path: Output path for the plot
    """
    if not delay_data:
        print("  Warning: No delay data available")
        return

    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    labels = list(delay_data.keys())
    colors = [get_provider_color(label) for label in labels]

    # Extract delay values
    delay_values = []
    for label in labels:
        df = delay_data[label]
        if 'Downlink_Delay_minutes' in df.columns:
            values = df['Downlink_Delay_minutes'].dropna().values
            delay_values.append(values)
        else:
            delay_values.append(np.array([]))

    # Panel 1: Box plots
    ax1 = axes[0]
    bp = ax1.boxplot(delay_values, labels=labels, patch_artist=True)

    for patch, color in zip(bp['boxes'], colors):
        patch.set_facecolor(color)
        patch.set_alpha(0.7)

    ax1.set_ylabel('Downlink Delay (minutes)', fontsize=11)
    ax1.set_title('Downlink Delay Distribution', fontsize=12, fontweight='bold')
    ax1.grid(axis='y', alpha=0.3)

    # Panel 2: Histogram overlay
    ax2 = axes[1]
    max_delay = max(v.max() if len(v) > 0 else 0 for v in delay_values)
    bins = np.linspace(0, min(max_delay, 200), 50)

    for i, (label, values) in enumerate(zip(labels, delay_values)):
        if len(values) > 0:
            ax2.hist(values, bins=bins, alpha=0.5, label=label, color=colors[i], edgecolor='black', linewidth=0.5)

    ax2.set_xlabel('Downlink Delay (minutes)', fontsize=11)
    ax2.set_ylabel('Frequency', fontsize=11)
    ax2.set_title('Downlink Delay Histogram', fontsize=12, fontweight='bold')
    ax2.legend(fontsize=9)
    ax2.grid(alpha=0.3)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, facecolor='white', bbox_inches='tight')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_delay_cdf_comparison(
    delay_data: Dict[str, pd.DataFrame],
    output_path: Path,
) -> None:
    """
    Generate cumulative distribution function comparison for delays.

    Args:
        delay_data: Dictionary mapping labels to delay DataFrames
        output_path: Output path for the plot
    """
    if not delay_data:
        print("  Warning: No delay data available")
        return

    fig, ax = plt.subplots(figsize=(10, 6))

    labels = list(delay_data.keys())
    colors = [get_provider_color(label) for label in labels]

    for label, color in zip(labels, colors):
        df = delay_data[label]
        if 'Downlink_Delay_minutes' in df.columns:
            values = df['Downlink_Delay_minutes'].dropna().sort_values()
            if len(values) > 0:
                cdf = np.arange(1, len(values) + 1) / len(values) * 100
                ax.plot(values, cdf, label=label, color=color, linewidth=2)

    # Add reference lines
    ax.axhline(y=50, color='gray', linestyle='--', alpha=0.5, label='50th percentile')
    ax.axhline(y=90, color='gray', linestyle=':', alpha=0.5, label='90th percentile')
    ax.axhline(y=95, color='gray', linestyle='-.', alpha=0.5, label='95th percentile')

    ax.set_xlabel('Downlink Delay (minutes)', fontsize=11)
    ax.set_ylabel('Cumulative Percentage (%)', fontsize=11)
    ax.set_title('Downlink Delay CDF Comparison', fontsize=12, fontweight='bold')
    ax.legend(fontsize=9, loc='lower right')
    ax.grid(alpha=0.3)
    ax.set_xlim(0, None)
    ax.set_ylim(0, 100)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, facecolor='white', bbox_inches='tight')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_throughput_comparison(
    results: List[ComparisonResult],
    output_path: Path,
) -> None:
    """
    Generate data throughput comparison chart.

    Args:
        results: List of ComparisonResult objects
        output_path: Output path for the plot
    """
    fig, axes = plt.subplots(1, 2, figsize=(14, 5))

    labels = [r.label for r in results]
    colors = [get_provider_color(label) for label in labels]
    x = np.arange(len(labels))
    width = 0.35

    # Panel 1: Collected vs Downlinked
    ax1 = axes[0]
    collected = [r.total_data_collected_gb for r in results]
    downlinked = [r.total_data_downlinked_gb for r in results]

    bars1 = ax1.bar(x - width/2, collected, width, label='Data Collected', color='lightblue', edgecolor='black')
    bars2 = ax1.bar(x + width/2, downlinked, width, label='Data Downlinked', color=colors, edgecolor='black')

    ax1.set_ylabel('Data Volume (GB)', fontsize=11)
    ax1.set_title('Data Collection vs Downlink', fontsize=12, fontweight='bold')
    ax1.set_xticks(x)
    ax1.set_xticklabels(labels, fontsize=10)
    ax1.legend(fontsize=9)
    ax1.grid(axis='y', alpha=0.3)

    # Add value labels
    for bar in bars1:
        height = bar.get_height()
        if height > 0:
            ax1.annotate(f'{height:.1f}',
                        xy=(bar.get_x() + bar.get_width() / 2, height),
                        xytext=(0, 3), textcoords="offset points",
                        ha='center', va='bottom', fontsize=9)

    for bar in bars2:
        height = bar.get_height()
        if height > 0:
            ax1.annotate(f'{height:.1f}',
                        xy=(bar.get_x() + bar.get_width() / 2, height),
                        xytext=(0, 3), textcoords="offset points",
                        ha='center', va='bottom', fontsize=9)

    # Panel 2: Downlink Efficiency
    ax2 = axes[1]
    efficiency = [r.downlink_efficiency_pct for r in results]

    bars3 = ax2.bar(x, efficiency, width * 1.5, label='Downlink Efficiency', color=colors, edgecolor='black', alpha=0.7)

    ax2.set_ylabel('Downlink Efficiency (%)', fontsize=11)
    ax2.set_title('Downlink Efficiency', fontsize=12, fontweight='bold')
    ax2.set_xticks(x)
    ax2.set_xticklabels(labels, fontsize=10)
    ax2.set_ylim(0, 110)
    ax2.grid(axis='y', alpha=0.3)

    # Add value labels
    for bar in bars3:
        height = bar.get_height()
        if height > 0:
            ax2.annotate(f'{height:.1f}%',
                        xy=(bar.get_x() + bar.get_width() / 2, height),
                        xytext=(0, 3), textcoords="offset points",
                        ha='center', va='bottom', fontsize=9)

    plt.tight_layout()
    plt.savefig(output_path, dpi=150, facecolor='white', bbox_inches='tight')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_summary_slide(
    results: List[ComparisonResult],
    ttc_data: Dict[str, Dict],
    delay_stats: Dict[str, Dict],
    opt_data: Dict[str, Dict],
    output_path: Path,
) -> None:
    """
    Generate summary slide with key findings focused on TT&C and Downlink.

    Args:
        results: List of ComparisonResult objects
        ttc_data: TT&C coverage data
        delay_stats: Downlink delay statistics
        opt_data: Optimization data (stations ranked by contacts)
        output_path: Output path for the plot
    """
    fig = plt.figure(figsize=(16, 9))

    # Left panel: TT&C findings
    ax1 = fig.add_axes([0.02, 0.05, 0.46, 0.88])
    ax1.axis('off')

    ttc_lines = ["TT&C COVERAGE FINDINGS", "=" * 45, ""]

    # Compare TT&C metrics
    ttc_lines.append("REQUIREMENT: 1 TT&C Pass per Revolution")
    ttc_lines.append("")

    for r in results:
        label = r.label
        data = ttc_data.get(label, {})
        ttc_lines.append(f"{label}:")
        ttc_lines.append(f"  Contacts/Rev: {data.get('Contacts per Revolution', 'N/A')}")
        ttc_lines.append(f"  Max Gap: {data.get('Maximum Gap', 'N/A')} min")
        ttc_lines.append(f"  P95 Gap: {data.get('P95 Gap', 'N/A')} min")
        ttc_lines.append(f"  Requirement Met: {data.get('1 Pass/Rev Requirement Met', 'N/A')}")
        ttc_lines.append("")

    # Optimization findings
    ttc_lines.append("-" * 45)
    ttc_lines.append("NETWORK OPTIMIZATION")
    ttc_lines.append("")

    for r in results:
        label = r.label
        opt = opt_data.get(label, {})
        stations = opt.get('stations_ranked', [])
        ttc_lines.append(f"{label} ({opt.get('total_stations', 0)} stations):")
        # Show top 3 stations by TT&C contacts
        ttc_lines.append("  Top stations by TT&C contacts:")
        for i, s in enumerate(stations[:3]):
            ka_str = "(Ka)" if s.get('ka_capable') else ""
            ttc_lines.append(f"    {i+1}. {s['name']}: {s['contacts']} contacts {ka_str}")
        ttc_lines.append("")

    ax1.text(0.02, 0.98, "\n".join(ttc_lines), transform=ax1.transAxes,
             fontsize=9, verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#e3f2fd', alpha=0.8))

    # Right panel: Downlink findings
    ax2 = fig.add_axes([0.52, 0.05, 0.46, 0.88])
    ax2.axis('off')

    dl_lines = ["PAYLOAD DOWNLINK FINDINGS", "=" * 45, ""]

    # Find best performer for delay
    delay_p95_values = [(r.label, delay_stats.get(r.label, {}).get('p95', 999)) for r in results]
    best_delay = min(delay_p95_values, key=lambda x: x[1])
    best_under30 = max(results, key=lambda r: delay_stats.get(r.label, {}).get('under_30min_pct', 0))

    dl_lines.append("DELAY COMPARISON")
    dl_lines.append("")

    for r in results:
        label = r.label
        stats = delay_stats.get(label, {})
        dl_lines.append(f"{label}:")
        dl_lines.append(f"  Median: {stats.get('median', 0):.1f} min")
        dl_lines.append(f"  P90: {stats.get('p90', 0):.1f} min")
        dl_lines.append(f"  P95: {stats.get('p95', 0):.1f} min")
        dl_lines.append(f"  < 30 min: {stats.get('under_30min_pct', 0):.0f}%")
        dl_lines.append(f"  < 60 min: {stats.get('under_60min_pct', 0):.0f}%")
        dl_lines.append("")

    dl_lines.append("-" * 45)
    dl_lines.append("KEY INSIGHTS")
    dl_lines.append("")
    dl_lines.append(f"Best P95 Delay: {best_delay[0]} ({best_delay[1]:.1f} min)")
    dl_lines.append(f"Best <30min %: {best_under30.label} ({delay_stats.get(best_under30.label, {}).get('under_30min_pct', 0):.0f}%)")
    dl_lines.append("")

    # Add recommendation
    dl_lines.append("-" * 45)
    dl_lines.append("RECOMMENDATION")
    dl_lines.append("")

    # Determine overall winner
    best_ttc_gap = min(results, key=lambda r: float(ttc_data.get(r.label, {}).get('P95 Gap', 999)))
    best_overall_delay = min(results, key=lambda r: delay_stats.get(r.label, {}).get('p95', 999))

    if best_ttc_gap.label == best_overall_delay.label:
        dl_lines.append(f"Overall Best: {best_ttc_gap.label}")
        dl_lines.append("  Excels in both TT&C and Downlink metrics")
    else:
        dl_lines.append(f"Best TT&C: {best_ttc_gap.label}")
        dl_lines.append(f"Best Downlink: {best_overall_delay.label}")
        dl_lines.append("  Trade-off depends on mission priorities")

    ax2.text(0.02, 0.98, "\n".join(dl_lines), transform=ax2.transAxes,
             fontsize=9, verticalalignment='top', fontfamily='monospace',
             bbox=dict(boxstyle='round', facecolor='#fff3e0', alpha=0.8))

    fig.suptitle('Summary & Key Findings', fontsize=14, fontweight='bold', y=0.98)

    plt.savefig(output_path, dpi=150, facecolor='white')
    plt.close(fig)
    print(f"  Generated: {output_path.name}")


def generate_all_comparison_plots(
    results: List[ComparisonResult],
    delay_data: Dict[str, pd.DataFrame],
    results_info: Dict[str, Dict],
    plots_dir: Path,
) -> None:
    """
    Generate all comparison plots.

    Args:
        results: List of ComparisonResult objects
        delay_data: Dictionary of delay DataFrames
        results_info: Dictionary with config info
        plots_dir: Output directory for plots
    """
    print("\nGenerating comparison plots...")

    # Load additional data for comparison
    ttc_data = load_ttc_data_from_excel(results_info)
    delay_stats = load_downlink_delay_stats(results_info)
    opt_data = load_optimization_data(results_info)
    gap_data = load_gap_data_from_excel(results_info)

    # Network map
    print("  - Network comparison map...")
    generate_network_comparison_map(results_info, plots_dir / 'network_map.png')

    # Metrics table (now with TT&C, delay, and gap data)
    print("  - Metrics comparison table...")
    generate_metrics_comparison_table(results, ttc_data, delay_stats, plots_dir / 'metrics_table.png', gap_data)

    # Contact comparison
    print("  - Contact comparison chart...")
    generate_contact_comparison_chart(results, plots_dir / 'contact_comparison.png')

    # Gap analysis comparison (if gap data available)
    if gap_data:
        print("  - Gap analysis comparison...")
        generate_gap_analysis_comparison(gap_data, ttc_data, plots_dir / 'gap_analysis.png')

    # Delay distribution
    print("  - Delay distribution comparison...")
    generate_delay_distribution_comparison(delay_data, plots_dir / 'delay_distribution.png')

    # Delay CDF
    print("  - Delay CDF comparison...")
    generate_delay_cdf_comparison(delay_data, plots_dir / 'delay_cdf.png')

    # Summary slide (now with TT&C, delay, and optimization data)
    print("  - Summary slide...")
    generate_summary_slide(results, ttc_data, delay_stats, opt_data, plots_dir / 'summary.png')


def build_comparison_presentation(
    plots_dir: Path,
    results: List[ComparisonResult],
    results_info: Dict[str, Dict],
    output_path: Path,
) -> Path:
    """
    Build comparison PowerPoint presentation.

    Args:
        plots_dir: Directory containing plot images
        results: List of ComparisonResult objects
        results_info: Dictionary with config info
        output_path: Output path for PowerPoint file

    Returns:
        Path to generated PowerPoint file
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide layouts
    blank_layout = prs.slide_layouts[6]  # Blank layout

    def add_title_slide(title: str, subtitle: str = ""):
        """Add a title slide."""
        slide = prs.slides.add_slide(blank_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.0), Inches(12.333), Inches(1.0)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.alignment = PP_ALIGN.CENTER

        return slide

    def add_image_slide(title: str, image_path: Path):
        """Add a slide with an image."""
        if not image_path.exists():
            print(f"  Warning: Image not found: {image_path}")
            return None

        slide = prs.slides.add_slide(blank_layout)

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.1), Inches(12.9), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True

        # Image
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.2), Inches(0.7),
            width=Inches(12.9)
        )

        return slide

    # Build presentation
    print("\n  Building PowerPoint slides...")

    # Slide 1: Title
    providers = [r.label for r in results]
    add_title_slide(
        "Ground Station Provider Comparison",
        f"Comparing: {', '.join(providers)}"
    )

    # Slide 2: Network map
    add_image_slide("Ground Station Network Comparison", plots_dir / 'network_map.png')

    # Slide 3: Metrics table
    add_image_slide("Key Metrics Comparison", plots_dir / 'metrics_table.png')

    # Slide 4: Contact comparison
    add_image_slide("Contact Volume & Success Rate", plots_dir / 'contact_comparison.png')

    # Slide 5: Gap analysis (if available)
    gap_analysis_path = plots_dir / 'gap_analysis.png'
    if gap_analysis_path.exists():
        add_image_slide("TT&C Pass Gap Analysis & Periodicity", gap_analysis_path)

    # Slide 6: Delay distribution
    add_image_slide("Downlink Delay Distribution", plots_dir / 'delay_distribution.png')

    # Slide 7: Delay CDF
    add_image_slide("Downlink Delay Cumulative Distribution", plots_dir / 'delay_cdf.png')

    # Slide 8: Summary
    add_image_slide("Summary & Key Findings", plots_dir / 'summary.png')

    # Save presentation
    prs.save(output_path)
    print(f"  Comparison PowerPoint saved to: {output_path}")

    return output_path


if __name__ == '__main__':
    # Test with existing results
    import sys

    print("Comparison Report Generator")
    print("=" * 50)
    print("This script is typically called from run_comparison.py")
    print("Run: python run_comparison.py")
